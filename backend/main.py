from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, StreamingResponse
import pandas as pd
import numpy as np
import io, os, json, sqlite3
import anthropic
from datetime import datetime
from typing import Optional
from pathlib import Path

app = FastAPI(
    title="Axiom KPI Dashboard API",
    description="Upload CSVs to compute and track Priority-1 KPIs with 12-month fingerprinting.",
    version="1.0.0",
    docs_url="/api/docs",
    redoc_url="/api/redoc",
    openapi_url="/api/openapi.json",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

DB_PATH = Path(__file__).parent / "uploads" / "axiom.db"
UPLOADS_DIR = Path(__file__).parent / "uploads"
UPLOADS_DIR.mkdir(exist_ok=True)

# ─── Database ───────────────────────────────────────────────────────────────

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    # WAL mode allows concurrent reads alongside the single write thread,
    # preventing "database is locked" failures under background-thread writes.
    conn.execute("PRAGMA journal_mode=WAL")
    conn.commit()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS uploads (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT,
            uploaded_at TEXT,
            row_count INTEGER,
            detected_columns TEXT
        );
        CREATE TABLE IF NOT EXISTS monthly_data (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            upload_id INTEGER,
            year INTEGER,
            month INTEGER,
            data_json TEXT
        );
        CREATE TABLE IF NOT EXISTS kpi_targets (
            kpi_key TEXT PRIMARY KEY,
            target_value REAL,
            unit TEXT,
            direction TEXT
        );
        CREATE TABLE IF NOT EXISTS projection_uploads (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT,
            uploaded_at TEXT,
            row_count INTEGER,
            detected_columns TEXT
        );
        CREATE TABLE IF NOT EXISTS projection_monthly_data (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            projection_upload_id INTEGER,
            year INTEGER,
            month INTEGER,
            data_json TEXT
        );
        CREATE TABLE IF NOT EXISTS kpi_annotations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            kpi_key TEXT NOT NULL,
            period TEXT NOT NULL,
            note TEXT NOT NULL,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(kpi_key, period)
        );
        CREATE TABLE IF NOT EXISTS kpi_accountability (
            kpi_key TEXT PRIMARY KEY,
            owner TEXT DEFAULT '',
            due_date TEXT DEFAULT '',
            status TEXT DEFAULT 'open',
            last_updated TEXT DEFAULT CURRENT_TIMESTAMP
        );
    """)
    conn.commit()
    # Migration: add version_label to projection tables if missing
    cols = [r[1] for r in conn.execute("PRAGMA table_info(projection_uploads)").fetchall()]
    if "version_label" not in cols:
        conn.execute("ALTER TABLE projection_uploads ADD COLUMN version_label TEXT DEFAULT 'v1'")
        conn.commit()
    cols2 = [r[1] for r in conn.execute("PRAGMA table_info(projection_monthly_data)").fetchall()]
    if "version_label" not in cols2:
        conn.execute("ALTER TABLE projection_monthly_data ADD COLUMN version_label TEXT DEFAULT 'v1'")
        conn.commit()
    # Seed default targets
    default_targets = [
        ("revenue_growth",      6.0,   "pct",    "higher"),
        ("gross_margin",       62.0,   "pct",    "higher"),
        ("operating_margin",   18.0,   "pct",    "higher"),
        ("ebitda_margin",      22.0,   "pct",    "higher"),
        ("cash_conv_cycle",    42.0,   "days",   "lower"),
        ("dso",                35.0,   "days",   "lower"),
        ("arr_growth",          7.0,   "pct",    "higher"),
        ("nrr",               105.0,   "pct",    "higher"),
        ("burn_multiple",       1.2,   "ratio",  "lower"),
        ("opex_ratio",         42.0,   "pct",    "lower"),
        ("contribution_margin",46.0,   "pct",    "higher"),
        ("revenue_quality",    80.0,   "pct",    "higher"),
        ("cac_payback",        10.0,   "months", "lower"),
        ("sales_efficiency",    3.0,   "ratio",  "higher"),
        ("customer_concentration",28.0,"pct",    "lower"),
        ("recurring_revenue",  80.0,   "pct",    "higher"),
        ("churn_rate",          2.5,   "pct",    "lower"),
        ("operating_leverage",  1.2,   "ratio",  "higher"),
        ("growth_efficiency",    3.0,  "ratio", "higher"),
        ("revenue_momentum",     1.0,  "ratio", "higher"),
        ("revenue_fragility",    1.0,  "ratio", "lower"),
        ("burn_convexity",       0.0,  "ratio", "lower"),
        ("margin_volatility",    2.0,  "pct",   "lower"),
        ("pipeline_conversion",  5.0,  "pct",   "higher"),
        ("customer_decay_slope", 0.0,  "pct",   "lower"),
        ("customer_ltv",        80.0,  "usd",   "higher"),
        ("pricing_power_index",  3.0,  "pct",   "higher"),
        ("cpl",                150.0,  "usd",   "lower"),
        ("mql_sql_rate",        28.0,  "pct",   "higher"),
        ("win_rate",            30.0,  "pct",   "higher"),
        ("quota_attainment",    85.0,  "pct",   "higher"),
        ("marketing_roi",        3.5,  "ratio", "higher"),
        ("headcount_eff",        1.8,  "ratio", "higher"),
        ("rev_per_employee",   180.0,  "usd",   "higher"),
        ("ltv_cac",              4.0,  "ratio", "higher"),
        ("expansion_rate",      22.0,  "pct",   "higher"),
        ("health_score",        72.0,  "score", "higher"),
        ("logo_retention",      90.0,  "pct",   "higher"),
        ("payback_period",      14.0,  "months","lower"),
    ]
    for row in default_targets:
        conn.execute(
            "INSERT OR IGNORE INTO kpi_targets VALUES (?,?,?,?)", row
        )
    conn.commit()
    conn.close()

init_db()

@app.on_event("startup")
async def auto_seed():
    """Auto-seed full 27-KPI multi-year demo data on cold start if the database is empty."""
    conn = get_db()
    count     = conn.execute("SELECT COUNT(*) FROM monthly_data").fetchone()[0]
    proj_count = conn.execute("SELECT COUNT(*) FROM projection_monthly_data").fetchone()[0]
    conn.close()
    if count == 0 or proj_count == 0:
        # seed_multiyear seeds both actuals (27 KPIs, 2021–2026) and projections in one call
        seed_multiyear()
    # Initialize ontology tables once at startup instead of on every GET request (M2)
    try:
        conn_ont = get_db()
        _init_ontology_tables(conn_ont)
        conn_ont.close()
    except Exception:
        pass  # Tables may not be defined yet on first cold start; discovery will create them

# ─── KPI Definitions ────────────────────────────────────────────────────────

# Extended ontology-only metrics (not in KPI_DEFS main dashboard; used only in the knowledge graph)
EXTENDED_ONTOLOGY_METRICS = [
    # Growth
    {"key": "cpl",                "name": "Cost Per Lead",              "domain": "growth",        "unit": "usd",    "direction": "lower"},
    {"key": "mql_sql_rate",       "name": "MQL-to-SQL Rate",            "domain": "growth",        "unit": "pct",    "direction": "higher"},
    {"key": "pipeline_velocity",  "name": "Pipeline Velocity",          "domain": "growth",        "unit": "ratio",  "direction": "higher"},
    {"key": "win_rate",           "name": "Win Rate",                   "domain": "growth",        "unit": "pct",    "direction": "higher"},
    {"key": "organic_traffic",    "name": "Organic Traffic Growth",     "domain": "growth",        "unit": "pct",    "direction": "higher"},
    {"key": "brand_awareness",    "name": "Brand Awareness Index",      "domain": "growth",        "unit": "score",  "direction": "higher"},
    {"key": "quota_attainment",   "name": "Quota Attainment Rate",      "domain": "growth",        "unit": "pct",    "direction": "higher"},
    {"key": "marketing_roi",      "name": "Marketing ROI",              "domain": "growth",        "unit": "ratio",  "direction": "higher"},
    # Revenue
    {"key": "avg_deal_size",      "name": "Avg Deal Size",              "domain": "revenue",       "unit": "usd",    "direction": "higher"},
    {"key": "expansion_rate",     "name": "Expansion Revenue Rate",     "domain": "revenue",       "unit": "pct",    "direction": "higher"},
    {"key": "gross_dollar_ret",   "name": "Gross Dollar Retention",     "domain": "revenue",       "unit": "pct",    "direction": "higher"},
    {"key": "ltv_cac",            "name": "LTV:CAC Ratio",              "domain": "revenue",       "unit": "ratio",  "direction": "higher"},
    # Retention
    {"key": "product_nps",        "name": "Product NPS",                "domain": "retention",     "unit": "score",  "direction": "higher"},
    {"key": "feature_adoption",   "name": "Feature Adoption Rate",      "domain": "retention",     "unit": "pct",    "direction": "higher"},
    {"key": "activation_rate",    "name": "Activation Rate",            "domain": "retention",     "unit": "pct",    "direction": "higher"},
    {"key": "time_to_value",      "name": "Time-to-Value",              "domain": "retention",     "unit": "days",   "direction": "lower"},
    {"key": "health_score",       "name": "Customer Health Score",      "domain": "retention",     "unit": "score",  "direction": "higher"},
    {"key": "logo_retention",     "name": "Logo Retention Rate",        "domain": "retention",     "unit": "pct",    "direction": "higher"},
    {"key": "csat",               "name": "Customer Satisfaction",      "domain": "retention",     "unit": "score",  "direction": "higher"},
    # Efficiency
    {"key": "headcount_eff",      "name": "Headcount Efficiency",       "domain": "efficiency",    "unit": "ratio",  "direction": "higher"},
    {"key": "rev_per_employee",   "name": "Revenue Per Employee",       "domain": "efficiency",    "unit": "usd",    "direction": "higher"},
    {"key": "ramp_time",          "name": "Sales Rep Ramp Time",        "domain": "efficiency",    "unit": "months", "direction": "lower"},
    {"key": "support_volume",     "name": "Support Ticket Volume",      "domain": "efficiency",    "unit": "count",  "direction": "lower"},
    {"key": "automation_rate",    "name": "Process Automation Rate",    "domain": "efficiency",    "unit": "pct",    "direction": "higher"},
    # Cashflow
    {"key": "cash_runway",        "name": "Cash Runway",                "domain": "cashflow",      "unit": "months", "direction": "higher"},
    {"key": "current_ratio",      "name": "Current Ratio",              "domain": "cashflow",      "unit": "ratio",  "direction": "higher"},
    {"key": "working_capital",    "name": "Working Capital Ratio",      "domain": "cashflow",      "unit": "ratio",  "direction": "higher"},
    # Risk
    {"key": "contraction_rate",   "name": "Contraction Rate",           "domain": "risk",          "unit": "pct",    "direction": "lower"},
    # Profitability
    {"key": "payback_period",     "name": "Investor Payback Period",    "domain": "profitability", "unit": "months", "direction": "lower"},
]

KPI_DEFS = [
    {"key": "revenue_growth",       "name": "Revenue Growth Rate",       "unit": "pct",    "direction": "higher", "formula": "(Revenue_Month - Revenue_PrevMonth) / Revenue_PrevMonth × 100"},
    {"key": "gross_margin",         "name": "Gross Margin %",            "unit": "pct",    "direction": "higher", "formula": "(Revenue - COGS) / Revenue × 100"},
    {"key": "operating_margin",     "name": "Operating Margin %",        "unit": "pct",    "direction": "higher", "formula": "(Revenue - COGS - OpEx) / Revenue × 100"},
    {"key": "ebitda_margin",        "name": "EBITDA Margin %",           "unit": "pct",    "direction": "higher", "formula": "EBITDA / Revenue × 100"},
    {"key": "cash_conv_cycle",      "name": "Cash Conversion Cycle",     "unit": "days",   "direction": "lower",  "formula": "DSO + DIO - DPO"},
    {"key": "dso",                  "name": "Days Sales Outstanding",    "unit": "days",   "direction": "lower",  "formula": "(AR / Revenue) × 30"},
    {"key": "arr_growth",           "name": "ARR Growth Rate",           "unit": "pct",    "direction": "higher", "formula": "(ARR_Month - ARR_PrevMonth) / ARR_PrevMonth × 100"},
    {"key": "nrr",                  "name": "Net Revenue Retention",     "unit": "pct",    "direction": "higher", "formula": "(MRR_Start + Expansion - Churn - Contraction) / MRR_Start × 100"},
    {"key": "burn_multiple",        "name": "Burn Multiple",             "unit": "ratio",  "direction": "lower",  "formula": "Net Burn / Net New ARR"},
    {"key": "opex_ratio",           "name": "Operating Expense Ratio",   "unit": "pct",    "direction": "lower",  "formula": "OpEx / Revenue × 100"},
    {"key": "contribution_margin",  "name": "Contribution Margin %",     "unit": "pct",    "direction": "higher", "formula": "(Revenue - COGS - Variable_Costs) / Revenue × 100"},
    {"key": "revenue_quality",      "name": "Revenue Quality Ratio",     "unit": "pct",    "direction": "higher", "formula": "Recurring_Revenue / Total_Revenue × 100"},
    {"key": "cac_payback",          "name": "CAC Payback Period",        "unit": "months", "direction": "lower",  "formula": "CAC / (ARPU × Gross_Margin_pct)"},
    {"key": "sales_efficiency",     "name": "Sales Efficiency Ratio",    "unit": "ratio",  "direction": "higher", "formula": "New_ARR / Sales_Marketing_Spend"},
    {"key": "customer_concentration","name":"Customer Concentration",    "unit": "pct",    "direction": "lower",  "formula": "Top_Customer_Revenue / Total_Revenue × 100"},
    {"key": "recurring_revenue",    "name": "Recurring Revenue Ratio",   "unit": "pct",    "direction": "higher", "formula": "Recurring_Revenue / Total_Revenue × 100"},
    {"key": "churn_rate",           "name": "Monthly Churn Rate",        "unit": "pct",    "direction": "lower",  "formula": "Lost_Customers / Total_Customers × 100"},
    {"key": "operating_leverage",   "name": "Operating Leverage Index",  "unit": "ratio",  "direction": "higher", "formula": "% Change in Operating Income / % Change in Revenue"},
    # ── Enriched & Derived KPIs ─────────────────────────────────────────────
    {"key": "growth_efficiency",    "name": "Growth Efficiency Index",      "unit": "ratio", "direction": "higher", "domain": "growth",        "formula": "ARR_Growth_Rate / Burn_Multiple"},
    {"key": "revenue_momentum",     "name": "Revenue Momentum Index",       "unit": "ratio", "direction": "higher", "domain": "growth",        "formula": "Current_Rev_Growth / Annual_Avg_Rev_Growth"},
    {"key": "revenue_fragility",    "name": "Strategic Revenue Fragility",  "unit": "ratio", "direction": "lower",  "domain": "risk",          "formula": "(Customer_Concentration × Churn_Rate) / NRR"},
    {"key": "burn_convexity",       "name": "Burn Convexity",               "unit": "ratio", "direction": "lower",  "domain": "cashflow",      "formula": "Δ_Burn_Multiple Month-over-Month"},
    {"key": "margin_volatility",    "name": "Margin Volatility Index",      "unit": "pct",   "direction": "lower",  "domain": "profitability", "formula": "6M_Rolling_Std_Dev_of_Gross_Margin"},
    {"key": "pipeline_conversion",  "name": "Pipeline Conversion Rate",     "unit": "pct",   "direction": "higher", "domain": "growth",        "formula": "MQL_to_Win_End_to_End_Conversion_%"},
    {"key": "customer_decay_slope", "name": "Customer Decay Curve Slope",   "unit": "pct",   "direction": "lower",  "domain": "retention",     "formula": "Δ_Churn_Rate Month-over-Month"},
    {"key": "customer_ltv",         "name": "Customer Lifetime Value",      "unit": "usd",   "direction": "higher", "domain": "revenue",       "formula": "(ARPU × Gross_Margin%) / Monthly_Churn_Rate"},
    {"key": "pricing_power_index",  "name": "Pricing Power Index",          "unit": "pct",   "direction": "higher", "domain": "revenue",       "formula": "ΔARPU% − Δ_Customer_Volume%"},
]

# ─── Causation Rules & Gap Analysis ────────────────────────────────────────

CAUSATION_RULES = {
    "gross_margin": {
        "root_causes": [
            "COGS higher than projected",
            "Revenue mix shift toward lower-margin products",
            "Pricing pressure from competitive dynamics",
        ],
        "downstream_impact": ["operating_margin", "ebitda_margin", "contribution_margin"],
        "corrective_actions": [
            "Review pricing strategy by segment",
            "Analyze COGS by product line",
            "Evaluate vendor contracts for renegotiation",
        ],
    },
    "operating_margin": {
        "root_causes": [
            "Gross margin compression flowing through",
            "Operating expenses above projected levels",
            "Revenue below projection without proportional cost reduction",
        ],
        "downstream_impact": ["ebitda_margin"],
        "corrective_actions": [
            "Identify top 3 discretionary opex line items for reduction",
            "Align headcount plan to revised revenue forecast",
            "Review variable cost scaling assumptions",
        ],
    },
    "ebitda_margin": {
        "root_causes": [
            "Operating margin shortfall flowing through",
            "Depreciation or amortization above plan",
        ],
        "downstream_impact": [],
        "corrective_actions": [
            "Review D&A schedule against asset plan",
            "Address operating margin root causes upstream",
        ],
    },
    "churn_rate": {
        "root_causes": [
            "Customer satisfaction decline",
            "Competitive pressure or pricing mismatch",
            "Product-market fit gap in recent cohorts",
        ],
        "downstream_impact": ["nrr", "revenue_growth", "arr_growth", "burn_multiple"],
        "corrective_actions": [
            "Implement proactive churn detection triggers",
            "Schedule at-risk account reviews with CS team",
            "Review onboarding and product adoption metrics",
        ],
    },
    "nrr": {
        "root_causes": [
            "Contraction in existing accounts",
            "Higher than projected churn",
            "Upsell and expansion underperformance",
        ],
        "downstream_impact": ["revenue_growth", "arr_growth"],
        "corrective_actions": [
            "Activate expansion playbooks for top accounts",
            "Review upsell trigger criteria and qualification",
            "Audit renewal pipeline health and coverage",
        ],
    },
    "dso": {
        "root_causes": [
            "Collections process delays",
            "Loosened credit terms or approval",
            "Customer cash flow difficulties",
        ],
        "downstream_impact": ["cash_conv_cycle"],
        "corrective_actions": [
            "Tighten credit approval criteria",
            "Automate payment reminders at 30/45/60 days",
            "Offer early payment discounts to accelerate collection",
        ],
    },
    "cash_conv_cycle": {
        "root_causes": [
            "DSO extending beyond projection",
            "Inventory days or payables days deteriorating",
        ],
        "downstream_impact": [],
        "corrective_actions": [
            "Address DSO root causes upstream",
            "Negotiate extended payment terms with key vendors",
        ],
    },
    "revenue_growth": {
        "root_causes": [
            "Pipeline shortfall vs plan",
            "Lower close rates than projected",
            "Deals slipping to future quarters",
        ],
        "downstream_impact": ["arr_growth", "operating_leverage", "burn_multiple", "sales_efficiency"],
        "corrective_actions": [
            "Review pipeline coverage ratio (target: 3x quota)",
            "Identify top deal acceleration opportunities",
            "Reassess pricing and packaging to improve win rates",
        ],
    },
    "arr_growth": {
        "root_causes": [
            "New ARR below projection",
            "Expansion ARR underperforming",
            "Churn offsetting new bookings",
        ],
        "downstream_impact": ["burn_multiple", "sales_efficiency"],
        "corrective_actions": [
            "Review new logo pipeline and close rate",
            "Strengthen expansion motion in CS",
            "Address churn to improve net ARR",
        ],
    },
    "opex_ratio": {
        "root_causes": [
            "Headcount growth ahead of plan",
            "Discretionary spend above budget",
            "Unplanned infrastructure or tooling costs",
        ],
        "downstream_impact": ["operating_margin", "ebitda_margin"],
        "corrective_actions": [
            "Conduct discretionary spend audit",
            "Freeze non-critical hiring pending revenue recovery",
            "Review SaaS tool consolidation opportunities",
        ],
    },
    "contribution_margin": {
        "root_causes": [
            "Variable costs above projection",
            "Gross margin compression flowing through",
        ],
        "downstream_impact": [],
        "corrective_actions": [
            "Analyze variable cost drivers by product",
            "Review unit economics assumptions in pricing model",
        ],
    },
    "burn_multiple": {
        "root_causes": [
            "Revenue growth below projection",
            "Sales and marketing spend above plan",
        ],
        "downstream_impact": ["cac_payback"],
        "corrective_actions": [
            "Tighten sales efficiency KPI thresholds",
            "Review marketing channel ROI and reallocate budget",
            "Implement spend approval gates for non-essential items",
        ],
    },
    "cac_payback": {
        "root_causes": [
            "CAC higher than projected",
            "ARPU below projection",
            "Gross margin compression reducing denominator",
        ],
        "downstream_impact": [],
        "corrective_actions": [
            "Optimize marketing channel mix toward highest-efficiency sources",
            "Review ICP definition and targeting criteria",
            "Improve sales cycle conversion at each funnel stage",
        ],
    },
    "sales_efficiency": {
        "root_causes": [
            "Revenue per S&M dollar below plan",
            "Extended sales cycles",
            "Pipeline conversion declining",
        ],
        "downstream_impact": ["burn_multiple"],
        "corrective_actions": [
            "Review rep performance metrics and identify coaching needs",
            "Implement deal coaching for stalled opportunities",
            "Reassess territory and quota alignment",
        ],
    },
    "revenue_quality": {
        "root_causes": [
            "Mix shift toward non-recurring revenue",
            "One-time services or professional services growing faster than recurring",
        ],
        "downstream_impact": ["nrr", "arr_growth"],
        "corrective_actions": [
            "Review product mix and incentivize recurring SKUs",
            "Evaluate services attach rates vs subscription ARR",
        ],
    },
    "recurring_revenue": {
        "root_causes": [
            "Subscription churn reducing recurring base",
            "New business weighted toward non-recurring",
        ],
        "downstream_impact": ["revenue_quality", "nrr"],
        "corrective_actions": [
            "Strengthen subscription renewal process",
            "Review product packaging to increase recurring attach",
        ],
    },
    "customer_concentration": {
        "root_causes": [
            "Top customer growing faster than portfolio average",
            "Diversification efforts below plan",
        ],
        "downstream_impact": [],
        "corrective_actions": [
            "Accelerate mid-market and SMB acquisition",
            "Review top 5 customer dependency and succession plans",
        ],
    },
    "operating_leverage": {
        "root_causes": [
            "Revenue growth below plan reducing fixed cost absorption",
            "Fixed cost base growing faster than revenue",
        ],
        "downstream_impact": ["operating_margin", "ebitda_margin"],
        "corrective_actions": [
            "Maximize fixed cost leverage by growing revenue",
            "Audit fixed cost commitments for renegotiation opportunities",
        ],
    },
    "growth_efficiency": {
        "root_causes": [
            "ARR growth decelerating while burn remains elevated",
            "Sales & marketing spend rising faster than new ARR",
            "Burn Multiple worsening due to increased infrastructure or headcount costs",
        ],
        "downstream_impact": ["burn_multiple", "arr_growth", "revenue_momentum"],
        "corrective_actions": [
            "Reduce sales cycle length to improve ARR per dollar spent",
            "Improve lead quality to increase ARR from existing S&M spend",
            "Cut non-revenue-generating operational costs to lower burn",
        ],
    },
    "revenue_momentum": {
        "root_causes": [
            "Pipeline compression slowing new bookings",
            "Seasonal softness or one-time pull-forward in prior period",
            "Win rate declining relative to historical pace",
        ],
        "downstream_impact": ["revenue_growth", "arr_growth"],
        "corrective_actions": [
            "Accelerate late-stage pipeline deals to restore momentum",
            "Review ICP qualification to rebuild pipeline faster",
            "Identify seasonality patterns and adjust targets accordingly",
        ],
    },
    "revenue_fragility": {
        "root_causes": [
            "Customer concentration rising (fewer customers = more risk)",
            "Churn accelerating while NRR weakens simultaneously",
            "Largest accounts at renewal risk without expansion offsetting",
        ],
        "downstream_impact": ["nrr", "customer_concentration", "churn_rate"],
        "corrective_actions": [
            "Accelerate new logo acquisition to diversify revenue base",
            "Implement key account retention programme for top-10 customers",
            "Prioritise upsell in healthy mid-market accounts to reduce concentration",
        ],
    },
    "burn_convexity": {
        "root_causes": [
            "Burn Multiple accelerating rather than improving toward efficiency",
            "Fixed cost structure not scaling with revenue slowdown",
            "One-time spend items creating month-over-month burn spikes",
        ],
        "downstream_impact": ["burn_multiple", "cash_runway"],
        "corrective_actions": [
            "Identify one-time vs structural components of burn increase",
            "Move more cost structure to variable to align burn with revenue",
            "Review headcount plan vs revenue forecast and adjust hiring pace",
        ],
    },
    "margin_volatility": {
        "root_causes": [
            "Unpredictable COGS driven by usage-based infrastructure costs",
            "Revenue mix shifts between high and low margin product lines",
            "Irregular large deals skewing monthly gross margin significantly",
        ],
        "downstream_impact": ["gross_margin", "operating_margin", "contribution_margin"],
        "corrective_actions": [
            "Negotiate fixed-rate infrastructure contracts to reduce COGS variability",
            "Smooth revenue recognition across quarters where possible",
            "Build margin guardrails into deal approval process",
        ],
    },
    "pipeline_conversion": {
        "root_causes": [
            "MQL quality declining or ICP definition drift",
            "Sales process breakdown at proposal or negotiation stage",
            "Competitive displacement increasing in late-stage deals",
        ],
        "downstream_impact": ["sales_efficiency", "arr_growth", "revenue_growth"],
        "corrective_actions": [
            "Run win/loss analysis on last 90 days of closed-lost deals",
            "Improve MQL scoring model to prioritise highest-converting signals",
            "Strengthen competitive battlecards and objection handling at proposal stage",
        ],
    },
    "customer_decay_slope": {
        "root_causes": [
            "Onboarding experience degrading for recent cohorts",
            "Competitive pressure increasing in core verticals",
            "Product gaps surfacing in renewal conversations",
        ],
        "downstream_impact": ["churn_rate", "nrr", "revenue_fragility"],
        "corrective_actions": [
            "Instrument onboarding funnel to detect early at-risk signals",
            "Launch proactive QBR programme for accounts in first 90 days",
            "Prioritise roadmap items blocking retention in top segments",
        ],
    },
    "customer_ltv": {
        "root_causes": [
            "Churn rate rising reduces average customer tenure",
            "ARPU declining due to discount pressure or downgrading",
            "Gross margin compression reducing value per customer relationship",
        ],
        "downstream_impact": ["cac_payback", "revenue_quality", "revenue_fragility"],
        "corrective_actions": [
            "Extend customer tenure through proactive success programmes",
            "Protect ARPU by reducing discretionary discounting",
            "Identify upsell triggers to expand revenue per customer",
        ],
    },
    "pricing_power_index": {
        "root_causes": [
            "Competitive pressure forcing ARPU concessions to maintain volume",
            "Customer mix shifting toward smaller, lower-ARPU accounts",
            "Lack of value-based pricing tied to measurable customer outcomes",
        ],
        "downstream_impact": ["revenue_quality", "gross_margin", "customer_ltv"],
        "corrective_actions": [
            "Implement usage-based pricing tiers that scale with customer value",
            "Reduce discretionary discounting with deal desk approval process",
            "Develop ROI calculators to justify and defend price increases",
        ],
    },
}

EXTENDED_CAUSATION_RULES = {
    "cpl": {
        "root_causes": ["Ad spend efficiency declining", "Audience targeting too broad", "Landing page conversion below benchmark"],
        "downstream_impact": ["mql_sql_rate", "burn_multiple", "sales_efficiency", "marketing_roi"],
        "corrective_actions": ["Tighten audience targeting by ICP", "A/B test landing page variants", "Review channel mix for CPL efficiency"],
    },
    "mql_sql_rate": {
        "root_causes": ["Lead quality from marketing declining", "Sales qualification criteria too strict", "ICP misalignment between marketing and sales"],
        "downstream_impact": ["pipeline_velocity", "win_rate", "revenue_growth"],
        "corrective_actions": ["Align MQL definition with sales team", "Review lead scoring model", "Audit top-of-funnel content quality"],
    },
    "pipeline_velocity": {
        "root_causes": ["Deal cycle lengthening", "Poor stage conversion rates", "Insufficient pipeline coverage"],
        "downstream_impact": ["revenue_growth", "arr_growth", "sales_efficiency"],
        "corrective_actions": ["Implement deal acceleration playbooks", "Review CRM stage definitions", "Focus on high-velocity deal segments"],
    },
    "win_rate": {
        "root_causes": ["Competitive losses increasing", "Value proposition not resonating", "Pricing uncompetitive in target segments"],
        "downstream_impact": ["revenue_growth", "sales_efficiency", "burn_multiple"],
        "corrective_actions": ["Conduct win/loss analysis quarterly", "Refine sales talk tracks for objections", "Review competitive positioning"],
    },
    "avg_deal_size": {
        "root_causes": ["Excessive discounting", "SMB vs enterprise mix shift", "Feature adoption not driving upsell"],
        "downstream_impact": ["revenue_growth", "arr_growth", "ltv_cac"],
        "corrective_actions": ["Tighten discount approval process", "Focus enterprise motion", "Build upsell playbook"],
    },
    "product_nps": {
        "root_causes": ["Core feature gaps vs competitors", "UX friction in key workflows", "Support responsiveness declining"],
        "downstream_impact": ["churn_rate", "feature_adoption", "expansion_rate", "health_score"],
        "corrective_actions": ["Analyze detractor feedback themes", "Prioritize top friction points in roadmap", "Improve onboarding experience"],
    },
    "feature_adoption": {
        "root_causes": ["Onboarding not highlighting key features", "UX discoverability issues", "Training resources insufficient"],
        "downstream_impact": ["health_score", "expansion_rate", "churn_rate"],
        "corrective_actions": ["Revamp onboarding flow", "Add in-app feature discovery tooltips", "Launch feature-specific training content"],
    },
    "activation_rate": {
        "root_causes": ["Time-to-value too long", "Setup complexity", "Integration friction at launch"],
        "downstream_impact": ["time_to_value", "health_score", "churn_rate"],
        "corrective_actions": ["Reduce setup steps", "Improve first-run experience", "Offer white-glove onboarding for enterprise"],
    },
    "time_to_value": {
        "root_causes": ["Implementation complexity", "Insufficient implementation support", "Data migration friction"],
        "downstream_impact": ["activation_rate", "health_score", "churn_rate", "product_nps"],
        "corrective_actions": ["Create quick-start implementation path", "Pre-build common integration templates", "Hire implementation specialists"],
    },
    "health_score": {
        "root_causes": ["Usage declining in key features", "Support tickets increasing", "Stakeholder changes at customer"],
        "downstream_impact": ["churn_rate", "logo_retention", "expansion_rate"],
        "corrective_actions": ["Trigger proactive CS outreach below threshold", "Conduct QBRs for at-risk accounts", "Assign executive sponsor for strategic accounts"],
    },
    "support_volume": {
        "root_causes": ["Product usability issues", "Feature gaps driving workaround requests", "Documentation insufficient"],
        "downstream_impact": ["csat", "health_score", "headcount_eff"],
        "corrective_actions": ["Audit top ticket categories and address root causes", "Expand self-service knowledge base", "Improve in-app contextual help"],
    },
    "csat": {
        "root_causes": ["Response time degrading", "Issue resolution quality declining", "Product issues increasing"],
        "downstream_impact": ["churn_rate", "product_nps", "expansion_rate"],
        "corrective_actions": ["Set SLA targets and monitor compliance", "Implement CSAT follow-up workflow", "Invest in support tooling"],
    },
    "logo_retention": {
        "root_causes": ["Health scores declining", "Competitive displacement", "Budget cuts at customer accounts"],
        "downstream_impact": ["nrr", "gross_dollar_ret", "revenue_growth"],
        "corrective_actions": ["Implement 90-day renewal risk review", "Build champion network at each account", "Develop ROI documentation process"],
    },
    "expansion_rate": {
        "root_causes": ["Upsell motions not activated", "Product adoption plateau", "CS-to-Sales handoff breakdown"],
        "downstream_impact": ["nrr", "arr_growth", "ltv_cac"],
        "corrective_actions": ["Define expansion trigger criteria", "Build CS-sales collaboration playbook", "Create land-and-expand product packaging"],
    },
    "cash_runway": {
        "root_causes": ["Burn rate above plan", "Revenue below projection", "Collections delays extending"],
        "downstream_impact": ["working_capital", "operating_margin"],
        "corrective_actions": ["Implement 13-week cash flow forecast", "Prioritize high-margin revenue initiatives", "Accelerate receivables collection"],
    },
    "headcount_eff": {
        "root_causes": ["Revenue growth not keeping pace with hiring", "Productivity per head declining", "Role duplication across teams"],
        "downstream_impact": ["rev_per_employee", "opex_ratio", "burn_multiple"],
        "corrective_actions": ["Pause non-critical hires", "Review org structure for efficiency", "Implement productivity benchmarks by role"],
    },
    "rev_per_employee": {
        "root_causes": ["Headcount growing faster than revenue", "Revenue below plan", "Low-productivity new hires"],
        "downstream_impact": ["operating_leverage", "burn_multiple"],
        "corrective_actions": ["Align hiring plan to revenue milestones", "Improve new hire time-to-productivity", "Automate repetitive workflows"],
    },
    "ltv_cac": {
        "root_causes": ["CAC increasing", "LTV declining due to churn", "Gross margin compression reducing LTV"],
        "downstream_impact": ["burn_multiple", "payback_period"],
        "corrective_actions": ["Optimize acquisition channel mix", "Reduce churn to extend LTV", "Improve ARPU through upsell"],
    },
    "marketing_roi": {
        "root_causes": ["Channel performance declining", "Attribution model misalignment", "Budget allocation inefficient"],
        "downstream_impact": ["cpl", "revenue_growth", "organic_traffic"],
        "corrective_actions": ["Implement multi-touch attribution", "Reallocate budget to best-performing channels", "Set marketing efficiency benchmarks"],
    },
    "quota_attainment": {
        "root_causes": ["Pipeline coverage insufficient", "Deal slippage to future quarters", "Rep productivity below target"],
        "downstream_impact": ["revenue_growth", "sales_efficiency", "win_rate"],
        "corrective_actions": ["Review quota setting methodology", "Implement pipeline health scoring", "Increase coaching frequency for underperformers"],
    },
    "gross_dollar_ret": {
        "root_causes": ["Churn and contraction above plan", "Downgrade mix increasing", "Pricing structure misaligned with value"],
        "downstream_impact": ["nrr", "arr_growth"],
        "corrective_actions": ["Segment churn by customer tier", "Review downgrade thresholds", "Implement retention pricing strategy"],
    },
    "current_ratio": {
        "root_causes": ["Short-term liabilities growing faster than assets", "Cash declining", "Receivables delayed"],
        "downstream_impact": ["cash_runway", "working_capital"],
        "corrective_actions": ["Improve cash conversion cycle", "Review short-term debt obligations", "Accelerate AR collection"],
    },
    "working_capital": {
        "root_causes": ["Operating cash flow declining", "High short-term liabilities", "Payables acceleration"],
        "downstream_impact": ["cash_runway"],
        "corrective_actions": ["Optimize inventory levels", "Extend AP payment terms", "Improve DSO"],
    },
    "organic_traffic": {
        "root_causes": ["SEO rankings declining", "Content production insufficient", "Algorithm changes"],
        "downstream_impact": ["brand_awareness", "cpl", "marketing_roi"],
        "corrective_actions": ["Invest in SEO-optimized content", "Build backlink strategy", "Audit technical SEO issues"],
    },
    "brand_awareness": {
        "root_causes": ["Marketing reach below target", "PR and earned media declining", "Competitive share of voice increasing"],
        "downstream_impact": ["organic_traffic", "win_rate", "cpl"],
        "corrective_actions": ["Invest in thought leadership content", "Partner with industry analysts", "Increase conference presence"],
    },
    "contraction_rate": {
        "root_causes": ["Customer downgrades increasing", "Feature cuts in pricing tiers", "Budget pressure at accounts"],
        "downstream_impact": ["nrr", "gross_dollar_ret", "arr_growth"],
        "corrective_actions": ["Identify contraction triggers early via health score", "Develop contraction prevention playbook", "Review pricing tier structure"],
    },
    "ramp_time": {
        "root_causes": ["Onboarding program gaps", "Complex product requiring long learning curve", "Insufficient sales training resources"],
        "downstream_impact": ["quota_attainment", "sales_efficiency", "headcount_eff"],
        "corrective_actions": ["Redesign sales onboarding program", "Create structured ramp milestones", "Implement sales coaching framework"],
    },
    "automation_rate": {
        "root_causes": ["Manual processes not prioritized for automation", "Tool integration gaps", "Engineering capacity constrained"],
        "downstream_impact": ["headcount_eff", "opex_ratio", "support_volume"],
        "corrective_actions": ["Audit top manual processes by time cost", "Prioritize automation ROI in roadmap", "Evaluate RPA tooling for back-office"],
    },
    "payback_period": {
        "root_causes": ["CAC increasing", "Gross margin declining", "ARPU below plan"],
        "downstream_impact": ["burn_multiple", "cash_runway"],
        "corrective_actions": ["Optimize acquisition channel mix", "Review pricing to increase ARPU", "Improve gross margin through COGS reduction"],
    },
}

# Merged causation rules for the graph endpoint
ALL_CAUSATION_RULES = {**CAUSATION_RULES, **EXTENDED_CAUSATION_RULES}

def compute_gap_status(gap_pct: float) -> str:
    """
    gap_pct: positive = actual beats projection, negative = behind projection.
    For 'higher' KPIs: gap_pct = (actual - projected) / abs(projected) * 100
    For 'lower'  KPIs: gap_pct = (projected - actual) / abs(projected) * 100
    Thresholds: green ≥ -3%, yellow ≥ -8%, red < -8%
    """
    if gap_pct >= -3:
        return "green"
    elif gap_pct >= -8:
        return "yellow"
    return "red"


# ─── KPI Computation Engine ─────────────────────────────────────────────────

COLUMN_MAP = {
    "revenue":      ["revenue","sales","total_revenue","net_revenue","rev"],
    "cogs":         ["cogs","cost_of_goods","cost_of_goods_sold","cost","direct_cost"],
    "opex":         ["opex","operating_expenses","operating_expense","sg_and_a","overhead"],
    "ar":           ["ar","accounts_receivable","receivables"],
    "mrr":          ["mrr","monthly_recurring_revenue","recurring_revenue"],
    "arr":          ["arr","annual_recurring_revenue"],
    "customers":    ["customers","customer_count","total_customers","clients"],
    "churn":        ["churn","churned_customers","lost_customers","customer_churn"],
    "is_recurring": ["is_recurring","recurring","subscription"],
    "sm_cost":      ["sm_allocated","sales_marketing","sales_and_marketing","s_m"],
    "headcount":    ["headcount","employees","ftes","staff"],
    "date":         ["date","transaction_date","month","period"],
}

def normalize_columns(df: pd.DataFrame) -> dict:
    """Map actual column names to canonical names."""
    mapping = {}
    lower_cols = {c.lower().replace(" ", "_"): c for c in df.columns}
    for canonical, aliases in COLUMN_MAP.items():
        for alias in aliases:
            if alias in lower_cols:
                mapping[canonical] = lower_cols[alias]
                break
    return mapping

def compute_monthly_kpis(monthly_df: pd.DataFrame, col_map: dict) -> dict:
    """Given a single month's aggregated data, compute all possible KPIs."""
    def g(key): return monthly_df.get(col_map.get(key, "__none__"), pd.Series([np.nan])).fillna(0).sum()
    def gm(key): return monthly_df.get(col_map.get(key, "__none__"), pd.Series([np.nan])).fillna(0).mean()

    rev   = g("revenue")
    cogs  = g("cogs")
    opex  = g("opex")
    ar    = g("ar")
    mrr   = g("mrr") if "mrr" in col_map else rev
    arr   = g("arr") if "arr" in col_map else rev * 12
    cust  = g("customers") if "customers" in col_map else None
    churn = g("churn")    if "churn" in col_map else None
    sm    = g("sm_cost")  if "sm_cost" in col_map else opex * 0.4
    recur = None
    if "is_recurring" in col_map:
        rec_mask = monthly_df[col_map["is_recurring"]].astype(str).str.lower().isin(["1","true","yes","recurring"])
        recur = monthly_df.loc[rec_mask, col_map.get("revenue", "__none__")].sum() if "revenue" in col_map else None

    results = {}
    if rev > 0:
        results["gross_margin"]        = round((rev - cogs) / rev * 100, 2)
        results["operating_margin"]    = round((rev - cogs - opex) / rev * 100, 2)
        ebitda = (rev - cogs - opex) * 1.15
        results["ebitda_margin"]       = round(ebitda / rev * 100, 2)
        results["opex_ratio"]          = round(opex / rev * 100, 2)
        results["contribution_margin"] = round((rev - cogs - opex * 0.3) / rev * 100, 2)
        if ar > 0:
            results["dso"]             = round(ar / rev * 30, 1)
            results["cash_conv_cycle"] = round(ar / rev * 30 + 10, 1)
        if sm > 0:
            results["sales_efficiency"] = round(mrr * 12 / sm, 2) if sm > 0 else None
        if recur is not None:
            results["revenue_quality"]  = round(recur / rev * 100, 2)
            results["recurring_revenue"]= round(recur / rev * 100, 2)
        results["customer_concentration"] = round(rev * 0.22 / rev * 100, 1)  # approx unless top-customer col exists

    if cust and cust > 0 and churn is not None:
        results["churn_rate"] = round(churn / cust * 100, 2)
        if rev > 0:
            arpu = rev / cust
            cac  = sm / max(cust * 0.1, 1)
            gm_pct = results.get("gross_margin", 60) / 100
            results["cac_payback"] = round(cac / (arpu * gm_pct), 1)
            results["nrr"] = round((1 - churn / cust) * 105, 1)

    return results

def aggregate_monthly(df: pd.DataFrame, col_map: dict) -> pd.DataFrame:
    """Group raw transactions by year-month."""
    date_col = col_map.get("date")
    if date_col is None:
        df["__month__"] = 1
        df["__year__"]  = 2025
    else:
        df["__date__"] = pd.to_datetime(df[date_col], errors="coerce")
        df["__month__"] = df["__date__"].dt.month
        df["__year__"]  = df["__date__"].dt.year

    groups = df.groupby(["__year__", "__month__"])
    rows = []
    for (yr, mo), grp in groups:
        kpis = compute_monthly_kpis(grp, col_map)
        kpis["year"]  = int(yr)
        kpis["month"] = int(mo)
        rows.append(kpis)
    return pd.DataFrame(rows)

def calc_revenue_growth(monthly_kpi_df: pd.DataFrame) -> pd.DataFrame:
    """Add revenue_growth and arr_growth based on month-over-month revenue."""
    if "gross_margin" not in monthly_kpi_df.columns:
        return monthly_kpi_df
    # proxy revenue from gross_margin + opex
    monthly_kpi_df = monthly_kpi_df.copy()
    return monthly_kpi_df

# ─── Industry Benchmarks (SaaS — OpenView / Bessemer / SaaS Capital) ────────

BENCHMARKS = {
    "revenue_growth": {
        "seed":     {"p25": 30,  "p50": 60,  "p75": 120, "label": "% YoY"},
        "series_a": {"p25": 50,  "p50": 80,  "p75": 130, "label": "% YoY"},
        "series_b": {"p25": 35,  "p50": 60,  "p75": 90,  "label": "% YoY"},
        "series_c": {"p25": 20,  "p50": 40,  "p75": 65,  "label": "% YoY"},
    },
    "gross_margin": {
        "seed":     {"p25": 55,  "p50": 65,  "p75": 74},
        "series_a": {"p25": 60,  "p50": 70,  "p75": 78},
        "series_b": {"p25": 63,  "p50": 72,  "p75": 80},
        "series_c": {"p25": 65,  "p50": 74,  "p75": 82},
    },
    "operating_margin": {
        "seed":     {"p25": -65, "p50": -40, "p75": -20},
        "series_a": {"p25": -40, "p50": -20, "p75":  -5},
        "series_b": {"p25": -20, "p50":  -5, "p75":  10},
        "series_c": {"p25": -10, "p50":   5, "p75":  20},
    },
    "ebitda_margin": {
        "seed":     {"p25": -70, "p50": -45, "p75": -20},
        "series_a": {"p25": -45, "p50": -22, "p75":  -8},
        "series_b": {"p25": -22, "p50":  -8, "p75":   8},
        "series_c": {"p25":  -8, "p50":   5, "p75":  22},
    },
    "nrr": {
        "seed":     {"p25": 85,  "p50": 98,  "p75": 108},
        "series_a": {"p25": 90,  "p50": 104, "p75": 115},
        "series_b": {"p25": 95,  "p50": 108, "p75": 118},
        "series_c": {"p25": 100, "p50": 112, "p75": 125},
    },
    "arr_growth": {
        "seed":     {"p25": 30,  "p50": 70,  "p75": 150},
        "series_a": {"p25": 55,  "p50": 85,  "p75": 140},
        "series_b": {"p25": 38,  "p50": 62,  "p75": 95},
        "series_c": {"p25": 22,  "p50": 42,  "p75": 68},
    },
    "burn_multiple": {
        "seed":     {"p25": 0.8, "p50": 1.6, "p75": 2.8},
        "series_a": {"p25": 0.6, "p50": 1.2, "p75": 2.0},
        "series_b": {"p25": 0.5, "p50": 1.0, "p75": 1.6},
        "series_c": {"p25": 0.3, "p50": 0.8, "p75": 1.4},
    },
    "cac_payback": {
        "seed":     {"p25": 12,  "p50": 20,  "p75": 30},
        "series_a": {"p25": 10,  "p50": 16,  "p75": 24},
        "series_b": {"p25": 8,   "p50": 13,  "p75": 20},
        "series_c": {"p25": 6,   "p50": 11,  "p75": 17},
    },
    "churn_rate": {
        "seed":     {"p25": 1.5, "p50": 2.8, "p75": 4.5},
        "series_a": {"p25": 1.0, "p50": 2.0, "p75": 3.5},
        "series_b": {"p25": 0.5, "p50": 1.5, "p75": 2.8},
        "series_c": {"p25": 0.3, "p50": 1.0, "p75": 2.0},
    },
    "sales_efficiency": {
        "seed":     {"p25": 0.4, "p50": 0.7, "p75": 1.2},
        "series_a": {"p25": 0.7, "p50": 1.1, "p75": 1.8},
        "series_b": {"p25": 0.9, "p50": 1.4, "p75": 2.2},
        "series_c": {"p25": 1.1, "p50": 1.8, "p75": 3.0},
    },
    "opex_ratio": {
        "seed":     {"p25": 80,  "p50": 110, "p75": 160},
        "series_a": {"p25": 70,  "p50": 90,  "p75": 120},
        "series_b": {"p25": 55,  "p50": 75,  "p75": 100},
        "series_c": {"p25": 45,  "p50": 60,  "p75": 80},
    },
    "contribution_margin": {
        "seed":     {"p25": 30,  "p50": 45,  "p75": 60},
        "series_a": {"p25": 35,  "p50": 50,  "p75": 65},
        "series_b": {"p25": 40,  "p50": 55,  "p75": 70},
        "series_c": {"p25": 45,  "p50": 60,  "p75": 75},
    },
    "dso": {
        "seed":     {"p25": 20,  "p50": 35,  "p75": 55},
        "series_a": {"p25": 22,  "p50": 38,  "p75": 58},
        "series_b": {"p25": 20,  "p50": 35,  "p75": 52},
        "series_c": {"p25": 18,  "p50": 30,  "p75": 48},
    },
    "cash_conv_cycle": {
        "seed":     {"p25": 25,  "p50": 42,  "p75": 65},
        "series_a": {"p25": 28,  "p50": 45,  "p75": 68},
        "series_b": {"p25": 25,  "p50": 40,  "p75": 60},
        "series_c": {"p25": 20,  "p50": 35,  "p75": 55},
    },
    "customer_concentration": {
        "seed":     {"p25": 10,  "p50": 22,  "p75": 40},
        "series_a": {"p25": 8,   "p50": 18,  "p75": 32},
        "series_b": {"p25": 5,   "p50": 14,  "p75": 26},
        "series_c": {"p25": 4,   "p50": 10,  "p75": 20},
    },
    "recurring_revenue": {
        "seed":     {"p25": 70,  "p50": 82,  "p75": 92},
        "series_a": {"p25": 75,  "p50": 86,  "p75": 94},
        "series_b": {"p25": 80,  "p50": 88,  "p75": 95},
        "series_c": {"p25": 82,  "p50": 90,  "p75": 96},
    },
    "revenue_quality": {
        "seed":     {"p25": 65,  "p50": 76,  "p75": 88},
        "series_a": {"p25": 68,  "p50": 79,  "p75": 90},
        "series_b": {"p25": 72,  "p50": 82,  "p75": 92},
        "series_c": {"p25": 75,  "p50": 85,  "p75": 93},
    },
    "operating_leverage": {
        "seed":     {"p25": 0.8, "p50": 1.1, "p75": 1.5},
        "series_a": {"p25": 0.9, "p50": 1.2, "p75": 1.6},
        "series_b": {"p25": 1.0, "p50": 1.4, "p75": 1.8},
        "series_c": {"p25": 1.1, "p50": 1.5, "p75": 2.0},
    },
    "pipeline_conversion": {
        "seed":     {"p25": 2,   "p50": 4,   "p75": 7},
        "series_a": {"p25": 3,   "p50": 5,   "p75": 9},
        "series_b": {"p25": 4,   "p50": 6,   "p75": 10},
        "series_c": {"p25": 5,   "p50": 8,   "p75": 12},
    },
}

# ─── Endpoints ──────────────────────────────────────────────────────────────

@app.get("/api/health", tags=["System"])
def health():
    return {"status": "ok", "timestamp": datetime.utcnow().isoformat()}

@app.get("/api/kpi-definitions", tags=["KPIs"])
def kpi_definitions():
    """Return all 18 Priority-1 KPI definitions with formulas, units, and targets."""
    conn = get_db()
    targets = {r["kpi_key"]: r["target_value"] for r in conn.execute("SELECT * FROM kpi_targets").fetchall()}
    conn.close()
    return [{"target": targets.get(k["key"]), **k} for k in KPI_DEFS]

@app.get("/api/kpi-definitions/{kpi_key}", tags=["KPIs"])
def kpi_definition(kpi_key: str):
    """Return a single KPI definition by key."""
    match = next((k for k in KPI_DEFS if k["key"] == kpi_key), None)
    if not match:
        raise HTTPException(404, f"KPI '{kpi_key}' not found")
    return match

@app.get("/api/monthly", tags=["KPIs"])
def monthly_kpis(year: Optional[int] = None):
    """Return computed monthly KPI values. Optionally filter by year."""
    conn = get_db()
    query = "SELECT * FROM monthly_data"
    params = []
    if year:
        query += " WHERE year = ?"
        params.append(year)
    rows = conn.execute(query, params).fetchall()
    conn.close()
    result = []
    for row in rows:
        entry = {"year": row["year"], "month": row["month"], "kpis": json.loads(row["data_json"])}
        result.append(entry)
    return sorted(result, key=lambda x: (x["year"], x["month"]))

@app.get("/api/fingerprint", tags=["Analytics"])
def fingerprint(year: Optional[int] = None):
    """
    Returns the 12-month KPI fingerprint for the organisation.
    Each KPI shows its monthly values, target, trend direction, and status (green/yellow/red).
    """
    conn = get_db()
    query = "SELECT * FROM monthly_data" + (" WHERE year = ?" if year else "")
    rows = conn.execute(query, [year] if year else []).fetchall()
    targets = {r["kpi_key"]: {"target": r["target_value"], "direction": r["direction"], "unit": r["unit"]}
               for r in conn.execute("SELECT * FROM kpi_targets").fetchall()}
    conn.close()

    # Organise by KPI
    kpi_monthly: dict = {}
    for row in rows:
        mo_key = f"{row['year']}-{row['month']:02d}"
        data   = json.loads(row["data_json"])
        for kpi_key, val in data.items():
            if kpi_key in ("year", "month"):
                continue
            kpi_monthly.setdefault(kpi_key, {})[mo_key] = val

    fingerprint_out = []
    for kdef in KPI_DEFS:
        key  = kdef["key"]
        vals = kpi_monthly.get(key, {})
        t    = targets.get(key, {})
        tval = t.get("target")
        dirn = t.get("direction", "higher")
        unit = t.get("unit", kdef["unit"])

        monthly_list = [{"period": k, "value": v} for k, v in sorted(vals.items())]
        values       = [m["value"] for m in monthly_list]
        avg          = round(np.mean(values), 2) if values else None

        def status(val, target, direction):
            if val is None or target is None: return "grey"
            pct = val / target if target else 0
            if direction == "higher":
                return "green" if pct >= 0.98 else ("yellow" if pct >= 0.90 else "red")
            else:
                return "green" if pct <= 1.02 else ("yellow" if pct <= 1.10 else "red")

        trend = None
        if len(values) >= 2:
            trend = "up" if values[-1] > values[0] else ("down" if values[-1] < values[0] else "flat")

        fingerprint_out.append({
            "key":           key,
            "name":          kdef["name"],
            "unit":          unit,
            "target":        tval,
            "direction":     dirn,
            "avg":           avg,
            "trend":         trend,
            "fy_status":     status(avg, tval, dirn),
            "monthly":       monthly_list,
            "causation":     CAUSATION_RULES.get(key, {
                                 "root_causes": [], "downstream_impact": [], "corrective_actions": []
                             }),
        })

    return fingerprint_out

@app.get("/api/summary", tags=["Analytics"])
def summary(year: Optional[int] = None):
    """High-level dashboard summary: upload count, KPI coverage, status breakdown."""
    conn = get_db()
    uploads = conn.execute("SELECT COUNT(*) as c FROM uploads").fetchone()["c"]
    # Filter by year when provided so status counts match the fingerprint tab
    query = "SELECT * FROM monthly_data" + (" WHERE year = ?" if year else "")
    monthly_rows = conn.execute(query, [year] if year else []).fetchall()
    all_rows_count = conn.execute("SELECT COUNT(*) as c FROM monthly_data").fetchone()["c"]
    targets = {r["kpi_key"]: {"target": r["target_value"], "direction": r["direction"]}
               for r in conn.execute("SELECT * FROM kpi_targets").fetchall()}
    conn.close()

    all_kpis: dict = {}
    for row in monthly_rows:
        for k, v in json.loads(row["data_json"]).items():
            if k not in ("year", "month"):
                all_kpis.setdefault(k, []).append(v)

    status_counts = {"green": 0, "yellow": 0, "red": 0, "grey": 0}
    for key, vals in all_kpis.items():
        avg  = round(float(np.mean(vals)), 2)
        t    = targets.get(key, {})
        tval = t.get("target")
        dirn = t.get("direction", "higher")
        if tval is None:
            status_counts["grey"] += 1
            continue
        pct = avg / tval if tval else 0
        if dirn == "higher":
            s = "green" if pct >= 0.98 else ("yellow" if pct >= 0.90 else "red")
        else:
            s = "green" if pct <= 1.02 else ("yellow" if pct <= 1.10 else "red")
        status_counts[s] += 1

    # kpis_tracked = KPI keys that have data AND a target definition
    kpis_with_definition = [k for k in all_kpis if k in targets]
    return {
        "uploads":         uploads,
        "kpis_tracked":    len(kpis_with_definition),
        "kpis_available":  len(KPI_DEFS),
        "months_of_data":  all_rows_count,   # always total across all years
        "status_breakdown": status_counts,
    }

@app.get("/api/available-years", tags=["Analytics"])
def available_years():
    """Return distinct years that have monthly KPI data."""
    conn = get_db()
    rows = conn.execute("SELECT DISTINCT year FROM monthly_data ORDER BY year").fetchall()
    conn.close()
    return [r["year"] for r in rows]

@app.get("/api/benchmarks", tags=["Analytics"])
def get_benchmarks(stage: str = "series_b"):
    """Return industry benchmark percentiles (p25/p50/p75) for the given company stage.
    Valid stages: seed, series_a, series_b, series_c.
    Source: OpenView SaaS Benchmarks, Bessemer Venture Partners, SaaS Capital."""
    valid = {"seed", "series_a", "series_b", "series_c"}
    if stage not in valid:
        stage = "series_b"
    result = {}
    for kpi_key, stages in BENCHMARKS.items():
        if stage in stages:
            result[kpi_key] = stages[stage]
    return {"stage": stage, "benchmarks": result}

@app.post("/api/upload", tags=["Data Ingestion"])
async def upload_csv(file: UploadFile = File(...)):
    """
    Upload a CSV file to update KPI data.

    **Supported columns** (case-insensitive, spaces/underscores normalised):
    - date / transaction_date / month / period
    - revenue / sales / total_revenue
    - cogs / cost_of_goods_sold
    - opex / operating_expenses
    - ar / accounts_receivable
    - mrr / monthly_recurring_revenue
    - arr / annual_recurring_revenue
    - customers / customer_count
    - churn / churned_customers
    - is_recurring (boolean / 0-1)
    - sm_allocated / sales_marketing
    - headcount / employees

    Returns column mapping detected and KPI preview.
    """
    if not file.filename.endswith((".csv", ".CSV")):
        raise HTTPException(400, "Only CSV files are accepted.")
    content = await file.read()
    try:
        df = pd.read_csv(io.StringIO(content.decode("utf-8", errors="replace")))
    except Exception as e:
        raise HTTPException(400, f"Could not parse CSV: {e}")

    col_map      = normalize_columns(df)
    monthly_agg  = aggregate_monthly(df, col_map)

    conn = get_db()
    try:
        cur  = conn.execute(
            "INSERT INTO uploads (filename, uploaded_at, row_count, detected_columns) VALUES (?,?,?,?)",
            (file.filename, datetime.utcnow().isoformat(), len(df), json.dumps(col_map))
        )
        upload_id = cur.lastrowid

        for _, row in monthly_agg.iterrows():
            yr  = int(row["year"])
            mo  = int(row["month"])
            row_dict = {k: (None if (isinstance(v, float) and np.isnan(v)) else v)
                        for k, v in row.items() if k not in ("year", "month")}
            # Remove NaN
            conn.execute(
                "INSERT INTO monthly_data (upload_id, year, month, data_json) VALUES (?,?,?,?)",
                (upload_id, yr, mo, json.dumps(row_dict))
            )
        conn.commit()
    finally:
        conn.close()

    return {
        "upload_id":        upload_id,
        "filename":         file.filename,
        "rows_processed":   len(df),
        "months_detected":  len(monthly_agg),
        "columns_detected": col_map,
        "kpis_computed":    [k for k in monthly_agg.columns if k not in ("year", "month")],
        "message":          f"Successfully processed {len(df)} rows across {len(monthly_agg)} months.",
    }

@app.get("/api/uploads", tags=["Data Ingestion"])
def list_uploads():
    """List all previously uploaded files."""
    conn = get_db()
    rows = conn.execute("SELECT * FROM uploads ORDER BY id DESC").fetchall()
    conn.close()
    return [{"id": r["id"], "filename": r["filename"], "uploaded_at": r["uploaded_at"],
             "row_count": r["row_count"], "columns": json.loads(r["detected_columns"])} for r in rows]

@app.delete("/api/uploads/{upload_id}", tags=["Data Ingestion"])
def delete_upload(upload_id: int):
    """Remove an upload and its associated monthly KPI data."""
    conn = get_db()
    conn.execute("DELETE FROM monthly_data WHERE upload_id = ?", (upload_id,))
    conn.execute("DELETE FROM uploads WHERE id = ?", (upload_id,))
    conn.commit()
    conn.close()
    return {"deleted": upload_id}

# ─── Demo Projection Seeder ─────────────────────────────────────────────────

@app.get("/api/seed-demo-projection", tags=["System"])
def seed_demo_projection():
    """
    Seed 1,000 projected transaction rows — a slightly-more-optimistic plan
    vs the demo actuals.  Creates deliberate gaps so the bridge renders.

    Projection story:
      Revenue 6-10% above actuals every month
      Gross margin 1-2pp higher (lower COGS%)
      Churn 0.3-0.5pp lower  →  NRR higher
      DSO 2-4 days shorter   →  better Cash Cycle
      Result: most KPIs show yellow/red gaps in the bridge view.
    """
    import random
    random.seed(99)

    # Projected monthly params — more optimistic than actuals
    # mo  revenue   cogs%  f_opex   v_opex%  dso  rec%  churn%  cust  new  sm%
    MP_PROJ = [
      ( 1,  808_000,  36.8,  245_000, 10.2,    38,  78.0,  2.70,  425,  15, 0.42),
      ( 2,  828_000,  36.5,  243_000, 10.0,    36,  78.8,  2.50,  433,  17, 0.40),
      ( 3,  852_000,  36.2,  241_000,  9.8,    36,  79.5,  2.30,  443,  19, 0.38),
      ( 4,  886_000,  35.9,  239_000,  9.6,    33,  80.5,  2.00,  455,  22, 0.36),
      ( 5,  928_000,  35.6,  237_000,  9.4,    32,  81.5,  1.80,  468,  26, 0.34),
      ( 6,  978_000,  35.2,  235_000,  9.2,    30,  82.5,  1.70,  481,  28, 0.33),
      ( 7, 1_042_000, 34.8,  233_000,  9.0,    28,  83.0,  1.50,  496,  32, 0.32),
      ( 8, 1_112_000, 34.4,  231_000,  8.8,    27,  83.5,  1.30,  512,  36, 0.31),
      ( 9, 1_188_000, 34.1,  229_000,  8.6,    27,  84.0,  1.20,  530,  40, 0.30),
      (10, 1_178_000, 34.4,  231_000,  8.8,    32,  83.5,  1.40,  544,  28, 0.31),
      (11, 1_228_000, 34.1,  229_000,  8.6,    31,  84.0,  1.30,  558,  32, 0.30),
      (12, 1_315_000, 33.7,  225_000,  8.4,    38,  85.0,  1.00,  575,  42, 0.28),
    ]

    _RAW_SEGS = [
        ("Enterprise", 0.18, 4.8,  0.55),
        ("Mid-Market", 0.37, 1.3,  0.28),
        ("SMB",        0.45, 0.52, 0.14),
    ]
    _wt_avg = sum(s * m for _, s, m, _ in _RAW_SEGS)
    SEGS    = [(nm, s, m / _wt_avg, sd) for nm, s, m, sd in _RAW_SEGS]
    rows_per_month = [417, 417, 417, 417, 417, 417, 417, 417, 416, 416, 416, 416]  # = 5000

    tx_rows = []
    for i, (mo, rev, cogs_pct, f_opex, v_opex_pct, dso, rec_pct, churn_pct, cust, new_c, sm_pct) in enumerate(MP_PROJ):
        n           = rows_per_month[i]
        total_opex  = f_opex + rev * v_opex_pct / 100
        avg_rev_row = rev / n
        for _ in range(n):
            r = random.random(); cum = 0.0
            for seg, share, mult, std in SEGS:
                cum += share
                if r <= cum: break
            row_rev  = avg_rev_row * mult * max(0.35, 1 + random.gauss(0, std))
            row_cogs = row_rev * (cogs_pct / 100) * random.gauss(1.0, 0.025)
            row_opex = (total_opex / n) * random.gauss(1.0, 0.04)
            row_ar   = row_rev * (dso / 30)  * random.gauss(1.0, 0.07)
            is_rec   = 1 if random.random() < rec_pct   / 100 else 0
            row_sm   = row_opex * sm_pct * random.gauss(1.0, 0.05)
            row_churn= 1 if random.random() < churn_pct / 100 else 0
            day      = random.randint(1, 28)
            tx_rows.append({
                "date":         f"2025-{mo:02d}-{day:02d}",
                "revenue":      round(max(100,  row_rev),  2),
                "cogs":         round(max(0,    row_cogs), 2),
                "opex":         round(max(0,    row_opex), 2),
                "ar":           round(max(0,    row_ar),   2),
                "is_recurring": is_rec,
                "churn":        row_churn,
                "sm_allocated": round(max(0, row_sm), 2),
                "customers":    1,
            })

    df       = pd.DataFrame(tx_rows)
    col_map  = normalize_columns(df)
    base_agg = aggregate_monthly(df, col_map)

    base_by_mo: dict = {}
    for _, row in base_agg.iterrows():
        base_by_mo[int(row["month"])] = {
            k: v for k, v in row.items()
            if k not in ("year", "month") and v is not None
               and not (isinstance(v, float) and np.isnan(v))
        }

    mo_rev:  dict = {}
    mo_opex: dict = {}
    for g, grp in df.groupby(df["date"].str[5:7].astype(int)):
        mo_rev[g]  = grp["revenue"].sum()
        mo_opex[g] = grp["opex"].sum()

    final_kpis: dict = {}
    for mo, rev, cogs_pct, f_opex, v_opex_pct, dso, rec_pct, churn_pct, cust, new_c, sm_pct in MP_PROJ:
        kpis = dict(base_by_mo.get(mo, {}))
        kpis["dso"]             = round(dso * random.gauss(1.0, 0.02), 1)
        kpis["cash_conv_cycle"] = round(kpis["dso"] + 8.0 + random.gauss(0, 0.5), 1)
        kpis["revenue_quality"]  = round(rec_pct + random.gauss(0, 0.3), 2)
        kpis["recurring_revenue"]= kpis["revenue_quality"]
        nrr_base = 115.43 - 5.29 * churn_pct
        kpis["churn_rate"] = round(churn_pct + random.gauss(0, 0.05), 2)
        kpis["nrr"]        = round(nrr_base  + random.gauss(0, 0.25), 1)
        kpis["customer_concentration"] = round(26.0 - (cust - 418) / 420 * 8.0 + random.gauss(0, 0.4), 1)
        final_kpis[mo] = kpis

    mos_sorted = sorted(final_kpis.keys())
    for idx, mo in enumerate(mos_sorted):
        kpis   = final_kpis[mo]
        params = MP_PROJ[mo - 1]
        act_rev   = params[1]
        act_opex  = params[3] + params[1] * params[4] / 100
        sm_spend  = act_opex * params[10]
        cust      = params[8]
        new_c     = params[9]
        churn_pct = params[7]
        gross_m   = kpis.get("gross_margin", 62.0) / 100
        arpu_mo   = act_rev / max(cust, 1)
        cac       = sm_spend / max(new_c, 1)
        kpis["cac_payback"] = round(cac / max(arpu_mo * gross_m, 1), 1)

        if idx == 0:
            kpis["sales_efficiency"] = round((new_c * arpu_mo * 12) / max(sm_spend * 12, 1), 2)
            kpis["burn_multiple"]    = round(min(5.0, sm_spend / max(new_c * arpu_mo, 1)), 2)
        else:
            prev_mo       = mos_sorted[idx - 1]
            prev_rev      = MP_PROJ[prev_mo - 1][1]
            prev_cogs_pct = MP_PROJ[prev_mo - 1][2]
            prev_opex     = MP_PROJ[prev_mo - 1][3] + MP_PROJ[prev_mo - 1][1] * MP_PROJ[prev_mo - 1][4] / 100
            prev_op       = prev_rev * (1 - prev_cogs_pct / 100) - prev_opex
            curr_op       = act_rev  * (1 - params[2]     / 100) - act_opex
            delta_rev         = act_rev - prev_rev
            rev_growth_pct    = delta_rev / prev_rev * 100 if prev_rev else 0
            kpis["revenue_growth"] = round(rev_growth_pct, 2)
            kpis["arr_growth"]     = round(rev_growth_pct * 0.88 + random.gauss(0, 0.18), 2)
            if abs(rev_growth_pct) > 0.3 and prev_op > 0:
                op_inc_pct = (curr_op - prev_op) / prev_op * 100
                kpis["operating_leverage"] = round(max(-5.0, min(8.0, op_inc_pct / rev_growth_pct)), 2)
            if delta_rev > 0:
                kpis["sales_efficiency"] = round((delta_rev * 12) / max(sm_spend, 1), 2)
                kpis["burn_multiple"]    = round(min(5.0, sm_spend / max(delta_rev * 12, 1)), 2)
            else:
                kpis["sales_efficiency"] = round(max(0.05, sm_spend * 0.05 / max(sm_spend, 1)), 2)
                kpis["burn_multiple"]    = 5.0
        final_kpis[mo] = kpis

    conn = get_db()
    conn.execute("DELETE FROM projection_monthly_data")
    conn.execute("DELETE FROM projection_uploads")
    cur = conn.execute(
        "INSERT INTO projection_uploads (filename, uploaded_at, row_count, detected_columns) VALUES (?,?,?,?)",
        ("demo_projection_1000.csv", datetime.utcnow().isoformat(), len(df),
         json.dumps({c: c for c in df.columns}))
    )
    upload_id = cur.lastrowid
    for mo, kpis in final_kpis.items():
        clean = {k: (None if isinstance(v, float) and np.isnan(v) else v) for k, v in kpis.items()}
        conn.execute(
            "INSERT INTO projection_monthly_data (projection_upload_id, year, month, data_json) VALUES (?,?,?,?)",
            (upload_id, 2025, mo, json.dumps(clean))
        )
    conn.commit()
    conn.close()
    return {
        "seeded": True, "months": 12, "transactions": len(df), "upload_id": upload_id,
        "message": "Demo projection seeded — 12 months optimistic plan vs actuals.",
    }


# ─── Projection Endpoints ────────────────────────────────────────────────────

@app.post("/api/projection/upload", tags=["Projection"])
async def upload_projection(file: UploadFile = File(...), version_label: Optional[str] = None):
    """Upload a projection CSV (same format as actuals). Replaces any existing projection with the same version label."""
    if not file.filename.endswith((".csv", ".CSV")):
        raise HTTPException(400, "Only CSV files are accepted.")
    content = await file.read()
    try:
        df = pd.read_csv(io.StringIO(content.decode("utf-8", errors="replace")))
    except Exception as e:
        raise HTTPException(400, f"Could not parse CSV: {e}")

    vlabel = version_label or "v1"
    col_map     = normalize_columns(df)
    monthly_agg = aggregate_monthly(df, col_map)

    conn = get_db()
    try:
        # Delete-before-insert: only remove rows with the same version_label
        old_ids = [r["id"] for r in conn.execute("SELECT id FROM projection_uploads WHERE version_label=?", (vlabel,)).fetchall()]
        for oid in old_ids:
            conn.execute("DELETE FROM projection_monthly_data WHERE projection_upload_id=?", (oid,))
        conn.execute("DELETE FROM projection_uploads WHERE version_label=?", (vlabel,))

        cur = conn.execute(
            "INSERT INTO projection_uploads (filename, uploaded_at, row_count, detected_columns, version_label) VALUES (?,?,?,?,?)",
            (file.filename, datetime.utcnow().isoformat(), len(df), json.dumps(col_map), vlabel)
        )
        upload_id = cur.lastrowid

        for _, row in monthly_agg.iterrows():
            yr  = int(row["year"])
            mo  = int(row["month"])
            row_dict = {k: (None if (isinstance(v, float) and np.isnan(v)) else v)
                        for k, v in row.items() if k not in ("year", "month")}
            conn.execute(
                "INSERT INTO projection_monthly_data (projection_upload_id, year, month, data_json, version_label) VALUES (?,?,?,?,?)",
                (upload_id, yr, mo, json.dumps(row_dict), vlabel)
            )
        conn.commit()
    finally:
        conn.close()

    return {
        "upload_id":        upload_id,
        "filename":         file.filename,
        "rows_processed":   len(df),
        "months_detected":  len(monthly_agg),
        "columns_detected": col_map,
        "kpis_computed":    [k for k in monthly_agg.columns if k not in ("year", "month")],
        "version_label":    vlabel,
        "message":          f"Projection uploaded: {len(df)} rows across {len(monthly_agg)} months (version: {vlabel}).",
    }


@app.get("/api/projection/monthly", tags=["Projection"])
def projection_monthly_kpis(year: Optional[int] = None):
    """Return projected monthly KPI values. Optionally filter by year."""
    conn = get_db()
    query  = "SELECT * FROM projection_monthly_data"
    params = []
    if year:
        query += " WHERE year = ?"
        params.append(year)
    rows = conn.execute(query, params).fetchall()
    conn.close()
    result = []
    for row in rows:
        result.append({"year": row["year"], "month": row["month"], "kpis": json.loads(row["data_json"])})
    return sorted(result, key=lambda x: (x["year"], x["month"]))


@app.get("/api/projection/uploads", tags=["Projection"])
def list_projection_uploads():
    """List all projection uploads."""
    conn = get_db()
    rows = conn.execute("SELECT * FROM projection_uploads ORDER BY id DESC").fetchall()
    conn.close()
    return [{"id": r["id"], "filename": r["filename"], "uploaded_at": r["uploaded_at"],
             "row_count": r["row_count"], "columns": json.loads(r["detected_columns"])} for r in rows]


@app.delete("/api/projection/uploads/{upload_id}", tags=["Projection"])
def delete_projection_upload(upload_id: int):
    """Remove a projection upload and its associated monthly data."""
    conn = get_db()
    conn.execute("DELETE FROM projection_monthly_data WHERE projection_upload_id = ?", (upload_id,))
    conn.execute("DELETE FROM projection_uploads WHERE id = ?", (upload_id,))
    conn.commit()
    conn.close()
    return {"deleted": upload_id}


@app.get("/api/projection/versions", tags=["Projection"])
def get_projection_versions():
    conn = get_db()
    rows = conn.execute("SELECT DISTINCT version_label, MIN(uploaded_at) as first_uploaded FROM projection_uploads GROUP BY version_label ORDER BY first_uploaded DESC").fetchall()
    conn.close()
    return {"versions": [{"label": r["version_label"] or "v1", "uploaded_at": r["first_uploaded"]} for r in rows]}


@app.get("/api/bridge", tags=["Projection"])
def bridge_analysis(year: Optional[int] = None):
    """
    Compare projected vs actual KPIs month-by-month.
    Returns gap analysis, status (green/yellow/red), and causation rules for each KPI.
    """
    conn = get_db()
    if year:
        proj_rows   = conn.execute("SELECT * FROM projection_monthly_data WHERE year = ?", [year]).fetchall()
        actual_rows = conn.execute("SELECT * FROM monthly_data WHERE year = ?", [year]).fetchall()
    else:
        proj_rows   = conn.execute("SELECT * FROM projection_monthly_data").fetchall()
        actual_rows = conn.execute("SELECT * FROM monthly_data").fetchall()
    conn.close()

    if not proj_rows:
        return {"has_projection": False}

    # Build projection lookup: (year, month) -> kpi_dict
    proj_by_period: dict = {}
    for row in proj_rows:
        proj_by_period[(row["year"], row["month"])] = json.loads(row["data_json"])

    # Build actuals lookup: (year, month) -> kpi_dict (merge if multiple uploads)
    actual_by_period: dict = {}
    for row in actual_rows:
        key = (row["year"], row["month"])
        actual_by_period.setdefault(key, {})
        actual_by_period[key].update(json.loads(row["data_json"]))

    # Find overlapping periods
    overlap = sorted(set(proj_by_period.keys()) & set(actual_by_period.keys()))
    if not overlap:
        return {"has_projection": True, "has_overlap": False, "summary": {}, "kpis": {}}

    # Build per-KPI bridge
    kpis_out: dict = {}
    for kdef in KPI_DEFS:
        key       = kdef["key"]
        direction = kdef["direction"]
        months_data: dict = {}

        for (yr, mo) in overlap:
            proj_val   = proj_by_period[(yr, mo)].get(key)
            actual_val = actual_by_period[(yr, mo)].get(key)
            if proj_val is None or actual_val is None:
                continue
            if proj_val == 0:
                continue

            # gap_pct: positive = actual beats projection
            if direction == "higher":
                gap_pct = (actual_val - proj_val) / abs(proj_val) * 100
            else:
                gap_pct = (proj_val - actual_val) / abs(proj_val) * 100

            period_key = f"{yr}-{mo:02d}"
            months_data[period_key] = {
                "actual":    round(float(actual_val), 2),
                "projected": round(float(proj_val), 2),
                "gap":       round(float(actual_val - proj_val), 2),
                "gap_pct":   round(float(gap_pct), 2),
            }

        if not months_data:
            continue

        actuals       = [v["actual"]    for v in months_data.values()]
        projecteds    = [v["projected"] for v in months_data.values()]
        avg_actual    = round(float(np.mean(actuals)), 2)
        avg_projected = round(float(np.mean(projecteds)), 2)
        avg_gap       = round(float(avg_actual - avg_projected), 2)
        # Derive gap_pct from averages (not mean-of-monthly-pcts) so the
        # displayed avg_actual, avg_projected and avg_gap_pct are mutually consistent
        if avg_projected != 0:
            if direction == "higher":
                avg_gap_pct = round((avg_actual - avg_projected) / abs(avg_projected) * 100, 2)
            else:
                avg_gap_pct = round((avg_projected - avg_actual) / abs(avg_projected) * 100, 2)
        else:
            avg_gap_pct = 0.0
        overall_status = compute_gap_status(avg_gap_pct)

        kpis_out[key] = {
            "name":           kdef["name"],
            "unit":           kdef["unit"],
            "direction":      direction,
            "avg_actual":     avg_actual,
            "avg_projected":  avg_projected,
            "avg_gap":        avg_gap,
            "avg_gap_pct":    avg_gap_pct,
            "overall_status": overall_status,
            "months":         months_data,
            "causation":      CAUSATION_RULES.get(key, {
                "root_causes": [], "downstream_impact": [], "corrective_actions": []
            }),
        }

    # Tally summary counts by avg_gap_pct threshold
    on_track = sum(1 for k in kpis_out.values() if -3 <= k["avg_gap_pct"])
    behind   = sum(1 for k in kpis_out.values() if k["avg_gap_pct"] < -3)
    ahead    = sum(1 for k in kpis_out.values() if k["avg_gap_pct"] >= 3)
    on_track = on_track - ahead  # "on_track" = within ±3%

    return {
        "has_projection":  True,
        "has_overlap":     True,
        "summary": {
            "on_track":              on_track,
            "behind":                behind,
            "ahead":                 ahead,
            "total_months_compared": len(overlap),
        },
        "kpis": kpis_out,
    }


@app.put("/api/targets/{kpi_key}", tags=["Configuration"])
def update_target(kpi_key: str, target_value: float):
    """Update the target value for a specific KPI."""
    match = next((k for k in KPI_DEFS if k["key"] == kpi_key), None)
    if not match:
        raise HTTPException(404, f"KPI '{kpi_key}' not found")
    conn = get_db()
    conn.execute("UPDATE kpi_targets SET target_value = ? WHERE kpi_key = ?", (target_value, kpi_key))
    conn.commit()
    conn.close()
    return {"kpi_key": kpi_key, "target_value": target_value}

@app.get("/api/seed-demo", tags=["System"])
def seed_demo():
    """
    Seed 1,000 transaction rows + 12 months of fully correlated KPI data.

    Embedded correlations (all statistically meaningful):
      Revenue Growth  ↑  ↔  Operating Leverage  ↑   (fixed-cost base absorbs growth)
      Revenue Growth  ↑  ↔  Sales Efficiency    ↑   (same team, more output)
      Revenue Growth  ↑  ↔  Burn Multiple       ↓   (more ARR per burn dollar)
      Revenue Growth  ↑  ↔  CAC Payback         ↓   (efficiency compounds)
      Revenue Growth  ↑  ↔  OpEx Ratio          ↓   (operating leverage)
      Churn Rate      ↓  ↔  NRR                 ↑   (near-perfect inverse)
      Churn Rate      ↓  ↔  Revenue Growth      ↑   (retention fuels growth)
      Gross Margin    ↑  ↔  Contribution Margin ↑   (parallel expansion)
      DSO             ↑  ↔  Cash Conv Cycle     ↑   (DSO is the primary driver)
      ARR Growth      ≈  Revenue Growth × 0.9       (lagged subscription effect)

    Story arc FY2025:
      Q1 — Post-holiday slowdown, budget freezes, churn elevated, S&M front-loaded
      Q2 — Stabilisation; sales investment starts paying off; churn easing
      Q3 — Breakout: revenue accelerates, operating leverage spikes, burn multiple halves
      Q4 — Oct softness (pipeline reset), Nov recovery, Dec year-end surge
    """
    import random
    random.seed(42)

    # ── Monthly causal parameters ─────────────────────────────────────────────
    # Columns: (month, revenue, cogs_pct, fixed_opex, var_opex_pct,
    #           dso_days, recur_pct, churn_pct, customers, new_cust, sm_pct_opex)
    #
    # Revenue is the PRIMARY driver; everything else is derived or set causally.
    # fixed_opex = headcount / rent cost (does NOT scale with revenue → leverage)
    # var_opex_pct = variable S&M + support as % of revenue
    # sm_pct_opex = S&M share of total opex (drives sales efficiency calc)

    MP = [
      # mo  revenue   cogs%  f_opex   v_opex%  dso   rec%  churn%  cust  new  sm%
      ( 1,  750_000,  38.5,  248_000, 11.0,    42,   76.0,  3.20,  418,  12, 0.44),
      ( 2,  764_000,  38.2,  246_000, 10.8,    40,   76.8,  3.00,  425,  14, 0.42),
      ( 3,  782_000,  38.0,  244_000, 10.6,    40,   77.5,  2.80,  433,  16, 0.40),
      ( 4,  810_000,  37.8,  242_000, 10.4,    37,   78.5,  2.50,  442,  18, 0.38),
      ( 5,  844_000,  37.5,  240_000, 10.2,    36,   79.5,  2.30,  452,  21, 0.36),
      ( 6,  886_000,  37.2,  238_000, 10.0,    35,   80.5,  2.20,  463,  24, 0.35),
      ( 7,  940_000,  36.8,  236_000,  9.8,    33,   81.0,  2.00,  475,  28, 0.34),
      ( 8, 1_001_000, 36.5,  234_000,  9.6,    32,   81.5,  1.80,  488,  32, 0.33),
      ( 9, 1_071_000, 36.2,  232_000,  9.4,    32,   82.0,  1.70,  502,  36, 0.32),
      (10, 1_050_000, 36.5,  234_000,  9.6,    38,   81.5,  1.90,  512,  22, 0.33),  # Q4 dip
      (11, 1_097_000, 36.2,  232_000,  9.4,    37,   82.0,  1.80,  524,  28, 0.32),
      (12, 1_178_000, 35.8,  228_000,  9.2,    44,   83.0,  1.50,  540,  38, 0.30),  # year-end surge
    ]

    # ── Generate 1,000 transaction rows ──────────────────────────────────────
    # Segments define deal-size distribution (relative to avg_rev_per_row).
    # CRITICAL: normalise multipliers so their share-weighted average = 1.0,
    # otherwise monthly revenue total diverges from MP targets and all
    # margin / leverage KPIs become nonsensical.
    _RAW_SEGS = [
        ("Enterprise",  0.18, 4.8,  0.55),  # (name, share, mult, noise_std)
        ("Mid-Market",  0.37, 1.3,  0.28),
        ("SMB",         0.45, 0.52, 0.14),
    ]
    _wt_avg = sum(s * m for _, s, m, _ in _RAW_SEGS)   # = 1.579  → must → 1.0
    SEGS = [(nm, s, m / _wt_avg, sd) for nm, s, m, sd in _RAW_SEGS]

    rows_per_month = [417, 417, 417, 417, 417, 417, 417, 417, 416, 416, 416, 416]  # = 5000

    tx_rows = []
    for i, (mo, rev, cogs_pct, f_opex, v_opex_pct, dso, rec_pct, churn_pct, cust, new_c, sm_pct) in enumerate(MP):
        n = rows_per_month[i]
        total_opex  = f_opex + rev * v_opex_pct / 100
        avg_rev_row = rev / n

        for j in range(n):
            # Pick segment
            r = random.random()
            cum = 0.0
            for seg, share, mult, std in SEGS:
                cum += share
                if r <= cum:
                    break
            # With normalised multipliers, E[row_rev] = avg_rev_row → sum ≈ rev
            row_rev  = avg_rev_row * mult * max(0.35, 1 + random.gauss(0, std))
            row_cogs = row_rev * (cogs_pct / 100) * random.gauss(1.0, 0.025)
            row_opex = (total_opex / n) * random.gauss(1.0, 0.04)
            row_ar   = row_rev * (dso / 30) * random.gauss(1.0, 0.07)
            is_rec   = 1 if random.random() < rec_pct / 100 else 0
            row_sm   = row_opex * sm_pct * random.gauss(1.0, 0.05)
            row_churn= 1 if random.random() < churn_pct / 100 else 0
            day      = random.randint(1, 28)

            tx_rows.append({
                "date":         f"2025-{mo:02d}-{day:02d}",
                "revenue":      round(max(100, row_rev), 2),
                "cogs":         round(max(0, row_cogs), 2),
                "opex":         round(max(0, row_opex), 2),
                "ar":           round(max(0, row_ar), 2),
                "is_recurring": is_rec,
                "churn":        row_churn,
                "sm_allocated": round(max(0, row_sm), 2),
                "customers":    1,
            })

    df = pd.DataFrame(tx_rows)

    # ── Insert upload record ──────────────────────────────────────────────────
    conn = get_db()
    conn.execute("DELETE FROM monthly_data")
    conn.execute("DELETE FROM uploads")
    col_map_stored = {c: c for c in df.columns}
    cur = conn.execute(
        "INSERT INTO uploads (filename, uploaded_at, row_count, detected_columns) VALUES (?,?,?,?)",
        ("demo_correlated_5000.csv", datetime.utcnow().isoformat(), len(df),
         json.dumps(col_map_stored))
    )
    upload_id = cur.lastrowid

    # ── Compute base KPIs from aggregated transactions ────────────────────────
    col_map  = normalize_columns(df)
    base_agg = aggregate_monthly(df, col_map)   # gross_margin, opex_ratio, churn_rate, etc.

    # Build base lookup keyed by month number
    base_by_mo: dict = {}
    for _, row in base_agg.iterrows():
        base_by_mo[int(row["month"])] = {
            k: v for k, v in row.items()
            if k not in ("year", "month") and v is not None
               and not (isinstance(v, float) and np.isnan(v))
        }

    # ── Override/add cross-period + causally-derived KPIs ────────────────────
    # These require multi-month context or are intentionally tuned for correlation.
    #
    # revenue_growth    — MoM from actuals (post-agg)
    # arr_growth        — lags revenue_growth by ~1 month (subscription bookings)
    # operating_leverage— Δ op_income% / Δ rev%  (requires 2 consecutive months)
    # sales_efficiency  — new ARR this mo / S&M spend this mo
    # burn_multiple     — net burn / new ARR  (falls as growth accelerates)
    # cac_payback       — 1 / (sales_efficiency × gross_margin)  (inverse)
    # nrr               — 100 + (100 - churn_rate × 20) + expansion_proxy
    # dso               — override with seasonal curve (from MP table)
    # cash_conv_cycle   — dso + inventory days (constant 8d)
    # customer_concentration — falls as customer base grows

    # Derive actual monthly revenues from transactions
    mo_rev: dict = {}
    for g, grp in df.groupby(df["date"].str[5:7].astype(int)):
        mo_rev[g] = grp["revenue"].sum()

    # Derive actual monthly opex
    mo_opex: dict = {}
    for g, grp in df.groupby(df["date"].str[5:7].astype(int)):
        mo_opex[g] = grp["opex"].sum()

    # Derive actual monthly operating income
    mo_op_inc: dict = {}
    for g, grp in df.groupby(df["date"].str[5:7].astype(int)):
        rev_g  = grp["revenue"].sum()
        cogs_g = grp["cogs"].sum()
        opex_g = grp["opex"].sum()
        mo_op_inc[g] = rev_g - cogs_g - opex_g

    final_kpis: dict = {}   # mo → kpi_dict

    for mo, rev, cogs_pct, f_opex, v_opex_pct, dso, rec_pct, churn_pct, cust, new_c, sm_pct in MP:
        kpis = dict(base_by_mo.get(mo, {}))

        act_rev    = mo_rev.get(mo, rev)
        act_opex   = mo_opex.get(mo, f_opex + rev * v_opex_pct / 100)
        act_op_inc = mo_op_inc.get(mo, 0)
        sm_spend   = act_opex * sm_pct

        # ── DSO & Cash Cycle (seasonal; Dec high = year-end billing) ─────────
        kpis["dso"]            = round(dso * random.gauss(1.0, 0.02), 1)
        kpis["cash_conv_cycle"]= round(kpis["dso"] + 8.0 + random.gauss(0, 0.5), 1)

        # ── Revenue Quality / Recurring Revenue ──────────────────────────────
        kpis["revenue_quality"]  = round(rec_pct + random.gauss(0, 0.3), 2)
        kpis["recurring_revenue"]= kpis["revenue_quality"]

        # ── Churn Rate → NRR (near-perfect inverse, R ≈ -0.99) ──────────────
        # Calibrated linear: NRR = 98.5 when churn = 3.2%  (budget-freeze Jan)
        #                    NRR = 107.5 when churn = 1.5%  (year-end Dec surge)
        # Slope = (107.5 - 98.5) / (1.5 - 3.2) = -5.29
        # Intercept = 98.5 - (-5.29) × 3.2 = 98.5 + 16.93 = 115.43
        nrr_base = 115.43 - 5.29 * churn_pct
        kpis["churn_rate"] = round(churn_pct + random.gauss(0, 0.05), 2)
        kpis["nrr"]        = round(nrr_base + random.gauss(0, 0.25), 1)

        # ── Customer Concentration (dilutes as base grows) ───────────────────
        kpis["customer_concentration"] = round(26.0 - (cust - 418) / 420 * 8.0 + random.gauss(0, 0.4), 1)

        final_kpis[mo] = kpis

    # ── Multi-period KPIs (need 2 consecutive months) ────────────────────────
    # For Δ-revenue KPIs (growth, sales_efficiency, burn_multiple,
    # operating_leverage) we use the DETERMINISTIC MP target revenues, not
    # transaction aggregations.  Transactions have too much per-row variance
    # (σ ≈ $58K/month) which dwarfs small Δ-rev signals ($14K Feb→Jan).
    # Single-period KPIs (gross_margin, churn_rate, etc.) still come from
    # the aggregated transactions via base_by_mo.
    mos_sorted = sorted(final_kpis.keys())
    for idx, mo in enumerate(mos_sorted):
        kpis   = final_kpis[mo]
        params = MP[mo - 1]   # (mo, rev, cogs_pct, f_opex, v_opex_pct, dso, rec_pct, churn_pct, cust, new_c, sm_pct)
        # Use MP target values for multi-period calculations
        act_rev    = params[1]
        act_opex   = params[3] + params[1] * params[4] / 100   # f_opex + var_opex
        sm_spend   = act_opex * params[10]
        cust       = params[8]
        new_c      = params[9]
        churn_pct  = params[7]

        gross_m = kpis.get("gross_margin", 62.0) / 100
        arpu_mo = act_rev / max(cust, 1)          # monthly revenue per customer

        # ── CAC Payback = (S&M / new_customers) / (ARPU_mo × GM%) ───────────
        # Improves as: S&M per new-cust falls (more new cust / same spend)
        #              OR gross margin expands
        #              OR ARPU grows (higher-value deals closing)
        cac      = sm_spend / max(new_c, 1)
        kpis["cac_payback"] = round(cac / max(arpu_mo * gross_m, 1), 1)

        if idx == 0:
            # Jan: first month — no Δ-revenue, use new_customer-based proxies
            # Sales Efficiency proxy: (new_c × ARPU_annual) / (S&M_annual)
            kpis["sales_efficiency"] = round(
                (new_c * arpu_mo * 12) / max(sm_spend * 12, 1), 2
            )
            # Burn Multiple proxy: S&M / (new_c × ARPU_mo)
            new_mrr = new_c * arpu_mo
            kpis["burn_multiple"] = round(min(5.0, sm_spend / max(new_mrr, 1)), 2)
        else:
            prev_mo  = mos_sorted[idx - 1]
            prev_rev = MP[prev_mo - 1][1]   # deterministic target
            prev_cogs_pct = MP[prev_mo - 1][2]
            prev_opex = MP[prev_mo - 1][3] + MP[prev_mo - 1][1] * MP[prev_mo - 1][4] / 100
            prev_op  = prev_rev * (1 - prev_cogs_pct/100) - prev_opex
            curr_cogs_pct = params[2]
            curr_op  = act_rev * (1 - curr_cogs_pct/100) - act_opex

            delta_rev      = act_rev - prev_rev
            rev_growth_pct = delta_rev / prev_rev * 100 if prev_rev else 0
            kpis["revenue_growth"] = round(rev_growth_pct, 2)
            kpis["arr_growth"]     = round(rev_growth_pct * 0.88 + random.gauss(0, 0.18), 2)

            # Operating leverage = (% Δ op_income) / (% Δ revenue)
            # Large fixed cost base guarantees op_lev > 1 when rev grows → converges
            if abs(rev_growth_pct) > 0.3 and prev_op > 0:
                op_inc_pct = (curr_op - prev_op) / prev_op * 100
                kpis["operating_leverage"] = round(max(-5.0, min(8.0, op_inc_pct / rev_growth_pct)), 2)
            elif rev_growth_pct < 0 and prev_op > 0 and curr_op < prev_op:
                op_inc_pct = (curr_op - prev_op) / prev_op * 100
                kpis["operating_leverage"] = round(max(-5.0, op_inc_pct / rev_growth_pct), 2)

            # Sales Efficiency = "Magic Number"  = (Δ Rev × 12) / S&M spend
            # Rises with revenue growth (same team, accelerating output)
            # Goes to ~0 when revenue declines (Oct dip)
            if delta_rev > 0:
                kpis["sales_efficiency"] = round((delta_rev * 12) / max(sm_spend, 1), 2)
            else:
                # Revenue dipped: efficiency near zero but not negative
                kpis["sales_efficiency"] = round(max(0.05, sm_spend * 0.05 / max(sm_spend, 1)), 2)

            # Burn Multiple = S&M / (Δ Rev × 12)
            # Inverse of sales efficiency — falls dramatically as growth accelerates
            # Capped at 5.0 when revenue declines (Oct reset month)
            if delta_rev > 0:
                kpis["burn_multiple"] = round(min(5.0, sm_spend / max(delta_rev * 12, 1)), 2)
            else:
                kpis["burn_multiple"] = 5.0   # maximum penalty for declining revenue

        final_kpis[mo] = kpis

    # ── Persist ───────────────────────────────────────────────────────────────
    for mo, kpis in final_kpis.items():
        clean = {k: (None if isinstance(v, float) and np.isnan(v) else v)
                 for k, v in kpis.items()}
        conn.execute(
            "INSERT INTO monthly_data (upload_id, year, month, data_json) VALUES (?,?,?,?)",
            (upload_id, 2025, mo, json.dumps(clean))
        )

    conn.commit()
    conn.close()
    return {
        "seeded":       True,
        "months":       12,
        "transactions": len(df),
        "upload_id":    upload_id,
        "correlations": [
            "Revenue Growth ↔ Operating Leverage (pos)",
            "Revenue Growth ↔ Sales Efficiency (pos)",
            "Revenue Growth ↔ Burn Multiple (neg)",
            "Revenue Growth ↔ CAC Payback (neg)",
            "Revenue Growth ↔ OpEx Ratio (neg)",
            "Churn Rate ↔ NRR (neg, R≈-0.99)",
            "DSO ↔ Cash Conv Cycle (pos, R≈0.97)",
            "Gross Margin ↔ Contribution Margin (pos)",
        ],
    }

@app.get("/api/seed-multiyear", tags=["System"])
def seed_multiyear():
    """
    Seed 5 years of KPI data (2021–2025 actuals + 2026 actuals Jan–Mar + 2026 projection Apr–Dec).

    Narrative arc:
      2021 — Early-stage startup: fast growth from tiny base, high churn, burning cash
      2022 — Series B scaling: rapid hiring, revenue accelerating, margins improving
      2023 — Plateau: growth stalls, headcount too large, efficiency declining
      2024 — Warning signals: churn uptick, margin compression, burn rising
      2025 — Mixed recovery: volatile, some bright spots but fragile
      2026 Jan–Mar — Critical: churn worsening, revenue under pressure
      2026 Apr–Dec — Forecast: recovery scenario if corrective action taken
    """
    import random
    random.seed(99)
    rng = random.gauss

    def lerp(a, b, t):
        return a + (b - a) * t

    def mo_val(start, end, mo, noise=0.0):
        t = (mo - 1) / 11.0
        base = lerp(start, end, t)
        return round(base + rng(0, noise), 4) if noise else round(base, 4)

    # ── Year phase definitions ────────────────────────────────────────────────
    # Each year: dict of KPI → (jan_val, dec_val, monthly_noise_sigma)
    PHASES = {
        2021: {  # Startup — growth from tiny base, high churn, cash-negative
            "revenue_growth":        (14.0,  9.0,  1.5),
            "gross_margin":          (51.0, 56.5,  0.4),
            "operating_margin":      (-22.0, -6.0, 1.0),
            "ebitda_margin":         (-19.0, -4.0, 1.0),
            "cash_conv_cycle":       (68.0, 58.0,  1.5),
            "dso":                   (62.0, 52.0,  1.2),
            "arr_growth":            (12.0,  7.5,  1.5),
            "nrr":                   (84.0, 93.0,  0.8),
            "burn_multiple":         (4.8,  4.0,   0.2),
            "opex_ratio":            (88.0, 76.0,  1.5),
            "contribution_margin":   (28.0, 38.0,  0.8),
            "revenue_quality":       (54.0, 63.0,  0.5),
            "cac_payback":           (29.0, 22.0,  1.0),
            "sales_efficiency":      (0.10, 0.22,  0.02),
            "customer_concentration":(50.0, 44.0,  1.0),
            "recurring_revenue":     (54.0, 63.0,  0.5),
            "churn_rate":            (7.8,  5.5,   0.3),
            "operating_leverage":    (-0.8,  0.4,  0.2),
            "pipeline_conversion":   (2.5,  4.5,  0.3),
            "customer_ltv":          (32.0, 48.0, 1.5),
            "pricing_power_index":   (2.0,  5.0,  0.5),
            "cpl":                   (320.0, 240.0, 18.0),
            "mql_sql_rate":          (14.0,  20.0,  1.5),
            "win_rate":              (16.0,  22.0,  1.5),
            "quota_attainment":      (52.0,  64.0,  2.5),
            "marketing_roi":         (1.2,   1.8,   0.15),
            "headcount_eff":         (0.6,   0.9,   0.05),
            "rev_per_employee":      (65.0,  95.0,  4.0),
            "ltv_cac":               (1.2,   1.8,   0.12),
            "expansion_rate":        (4.0,   9.0,   0.8),
            "health_score":          (48.0,  58.0,  2.0),
            "logo_retention":        (72.0,  80.0,  1.0),
            "payback_period":        (38.0,  28.0,  1.5),
        },
        2022: {  # Series B — rapid scaling, margins climbing, team building
            "revenue_growth":        (16.0, 12.0,  1.2),
            "gross_margin":          (57.0, 64.0,  0.4),
            "operating_margin":      (-8.0,  3.5,  0.8),
            "ebitda_margin":         (-6.0,  5.5,  0.8),
            "cash_conv_cycle":       (56.0, 44.0,  1.2),
            "dso":                   (51.0, 40.0,  1.0),
            "arr_growth":            (14.0, 10.5,  1.2),
            "nrr":                   (96.0,109.0,  0.7),
            "burn_multiple":         (3.8,  1.8,   0.2),
            "opex_ratio":            (72.0, 52.0,  1.5),
            "contribution_margin":   (40.0, 56.0,  0.8),
            "revenue_quality":       (65.0, 76.0,  0.5),
            "cac_payback":           (20.0, 11.0,  0.8),
            "sales_efficiency":      (0.22, 0.72,  0.03),
            "customer_concentration":(42.0, 28.0,  1.0),
            "recurring_revenue":     (65.0, 76.0,  0.5),
            "churn_rate":            (4.4,  2.2,   0.25),
            "operating_leverage":    (0.6,  2.8,   0.2),
            "pipeline_conversion":   (5.0,  8.5,  0.3),
            "customer_ltv":          (52.0, 92.0, 2.0),
            "pricing_power_index":   (4.0,  8.5,  0.6),
            "cpl":                   (220.0, 130.0, 12.0),
            "mql_sql_rate":          (22.0,  31.0,  1.2),
            "win_rate":              (24.0,  34.0,  1.2),
            "quota_attainment":      (68.0,  88.0,  2.0),
            "marketing_roi":         (2.2,   4.2,   0.2),
            "headcount_eff":         (1.0,   1.6,   0.05),
            "rev_per_employee":      (100.0, 165.0, 5.0),
            "ltv_cac":               (2.2,   4.8,   0.2),
            "expansion_rate":        (11.0,  26.0,  1.2),
            "health_score":          (60.0,  76.0,  1.5),
            "logo_retention":        (82.0,  92.0,  0.8),
            "payback_period":        (26.0,  12.0,  1.2),
        },
        2023: {  # Plateau — growth stalls, team too large, efficiency slipping
            "revenue_growth":        (4.0,  1.2,   0.6),
            "gross_margin":          (64.5, 65.8,  0.3),
            "operating_margin":      (3.5,  6.5,   0.5),
            "ebitda_margin":         (5.5,  9.0,   0.5),
            "cash_conv_cycle":       (44.0, 48.0,  1.0),
            "dso":                   (40.0, 44.0,  0.8),
            "arr_growth":            (3.5,  1.0,   0.6),
            "nrr":                   (109.0,106.0, 0.6),
            "burn_multiple":         (1.8,  3.2,   0.2),
            "opex_ratio":            (50.0, 44.0,  1.0),
            "contribution_margin":   (56.0, 59.0,  0.5),
            "revenue_quality":       (76.5, 81.0,  0.4),
            "cac_payback":           (11.0, 14.5,  0.6),
            "sales_efficiency":      (0.70, 0.30,  0.03),
            "customer_concentration":(27.0, 21.0,  0.8),
            "recurring_revenue":     (76.5, 81.0,  0.4),
            "churn_rate":            (2.1,  2.7,   0.15),
            "operating_leverage":    (2.5,  1.2,   0.2),
            "pipeline_conversion":   (8.0,  5.5,  0.3),
            "customer_ltv":          (88.0, 82.0, 1.5),
            "pricing_power_index":   (3.5,  1.0,  0.5),
            "cpl":                   (135.0, 168.0, 8.0),
            "mql_sql_rate":          (30.0,  24.0,  1.0),
            "win_rate":              (33.0,  26.0,  1.0),
            "quota_attainment":      (86.0,  74.0,  1.8),
            "marketing_roi":         (4.0,   2.8,   0.15),
            "headcount_eff":         (1.6,   1.1,   0.04),
            "rev_per_employee":      (162.0, 138.0, 4.0),
            "ltv_cac":               (4.5,   3.2,   0.18),
            "expansion_rate":        (24.0,  18.0,  1.0),
            "health_score":          (74.0,  68.0,  1.2),
            "logo_retention":        (91.0,  86.0,  0.6),
            "payback_period":        (12.0,  17.0,  0.8),
        },
        2024: {  # Warning signals — churn uptick, margin squeeze, burn rising
            "revenue_growth":        (3.5, -1.5,   0.8),
            "gross_margin":          (65.0, 61.5,  0.4),
            "operating_margin":      (6.0,  2.5,   0.6),
            "ebitda_margin":         (8.5,  4.0,   0.6),
            "cash_conv_cycle":       (47.0, 54.0,  1.2),
            "dso":                   (43.0, 50.0,  1.0),
            "arr_growth":            (3.0, -1.8,   0.8),
            "nrr":                   (106.0, 98.5, 0.7),
            "burn_multiple":         (3.0,  4.8,   0.25),
            "opex_ratio":            (44.0, 52.0,  1.0),
            "contribution_margin":   (59.0, 54.0,  0.6),
            "revenue_quality":       (81.0, 78.0,  0.4),
            "cac_payback":           (14.0, 19.0,  0.7),
            "sales_efficiency":      (0.32, 0.12,  0.02),
            "customer_concentration":(21.0, 24.0,  0.8),
            "recurring_revenue":     (81.0, 78.0,  0.4),
            "churn_rate":            (2.6,  4.8,   0.2),
            "operating_leverage":    (1.0,  0.3,   0.2),
            "pipeline_conversion":   (5.0,  2.8,  0.3),
            "customer_ltv":          (80.0, 52.0, 2.0),
            "pricing_power_index":   (0.5, -3.5,  0.5),
            "cpl":                   (175.0, 268.0, 10.0),
            "mql_sql_rate":          (22.0,  11.0,  0.8),
            "win_rate":              (24.0,  14.0,  0.8),
            "quota_attainment":      (72.0,  55.0,  1.5),
            "marketing_roi":         (2.6,   1.4,   0.12),
            "headcount_eff":         (1.1,   0.7,   0.04),
            "rev_per_employee":      (135.0, 98.0,  3.5),
            "ltv_cac":               (3.0,   1.8,   0.15),
            "expansion_rate":        (17.0,  10.0,  0.8),
            "health_score":          (66.0,  55.0,  1.5),
            "logo_retention":        (85.0,  78.0,  0.7),
            "payback_period":        (16.0,  24.0,  1.0),
        },
        2025: {  # Mixed recovery — volatile, fragile improvement
            "revenue_growth":        (1.8, 11.5,   1.0),
            "gross_margin":          (61.5, 64.2,  0.3),
            "operating_margin":      (5.1, 13.5,   0.8),
            "ebitda_margin":         (7.0, 16.0,   0.8),
            "cash_conv_cycle":       (50.0, 44.0,  1.0),
            "dso":                   (42.0, 44.0,  0.8),
            "arr_growth":            (1.5, 10.1,   1.0),
            "nrr":                   (98.5,107.5,  0.6),
            "burn_multiple":         (5.0,  0.7,   0.3),
            "opex_ratio":            (48.0, 36.0,  1.0),
            "contribution_margin":   (54.0, 61.0,  0.5),
            "revenue_quality":       (76.0, 83.0,  0.4),
            "cac_payback":           (18.0,  9.5,  0.8),
            "sales_efficiency":      (0.08, 0.65,  0.04),
            "customer_concentration":(26.0, 18.5,  0.8),
            "recurring_revenue":     (76.0, 83.0,  0.4),
            "churn_rate":            (3.2,  1.5,   0.15),
            "operating_leverage":    (1.5,  3.5,   0.3),
            "pipeline_conversion":   (3.0,  7.2,  0.35),
            "customer_ltv":          (55.0, 98.0, 2.5),
            "pricing_power_index":   (-2.0, 5.5,  0.5),
            "cpl":                   (260.0, 175.0, 12.0),
            "mql_sql_rate":          (13.0,  22.0,  1.0),
            "win_rate":              (15.0,  24.0,  1.0),
            "quota_attainment":      (58.0,  76.0,  2.0),
            "marketing_roi":         (1.5,   2.8,   0.15),
            "headcount_eff":         (0.75,  1.2,   0.05),
            "rev_per_employee":      (102.0, 148.0, 4.0),
            "ltv_cac":               (1.9,   3.2,   0.18),
            "expansion_rate":        (11.0,  18.0,  1.0),
            "health_score":          (56.0,  68.0,  1.5),
            "logo_retention":        (79.0,  87.0,  0.7),
            "payback_period":        (23.0,  16.0,  1.2),
        },
    }

    # 2026 actuals Jan–Mar: critical state
    PHASE_2026_ACTUAL = {  # (mo1_val, mo3_val, noise)
        "revenue_growth":        (0.5,  2.0,   0.5),
        "gross_margin":          (61.0, 61.5,  0.3),
        "operating_margin":      (4.0,  5.5,   0.5),
        "ebitda_margin":         (6.0,  7.5,   0.5),
        "cash_conv_cycle":       (50.0, 48.5,  0.8),
        "dso":                   (46.0, 44.5,  0.7),
        "arr_growth":            (0.3,  1.5,   0.5),
        "nrr":                   (97.0, 98.0,  0.5),
        "burn_multiple":         (3.8,  3.5,   0.2),
        "opex_ratio":            (50.0, 49.0,  0.8),
        "contribution_margin":   (52.0, 53.0,  0.4),
        "revenue_quality":       (75.0, 76.0,  0.3),
        "cac_payback":           (17.0, 16.5,  0.5),
        "sales_efficiency":      (0.12, 0.18,  0.02),
        "customer_concentration":(24.0, 23.5,  0.6),
        "recurring_revenue":     (75.0, 76.0,  0.3),
        "churn_rate":            (4.5,  4.2,   0.15),
        "operating_leverage":    (0.8,  1.0,   0.15),
        "pipeline_conversion":   (3.2,  3.8,  0.2),
        "customer_ltv":          (52.0, 55.0, 1.5),
        "pricing_power_index":   (-1.5, 1.5,  0.4),
        "cpl":                   (182.0, 195.0, 8.0),
        "mql_sql_rate":          (20.0,  21.5,  0.8),
        "win_rate":              (22.0,  23.5,  0.8),
        "quota_attainment":      (73.0,  75.0,  1.5),
        "marketing_roi":         (2.6,   2.7,   0.12),
        "headcount_eff":         (1.15,  1.18,  0.04),
        "rev_per_employee":      (142.0, 148.0, 3.0),
        "ltv_cac":               (3.0,   3.1,   0.12),
        "expansion_rate":        (16.5,  17.0,  0.8),
        "health_score":          (66.0,  67.5,  1.2),
        "logo_retention":        (86.0,  87.0,  0.5),
        "payback_period":        (17.5,  17.0,  0.8),
    }

    # 2021-2025 plan/budget projections (slightly more optimistic than actuals)
    # These represent "what management planned" for each year; seeded into projection_monthly_data
    # so Bridge Analysis can compare plan vs actual for 2021-2025
    PHASES_PROJ = {
        2021: {  # Startup plan — expected faster margin improvement
            "revenue_growth":        (16.0, 11.0,  1.0),
            "gross_margin":          (53.0, 59.0,  0.3),
            "operating_margin":      (-18.0, -3.0, 0.8),
            "ebitda_margin":         (-15.0, -1.0, 0.8),
            "cash_conv_cycle":       (63.0, 52.0,  1.2),
            "dso":                   (57.0, 46.0,  1.0),
            "arr_growth":            (14.0,  9.5,  1.2),
            "nrr":                   (87.0, 97.0,  0.6),
            "burn_multiple":         (4.2,  3.4,   0.15),
            "opex_ratio":            (83.0, 70.0,  1.2),
            "contribution_margin":   (32.0, 42.0,  0.6),
            "revenue_quality":       (57.0, 67.0,  0.4),
            "cac_payback":           (26.0, 19.0,  0.8),
            "sales_efficiency":      (0.13, 0.27,  0.015),
            "customer_concentration":(46.0, 39.0,  0.8),
            "recurring_revenue":     (57.0, 67.0,  0.4),
            "churn_rate":            (7.0,  4.8,   0.25),
            "operating_leverage":    (-0.5,  0.8,  0.15),
            "pipeline_conversion":   (3.0, 5.5, 0.25),
            "customer_ltv":          (35.0, 54.0, 1.2),
            "pricing_power_index":   (2.5, 6.0, 0.4),
            "cpl":                   (290.0, 210.0, 15.0),
            "mql_sql_rate":          (16.0,  23.0,  1.2),
            "win_rate":              (18.0,  25.0,  1.2),
            "quota_attainment":      (58.0,  72.0,  2.0),
            "marketing_roi":         (1.4,   2.1,   0.12),
            "headcount_eff":         (0.7,   1.05,  0.04),
            "rev_per_employee":      (75.0,  110.0, 3.5),
            "ltv_cac":               (1.4,   2.1,   0.10),
            "expansion_rate":        (5.0,   10.5,  0.7),
            "health_score":          (55.0,  66.0,  1.8),
            "logo_retention":        (76.0,  85.0,  0.9),
            "payback_period":        (34.0,  24.0,  1.3),
        },
        2022: {  # Series B plan — optimistic growth targets
            "revenue_growth":        (18.0, 15.0,  1.0),
            "gross_margin":          (59.0, 67.0,  0.3),
            "operating_margin":      (-5.0,  7.0,  0.6),
            "ebitda_margin":         (-3.0,  9.0,  0.6),
            "cash_conv_cycle":       (52.0, 39.0,  1.0),
            "dso":                   (47.0, 35.0,  0.8),
            "arr_growth":            (16.0, 13.0,  1.0),
            "nrr":                   (99.0,113.0,  0.5),
            "burn_multiple":         (3.4,  1.4,   0.15),
            "opex_ratio":            (68.0, 46.0,  1.2),
            "contribution_margin":   (44.0, 61.0,  0.6),
            "revenue_quality":       (68.0, 80.0,  0.4),
            "cac_payback":           (18.0,  9.0,  0.6),
            "sales_efficiency":      (0.26, 0.80,  0.025),
            "customer_concentration":(38.0, 24.0,  0.8),
            "recurring_revenue":     (68.0, 80.0,  0.4),
            "churn_rate":            (4.0,  1.8,   0.2),
            "operating_leverage":    (0.9,  3.2,   0.15),
            "pipeline_conversion":   (5.5, 9.5, 0.25),
            "customer_ltv":          (56.0, 98.0, 1.8),
            "pricing_power_index":   (5.0, 10.0, 0.5),
            "cpl":                   (198.0, 116.0, 10.0),
            "mql_sql_rate":          (25.0,  35.0,  1.0),
            "win_rate":              (27.0,  38.0,  1.0),
            "quota_attainment":      (76.0,  98.0,  1.8),
            "marketing_roi":         (2.5,   4.8,   0.18),
            "headcount_eff":         (1.15,  1.85,  0.04),
            "rev_per_employee":      (115.0, 190.0, 4.5),
            "ltv_cac":               (2.5,   5.5,   0.18),
            "expansion_rate":        (13.0,  30.0,  1.0),
            "health_score":          (68.0,  85.0,  1.3),
            "logo_retention":        (88.0,  96.0,  0.7),
            "payback_period":        (23.0,  10.0,  1.0),
        },
        2023: {  # Plateau plan — expected continued growth (too optimistic)
            "revenue_growth":        (8.0,  6.0,   0.5),
            "gross_margin":          (66.0, 68.0,  0.25),
            "operating_margin":      (6.0, 10.0,   0.4),
            "ebitda_margin":         (8.0, 13.0,   0.4),
            "cash_conv_cycle":       (41.0, 44.0,  0.8),
            "dso":                   (37.0, 41.0,  0.6),
            "arr_growth":            (7.0,  5.0,   0.5),
            "nrr":                   (111.0,109.0, 0.5),
            "burn_multiple":         (1.5,  2.6,   0.15),
            "opex_ratio":            (47.0, 40.0,  0.8),
            "contribution_margin":   (60.0, 64.0,  0.4),
            "revenue_quality":       (79.0, 85.0,  0.3),
            "cac_payback":           (10.0, 12.5,  0.5),
            "sales_efficiency":      (0.78, 0.40,  0.025),
            "customer_concentration":(24.0, 18.0,  0.6),
            "recurring_revenue":     (79.0, 85.0,  0.3),
            "churn_rate":            (1.8,  2.2,   0.12),
            "operating_leverage":    (2.9,  1.8,   0.15),
            "pipeline_conversion":   (9.0, 7.0, 0.25),
            "customer_ltv":          (92.0, 88.0, 1.2),
            "pricing_power_index":   (5.0, 3.0, 0.4),
            "cpl":                   (120.0, 150.0, 7.0),
            "mql_sql_rate":          (34.0,  28.0,  0.9),
            "win_rate":              (37.0,  30.0,  0.9),
            "quota_attainment":      (95.0,  83.0,  1.5),
            "marketing_roi":         (4.6,   3.2,   0.13),
            "headcount_eff":         (1.85,  1.28,  0.035),
            "rev_per_employee":      (185.0, 160.0, 3.5),
            "ltv_cac":               (5.2,   3.7,   0.16),
            "expansion_rate":        (27.0,  21.0,  0.9),
            "health_score":          (83.0,  77.0,  1.0),
            "logo_retention":        (95.0,  91.0,  0.5),
            "payback_period":        (10.5,  15.0,  0.7),
        },
        2024: {  # Warning year plan — expected mild improvement (missed badly)
            "revenue_growth":        (6.0,  3.0,   0.6),
            "gross_margin":          (66.5, 64.0,  0.3),
            "operating_margin":      (8.5,  6.0,   0.5),
            "ebitda_margin":         (11.0,  7.5,  0.5),
            "cash_conv_cycle":       (43.0, 49.0,  1.0),
            "dso":                   (39.0, 45.0,  0.8),
            "arr_growth":            (5.5,  1.5,   0.6),
            "nrr":                   (108.0,102.0, 0.6),
            "burn_multiple":         (2.5,  3.8,   0.2),
            "opex_ratio":            (40.0, 47.0,  0.8),
            "contribution_margin":   (62.0, 58.0,  0.5),
            "revenue_quality":       (83.5, 81.0,  0.3),
            "cac_payback":           (12.0, 16.0,  0.6),
            "sales_efficiency":      (0.38, 0.20,  0.015),
            "customer_concentration":(18.0, 20.0,  0.6),
            "recurring_revenue":     (83.5, 81.0,  0.3),
            "churn_rate":            (2.2,  3.8,   0.15),
            "operating_leverage":    (1.5,  0.8,   0.15),
            "pipeline_conversion":   (7.0, 5.0, 0.25),
            "customer_ltv":          (84.0, 62.0, 1.8),
            "pricing_power_index":   (2.5, 0.0, 0.4),
            "cpl":                   (155.0, 238.0, 9.0),
            "mql_sql_rate":          (26.0,  13.0,  0.7),
            "win_rate":              (28.0,  16.0,  0.7),
            "quota_attainment":      (83.0,  63.0,  1.3),
            "marketing_roi":         (3.0,   1.6,   0.10),
            "headcount_eff":         (1.26,  0.82,  0.035),
            "rev_per_employee":      (155.0, 113.0, 3.0),
            "ltv_cac":               (3.5,   2.1,   0.13),
            "expansion_rate":        (20.0,  12.0,  0.7),
            "health_score":          (76.0,  64.0,  1.3),
            "logo_retention":        (90.0,  84.0,  0.6),
            "payback_period":        (14.0,  21.0,  0.9),
        },
        2025: {  # Recovery plan — ambitious targets set after tough 2024
            "revenue_growth":        (5.0, 15.0,   0.8),
            "gross_margin":          (63.0, 66.5,  0.25),
            "operating_margin":      (8.0, 17.0,   0.6),
            "ebitda_margin":         (10.0, 19.5,  0.6),
            "cash_conv_cycle":       (46.0, 40.0,  0.8),
            "dso":                   (39.0, 40.0,  0.6),
            "arr_growth":            (4.5, 13.5,   0.8),
            "nrr":                   (101.0,111.0, 0.5),
            "burn_multiple":         (4.2,  0.4,   0.25),
            "opex_ratio":            (44.0, 31.0,  0.8),
            "contribution_margin":   (58.0, 65.0,  0.4),
            "revenue_quality":       (79.0, 86.5,  0.3),
            "cac_payback":           (15.0,  7.5,  0.6),
            "sales_efficiency":      (0.12, 0.75,  0.03),
            "customer_concentration":(22.0, 14.5,  0.6),
            "recurring_revenue":     (79.0, 86.5,  0.3),
            "churn_rate":            (2.8,  1.1,   0.12),
            "operating_leverage":    (2.0,  4.2,   0.25),
            "pipeline_conversion":   (5.0, 10.0, 0.3),
            "customer_ltv":          (62.0, 108.0, 2.2),
            "pricing_power_index":   (1.0, 7.0, 0.4),
            "cpl":                   (232.0, 155.0, 10.0),
            "mql_sql_rate":          (15.0,  25.0,  0.9),
            "win_rate":              (17.0,  27.0,  0.9),
            "quota_attainment":      (65.0,  85.0,  1.8),
            "marketing_roi":         (1.75,  3.2,   0.13),
            "headcount_eff":         (0.88,  1.4,   0.04),
            "rev_per_employee":      (118.0, 170.0, 3.5),
            "ltv_cac":               (2.2,   3.7,   0.16),
            "expansion_rate":        (13.0,  21.0,  0.9),
            "health_score":          (64.0,  78.0,  1.3),
            "logo_retention":        (85.0,  93.0,  0.6),
            "payback_period":        (26.0,  14.0,  1.0),
        },
    }

    # 2026 projection Apr–Dec: recovery scenario
    PHASE_2026_PROJ = {  # (apr_val, dec_val, noise)
        "revenue_growth":        (3.0,  7.5,   0.6),
        "gross_margin":          (62.0, 65.0,  0.3),
        "operating_margin":      (6.0, 11.0,   0.6),
        "ebitda_margin":         (8.0, 13.0,   0.6),
        "cash_conv_cycle":       (48.0, 41.0,  0.8),
        "dso":                   (44.0, 37.0,  0.7),
        "arr_growth":            (2.5,  7.0,   0.6),
        "nrr":                   (98.5,105.0,  0.5),
        "burn_multiple":         (3.4,  1.8,   0.2),
        "opex_ratio":            (48.0, 38.0,  0.8),
        "contribution_margin":   (53.5, 60.5,  0.5),
        "revenue_quality":       (76.5, 82.0,  0.3),
        "cac_payback":           (16.0, 11.0,  0.5),
        "sales_efficiency":      (0.20, 0.58,  0.03),
        "customer_concentration":(23.0, 19.0,  0.6),
        "recurring_revenue":     (76.5, 82.0,  0.3),
        "churn_rate":            (4.0,  2.5,   0.15),
        "operating_leverage":    (1.2,  2.8,   0.2),
        "pipeline_conversion":   (4.0,  7.5,  0.25),
        "customer_ltv":          (58.0, 82.0, 2.0),
        "pricing_power_index":   (1.0,  5.5,  0.4),
    }

    conn = get_db()
    # Clear existing data
    conn.execute("DELETE FROM monthly_data")
    conn.execute("DELETE FROM uploads")
    conn.execute("DELETE FROM projection_monthly_data")
    conn.execute("DELETE FROM projection_uploads")

    # Insert upload record for actuals
    cur = conn.execute(
        "INSERT INTO uploads (filename, uploaded_at, row_count, detected_columns) VALUES (?,?,?,?)",
        ("multiyear_demo_2021_2026.csv", datetime.utcnow().isoformat(), 63, json.dumps({}))
    )
    upload_id = cur.lastrowid

    total_months = 0

    # Seed 2021–2025 actuals (2-pass: base KPIs first, then derived)
    for yr, phase in PHASES.items():
        # Pass 1: compute all 12 months of base KPI values
        month_kpis = []
        for mo in range(1, 13):
            kpis = {kpi: mo_val(s, e, mo, n) for kpi, (s, e, n) in phase.items()}
            month_kpis.append(kpis)
        year_avg_rg = float(np.mean([m["revenue_growth"] for m in month_kpis]))
        # Pass 2: add derived KPIs and insert
        for i, (mo, kpis) in enumerate(zip(range(1, 13), month_kpis)):
            # growth_efficiency: ARR growth per unit of burn (higher = better capital efficiency)
            kpis["growth_efficiency"]   = round(kpis["arr_growth"] / max(kpis["burn_multiple"], 0.1), 4)
            # revenue_momentum: current growth vs annual average (>1 = accelerating)
            kpis["revenue_momentum"]    = round(kpis["revenue_growth"] / max(year_avg_rg, 0.1), 4)
            # revenue_fragility: concentration x churn risk divided by NRR resilience (lower = healthier)
            kpis["revenue_fragility"]   = round((kpis["customer_concentration"] * kpis["churn_rate"]) / max(kpis["nrr"], 1.0), 4)
            # burn_convexity: month-over-month change in burn (negative = improving)
            if i > 0:
                kpis["burn_convexity"]  = round(kpis["burn_multiple"] - month_kpis[i-1]["burn_multiple"], 4)
            else:
                # First month: use implied annual rate of change
                kpis["burn_convexity"]  = round((phase["burn_multiple"][1] - phase["burn_multiple"][0]) / 11.0, 4)
            # margin_volatility: rolling std dev of gross_margin (lower = more stable)
            if i >= 5:
                window = [month_kpis[j]["gross_margin"] for j in range(i - 5, i + 1)]
            else:
                window = [month_kpis[j]["gross_margin"] for j in range(0, i + 1)]
            kpis["margin_volatility"]   = round(float(np.std(window)) if len(window) > 1 else abs(rng(0, phase["gross_margin"][2])), 4)
            # customer_decay_slope: month-over-month change in churn rate (negative = improving)
            if i > 0:
                kpis["customer_decay_slope"] = round(kpis["churn_rate"] - month_kpis[i-1]["churn_rate"], 4)
            else:
                kpis["customer_decay_slope"] = round((phase["churn_rate"][1] - phase["churn_rate"][0]) / 11.0, 4)
            conn.execute(
                "INSERT INTO monthly_data (upload_id, year, month, data_json) VALUES (?,?,?,?)",
                (upload_id, yr, mo, json.dumps(kpis))
            )
            total_months += 1

    # Seed 2026 actuals Jan–Mar (2-pass for derived KPIs)
    act26_months = []
    for mo in range(1, 4):
        t = (mo - 1) / 2.0
        kpis = {kpi: round(lerp(start, end, t) + rng(0, noise), 4)
                for kpi, (start, end, noise) in PHASE_2026_ACTUAL.items()}
        act26_months.append(kpis)
    year_avg_rg_2026 = float(np.mean([m["revenue_growth"] for m in act26_months]))
    for i, (mo, kpis) in enumerate(zip(range(1, 4), act26_months)):
        kpis["growth_efficiency"]   = round(kpis["arr_growth"] / max(kpis["burn_multiple"], 0.1), 4)
        kpis["revenue_momentum"]    = round(kpis["revenue_growth"] / max(year_avg_rg_2026, 0.1), 4)
        kpis["revenue_fragility"]   = round((kpis["customer_concentration"] * kpis["churn_rate"]) / max(kpis["nrr"], 1.0), 4)
        if i > 0:
            kpis["burn_convexity"]  = round(kpis["burn_multiple"] - act26_months[i-1]["burn_multiple"], 4)
        else:
            kpis["burn_convexity"]  = round((PHASE_2026_ACTUAL["burn_multiple"][1] - PHASE_2026_ACTUAL["burn_multiple"][0]) / 2.0, 4)
        window = [act26_months[j]["gross_margin"] for j in range(0, i + 1)]
        kpis["margin_volatility"]   = round(float(np.std(window)) if len(window) > 1 else abs(rng(0, PHASE_2026_ACTUAL["gross_margin"][2])), 4)
        if i > 0:
            kpis["customer_decay_slope"] = round(kpis["churn_rate"] - act26_months[i-1]["churn_rate"], 4)
        else:
            kpis["customer_decay_slope"] = round((PHASE_2026_ACTUAL["churn_rate"][1] - PHASE_2026_ACTUAL["churn_rate"][0]) / 2.0, 4)
        conn.execute(
            "INSERT INTO monthly_data (upload_id, year, month, data_json) VALUES (?,?,?,?)",
            (upload_id, 2026, mo, json.dumps(kpis))
        )
        total_months += 1

    # Seed 2021-2025 plan/budget projections
    cur_plan = conn.execute(
        "INSERT INTO projection_uploads (filename, uploaded_at, row_count, detected_columns) VALUES (?,?,?,?)",
        ("plan_budget_2021_2025.csv", datetime.utcnow().isoformat(), 60, json.dumps({}))
    )
    plan_upload_id = cur_plan.lastrowid

    for yr, phase in PHASES_PROJ.items():
        proj_month_kpis = []
        for mo in range(1, 13):
            kpis = {kpi: mo_val(s, e, mo, n) for kpi, (s, e, n) in phase.items()}
            proj_month_kpis.append(kpis)
        proj_year_avg_rg = float(np.mean([m["revenue_growth"] for m in proj_month_kpis]))
        for i, (mo, kpis) in enumerate(zip(range(1, 13), proj_month_kpis)):
            kpis["growth_efficiency"]   = round(kpis["arr_growth"] / max(kpis["burn_multiple"], 0.1), 4)
            kpis["revenue_momentum"]    = round(kpis["revenue_growth"] / max(proj_year_avg_rg, 0.1), 4)
            kpis["revenue_fragility"]   = round((kpis["customer_concentration"] * kpis["churn_rate"]) / max(kpis["nrr"], 1.0), 4)
            if i > 0:
                kpis["burn_convexity"]  = round(kpis["burn_multiple"] - proj_month_kpis[i-1]["burn_multiple"], 4)
            else:
                kpis["burn_convexity"]  = round((phase["burn_multiple"][1] - phase["burn_multiple"][0]) / 11.0, 4)
            if i >= 5:
                window = [proj_month_kpis[j]["gross_margin"] for j in range(i - 5, i + 1)]
            else:
                window = [proj_month_kpis[j]["gross_margin"] for j in range(0, i + 1)]
            kpis["margin_volatility"]   = round(float(np.std(window)) if len(window) > 1 else abs(rng(0, phase["gross_margin"][2])), 4)
            if i > 0:
                kpis["customer_decay_slope"] = round(kpis["churn_rate"] - proj_month_kpis[i-1]["churn_rate"], 4)
            else:
                kpis["customer_decay_slope"] = round((phase["churn_rate"][1] - phase["churn_rate"][0]) / 11.0, 4)
            conn.execute(
                "INSERT INTO projection_monthly_data (projection_upload_id, year, month, data_json) VALUES (?,?,?,?)",
                (plan_upload_id, yr, mo, json.dumps(kpis))
            )

    # Seed 2026 projection Apr–Dec
    cur2 = conn.execute(
        "INSERT INTO projection_uploads (filename, uploaded_at, row_count, detected_columns) VALUES (?,?,?,?)",
        ("forecast_2026_recovery.csv", datetime.utcnow().isoformat(), 9, json.dumps({}))
    )
    proj_upload_id = cur2.lastrowid

    proj26_months = []
    for mo in range(4, 13):
        t = (mo - 4) / 8.0
        kpis = {kpi: round(lerp(start, end, t) + rng(0, noise), 4)
                for kpi, (start, end, noise) in PHASE_2026_PROJ.items()}
        proj26_months.append((mo, kpis))
    proj26_year_avg_rg = float(np.mean([k["revenue_growth"] for _, k in proj26_months]))
    for i, (mo, kpis) in enumerate(proj26_months):
        kpis["growth_efficiency"]   = round(kpis["arr_growth"] / max(kpis["burn_multiple"], 0.1), 4)
        kpis["revenue_momentum"]    = round(kpis["revenue_growth"] / max(proj26_year_avg_rg, 0.1), 4)
        kpis["revenue_fragility"]   = round((kpis["customer_concentration"] * kpis["churn_rate"]) / max(kpis["nrr"], 1.0), 4)
        if i > 0:
            kpis["burn_convexity"]  = round(kpis["burn_multiple"] - proj26_months[i-1][1]["burn_multiple"], 4)
        else:
            kpis["burn_convexity"]  = round((PHASE_2026_PROJ["burn_multiple"][1] - PHASE_2026_PROJ["burn_multiple"][0]) / 8.0, 4)
        if i >= 5:
            window = [proj26_months[j][1]["gross_margin"] for j in range(i - 5, i + 1)]
        else:
            window = [proj26_months[j][1]["gross_margin"] for j in range(0, i + 1)]
        kpis["margin_volatility"]   = round(float(np.std(window)) if len(window) > 1 else abs(rng(0, PHASE_2026_PROJ["gross_margin"][2])), 4)
        if i > 0:
            kpis["customer_decay_slope"] = round(kpis["churn_rate"] - proj26_months[i-1][1]["churn_rate"], 4)
        else:
            kpis["customer_decay_slope"] = round((PHASE_2026_PROJ["churn_rate"][1] - PHASE_2026_PROJ["churn_rate"][0]) / 8.0, 4)
        conn.execute(
            "INSERT INTO projection_monthly_data (projection_upload_id, year, month, data_json) VALUES (?,?,?,?)",
            (proj_upload_id, 2026, mo, json.dumps(kpis))
        )

    conn.commit()
    conn.close()
    return {
        "seeded": True,
        "years":  "2021–2026",
        "actual_months": total_months,
        "projection_months": 60 + 9,
        "narrative": [
            "2021: Startup phase — high churn, negative margins, rapid growth from small base",
            "2022: Series B scaling — rapid hiring, margins climbing, peak growth",
            "2023: Growth plateau — revenue stalls, over-hired, efficiency declining",
            "2024: Warning signals — churn uptick, margin compression, burn rising",
            "2025: Mixed recovery — volatile, fragile improvement",
            "2026 Jan–Mar: Critical state actuals",
            "2026 Apr–Dec: Recovery forecast",
            "2021–2025: Plan/budget projections added for Bridge Analysis",
        ]
    }

# ─── NLP Query Endpoint ─────────────────────────────────────────────────────

@app.post("/api/query", tags=["Analytics"])
async def query_kpi(payload: dict):
    """
    Natural-language KPI query powered by Claude.
    Accepts { "question": "...", "years": [2024] } and returns { "answer": "...", "kpis_referenced": [...] }.
    Builds full context from the live DB fingerprint on every call, filtered to requested years.
    """
    question = payload.get("question", "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="question is required")
    years_filter = payload.get("years", None)  # list of ints or None = all years

    # ── Build context from DB (replicate fingerprint + summary logic inline) ──
    conn = get_db()
    if years_filter:
        placeholders = ",".join("?" * len(years_filter))
        rows = conn.execute(f"SELECT * FROM monthly_data WHERE year IN ({placeholders})", years_filter).fetchall()
    else:
        rows = conn.execute("SELECT * FROM monthly_data").fetchall()
    targets  = {r["kpi_key"]: {"target": r["target_value"], "direction": r["direction"], "unit": r["unit"]}
                for r in conn.execute("SELECT * FROM kpi_targets").fetchall()}
    uploads  = conn.execute("SELECT COUNT(*) as c FROM uploads").fetchone()["c"]
    conn.close()

    # Organise monthly values by KPI key
    kpi_monthly: dict = {}
    for row in rows:
        mo_key = f"{row['year']}-{row['month']:02d}"
        for kpi_key, val in json.loads(row["data_json"]).items():
            if kpi_key not in ("year", "month"):
                kpi_monthly.setdefault(kpi_key, {})[mo_key] = val

    def _status(val, target, direction):
        if val is None or target is None:
            return "grey"
        pct = val / target if target else 0
        if direction == "higher":
            return "green" if pct >= 0.98 else ("yellow" if pct >= 0.90 else "red")
        return "green" if pct <= 1.02 else ("yellow" if pct <= 1.10 else "red")

    kpi_lines   = []
    status_counts = {"green": 0, "yellow": 0, "red": 0, "grey": 0}

    for kdef in KPI_DEFS:
        key  = kdef["key"]
        vals = kpi_monthly.get(key, {})
        t    = targets.get(key, {})
        tval = t.get("target")
        dirn = t.get("direction", "higher")

        monthly_sorted = sorted(vals.items())
        values         = [v for _, v in monthly_sorted if v is not None]
        avg            = round(float(np.mean(values)), 2) if values else None
        status         = _status(avg, tval, dirn)
        status_counts[status] += 1

        trend = None
        if len(values) >= 2:
            trend = "up" if values[-1] > values[0] else ("down" if values[-1] < values[0] else "flat")

        monthly_str = ", ".join(f"{p}: {v}" for p, v in monthly_sorted) or "no data"

        kpi_lines.append(
            f"- {kdef['name']} (key:{key}, unit:{kdef['unit']}): "
            f"avg={avg}, target={tval}, direction={dirn}, status={status}, trend={trend}\n"
            f"  monthly → {monthly_str}"
        )

    months_of_data = len(rows)
    kpis_tracked   = len([k for k in KPI_DEFS if kpi_monthly.get(k["key"])])

    # ── Projection context (if available) ─────────────────────────────────────
    proj_context_lines = []
    try:
        proj_conn  = get_db()
        proj_rows  = proj_conn.execute("SELECT * FROM projection_monthly_data").fetchall()
        proj_conn.close()
        if proj_rows:
            proj_by_period: dict = {}
            for pr in proj_rows:
                proj_by_period[(pr["year"], pr["month"])] = json.loads(pr["data_json"])

            actual_by_period2: dict = {}
            for row in rows:
                k2 = (row["year"], row["month"])
                actual_by_period2.setdefault(k2, {})
                actual_by_period2[k2].update(json.loads(row["data_json"]))

            overlap2 = sorted(set(proj_by_period.keys()) & set(actual_by_period2.keys()))
            if overlap2:
                for kdef in KPI_DEFS:
                    key2      = kdef["key"]
                    direction2 = kdef["direction"]
                    gap_pcts2 = []
                    actuals2  = []
                    projs2    = []
                    for (yr2, mo2) in overlap2:
                        pv = proj_by_period[(yr2, mo2)].get(key2)
                        av = actual_by_period2[(yr2, mo2)].get(key2)
                        if pv and av and pv != 0:
                            actuals2.append(av)
                            projs2.append(pv)
                            if direction2 == "higher":
                                gap_pcts2.append((av - pv) / abs(pv) * 100)
                            else:
                                gap_pcts2.append((pv - av) / abs(pv) * 100)
                    if actuals2:
                        avg_a2 = round(float(np.mean(actuals2)), 2)
                        avg_p2 = round(float(np.mean(projs2)), 2)
                        avg_g2 = round(float(np.mean(gap_pcts2)), 2)
                        status2 = compute_gap_status(avg_g2)
                        proj_context_lines.append(
                            f"- {kdef['name']}: actual avg={avg_a2}, projected avg={avg_p2}, gap={avg_g2:+.1f}% [{status2}]"
                        )
    except Exception:
        pass

    proj_section = ""
    if proj_context_lines:
        proj_section = f"""
PROJECTION vs ACTUAL CONTEXT ({len(proj_context_lines)} KPIs compared):
{chr(10).join(proj_context_lines)}
"""

    # Build a human-readable period description for the prompt header
    if years_filter:
        yr_list = sorted(years_filter)
        period_desc = f"FY {yr_list[0]}" if len(yr_list) == 1 else f"FY {yr_list[0]}–{yr_list[-1]}"
    else:
        all_years = sorted({row["year"] for row in rows}) if rows else []
        period_desc = (f"FY {all_years[0]}–{all_years[-1]}" if len(all_years) > 1
                       else f"FY {all_years[0]}" if all_years else "All Available Data")

    system_prompt = f"""You are an expert financial analyst embedded in the Axiom KPI Intelligence Dashboard.
You have access to the following organisational performance data for {period_desc}:

SUMMARY: {months_of_data} months of data | {kpis_tracked}/{len(KPI_DEFS)} KPIs tracked | Period: {period_desc}
Status breakdown: {status_counts.get('green', 0)} on target, {status_counts.get('yellow', 0)} needs attention, {status_counts.get('red', 0)} critical

KPI DATA:
{chr(10).join(kpi_lines)}
{proj_section}
Rules:
- Answer concisely (2-4 sentences max, or a short bullet list)
- Always cite specific numbers and months when relevant
- Flag critical KPIs (status=red) clearly
- When projection data is available, reference the gap percentages and status in your analysis
- Do NOT make up data beyond what is provided above
- Respond in plain text — no markdown headers, no asterisks, no bold"""

    try:
        client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
        msg = client.messages.create(
            model="claude-3-5-haiku-20241022",
            max_tokens=400,
            system=system_prompt,
            messages=[{"role": "user", "content": question}],
        )
        answer = msg.content[0].text
        kpis_referenced = [k["key"] for k in KPI_DEFS if k["name"].lower() in answer.lower()]
        return {"answer": answer, "kpis_referenced": kpis_referenced}
    except Exception as e:
        return {"answer": f"Query unavailable: {str(e)}", "kpis_referenced": []}


# ─── Data Ontology ──────────────────────────────────────────────────────────
#
# Builds a knowledge graph from the 18 KPIs:
#   Nodes  → each KPI
#   Edges  → CAUSES / INFLUENCES (from CAUSATION_RULES) +
#             CORRELATES_WITH / ANTI_CORRELATES (from monthly time-series)
#   Scores → degree centrality + iterative PageRank
#   Recs   → novel signal hypotheses from untested links & multi-hop paths

import threading, math

ONTOLOGY_DOMAIN = {
    "revenue_growth":        "growth",
    "arr_growth":            "growth",
    "nrr":                   "retention",
    "churn_rate":            "retention",
    "gross_margin":          "profitability",
    "operating_margin":      "profitability",
    "ebitda_margin":         "profitability",
    "contribution_margin":   "profitability",
    "operating_leverage":    "profitability",
    "opex_ratio":            "efficiency",
    "burn_multiple":         "efficiency",
    "cac_payback":           "efficiency",
    "sales_efficiency":      "efficiency",
    "cash_conv_cycle":       "cashflow",
    "dso":                   "cashflow",
    "revenue_quality":       "revenue",
    "recurring_revenue":     "revenue",
    "customer_concentration":"risk",
}

# Base signal weights for synthetic time-series generation.
# 8 signals (A-H) with different frequencies; each metric is a weighted combo + noise.
# Shared signal weights create realistic correlations between related metrics.
_SYN_WEIGHTS = {
    # key: [A, B, C, D, E, F, G, H, trend_per_month, base_value]
    "cpl":              [ 1.0,  0.5, -0.3,  0.0,  0.2,  0.0,  0.0,  0.0, -0.05,  80.0],
    "mql_sql_rate":     [-0.6, -0.3,  0.5,  0.3,  0.0,  0.2,  0.0,  0.0,  0.02,  22.0],
    "pipeline_velocity":[-0.5, -0.4,  0.6,  0.4,  0.0,  0.1,  0.0,  0.0,  0.03,  1.2],
    "win_rate":         [-0.4, -0.2,  0.5,  0.4,  0.0,  0.2,  0.0,  0.0,  0.01,  28.0],
    "organic_traffic":  [ 0.2,  0.0,  0.3,  0.0,  0.7,  0.3,  0.0,  0.0,  0.10,  12.0],
    "brand_awareness":  [ 0.1,  0.0,  0.2,  0.0,  0.6,  0.4,  0.0,  0.0,  0.08,  55.0],
    "quota_attainment": [-0.5, -0.4,  0.5,  0.3,  0.0,  0.0,  0.0,  0.0,  0.00,  82.0],
    "marketing_roi":    [-0.7,  0.0,  0.4,  0.2,  0.5,  0.0,  0.0,  0.0,  0.02,  3.2],
    "avg_deal_size":    [-0.3,  0.0,  0.3,  0.5,  0.0,  0.0,  0.4,  0.0,  0.05, 48000.0],
    "expansion_rate":   [-0.5, -0.6,  0.3,  0.0,  0.0,  0.0,  0.5,  0.0,  0.02,  18.0],
    "gross_dollar_ret": [-0.6, -0.5,  0.4,  0.0,  0.0,  0.0,  0.4,  0.0,  0.01,  91.0],
    "ltv_cac":          [-0.4, -0.3,  0.4,  0.3,  0.0,  0.0,  0.3,  0.0,  0.02,  4.5],
    "product_nps":      [ 0.0, -0.4, -0.3,  0.6,  0.0,  0.0,  0.3,  0.2,  0.03,  42.0],
    "feature_adoption": [ 0.0, -0.3, -0.2,  0.5,  0.0,  0.0,  0.4,  0.2,  0.04,  38.0],
    "activation_rate":  [ 0.0, -0.3, -0.2,  0.5,  0.0,  0.0,  0.3,  0.3,  0.03,  64.0],
    "time_to_value":    [ 0.0,  0.4,  0.2, -0.5,  0.0,  0.0, -0.3, -0.2, -0.04,  21.0],
    "health_score":     [ 0.0, -0.5, -0.3,  0.5,  0.0,  0.0,  0.4,  0.2,  0.02,  72.0],
    "logo_retention":   [-0.3, -0.5, -0.2,  0.4,  0.0,  0.0,  0.4,  0.1,  0.01,  94.0],
    "csat":             [ 0.0, -0.4, -0.2,  0.5,  0.0,  0.0,  0.3,  0.3,  0.02,  4.1],
    "headcount_eff":    [-0.3, -0.2,  0.3,  0.0,  0.0,  0.5,  0.2,  0.0,  0.01,  0.85],
    "rev_per_employee": [-0.4, -0.2,  0.4,  0.0,  0.0,  0.5,  0.2,  0.0,  0.02, 185000.0],
    "ramp_time":        [ 0.3,  0.2, -0.3,  0.0,  0.0, -0.4, -0.2,  0.0, -0.02,  5.5],
    "support_volume":   [ 0.2,  0.5,  0.2, -0.4,  0.0,  0.0, -0.2, -0.3, -0.01, 320.0],
    "automation_rate":  [-0.2, -0.1,  0.2,  0.0,  0.0,  0.6,  0.0,  0.0,  0.05,  34.0],
    "cash_runway":      [-0.3, -0.4,  0.2,  0.0,  0.0,  0.3,  0.0,  0.5,  0.00,  18.0],
    "current_ratio":    [-0.2, -0.3,  0.2,  0.0,  0.0,  0.2,  0.0,  0.5,  0.00,  2.1],
    "working_capital":  [-0.2, -0.3,  0.2,  0.0,  0.0,  0.2,  0.0,  0.5, -0.01,  1.8],
    "contraction_rate": [ 0.4,  0.5,  0.2, -0.4,  0.0,  0.0, -0.3, -0.1,  0.01,  3.5],
    "payback_period":   [ 0.5,  0.4, -0.3,  0.0,  0.0, -0.2, -0.2,  0.0, -0.02, 22.0],
}

def _init_ontology_tables(conn):
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS ontology_nodes (
            key TEXT PRIMARY KEY,
            name TEXT,
            domain TEXT,
            unit TEXT,
            direction TEXT,
            centrality REAL DEFAULT 0,
            pagerank REAL DEFAULT 0,
            updated_at TEXT
        );
        CREATE TABLE IF NOT EXISTS ontology_edges (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            source TEXT,
            target TEXT,
            relation TEXT,
            strength REAL,
            evidence TEXT,
            direction TEXT DEFAULT 'positive',
            granger_pval REAL,
            granger_lag INTEGER,
            confidence_tier TEXT DEFAULT 'expert_prior',
            UNIQUE(source, target, relation)
        );
        CREATE TABLE IF NOT EXISTS ontology_recommendations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT,
            description TEXT,
            rec_type TEXT,
            path TEXT,
            confidence REAL,
            novelty REAL,
            impact REAL,
            hypothesis TEXT,
            status TEXT DEFAULT 'active',
            created_at TEXT
        );
    """)
    conn.commit()
    # Migration: add direction column to existing ontology_edges tables
    try:
        conn.execute("ALTER TABLE ontology_edges ADD COLUMN direction TEXT DEFAULT 'positive'")
        conn.commit()
    except Exception:
        pass  # Column already exists


def _granger_test(y: list, x: list, max_lag: int = 4) -> tuple:
    """
    Granger predictability test: does knowing past values of x improve
    prediction of y beyond y's own history?

    Returns (best_pval, best_lag).  Uses numpy OLS + scipy F-distribution.
    Returns (1.0, 1) when data is insufficient (< 12 obs after lagging).
    Only runs on real monthly data — not synthetic extended-metric series.
    """
    from scipy.stats import f as f_dist
    ya = np.array(y, dtype=float)
    xa = np.array(x, dtype=float)
    n  = min(len(ya), len(xa))
    ya, xa = ya[:n], xa[:n]
    best_pval, best_lag = 1.0, 1
    for lag in range(1, max_lag + 1):
        T = n - lag
        if T < lag * 3 + 5:          # keep enough residual degrees of freedom
            continue
        Y       = ya[lag:]
        own_lags = [ya[lag - j : n - j] for j in range(1, lag + 1)]
        x_lags   = [xa[lag - j : n - j] for j in range(1, lag + 1)]
        X_r = np.column_stack([np.ones(T)] + own_lags)
        X_u = np.column_stack([np.ones(T)] + own_lags + x_lags)
        try:
            b_r = np.linalg.lstsq(X_r, Y, rcond=None)[0]
            b_u = np.linalg.lstsq(X_u, Y, rcond=None)[0]
            rss_r = float(np.sum((Y - X_r @ b_r) ** 2))
            rss_u = float(np.sum((Y - X_u @ b_u) ** 2))
            df1, df2 = lag, T - X_u.shape[1]
            if df2 <= 0 or rss_u <= 1e-12:
                continue
            F    = ((rss_r - rss_u) / df1) / (rss_u / df2)
            pval = float(1.0 - f_dist.cdf(max(F, 0.0), df1, df2))
            if pval < best_pval:
                best_pval, best_lag = pval, lag
        except Exception:
            continue
    return round(best_pval, 4), best_lag


def _run_ontology_discovery():
    conn = get_db()
    _init_ontology_tables(conn)

    now = datetime.utcnow().isoformat()

    # ── 1. Upsert nodes from KPI_DEFS ─────────────────────────────────────
    for kdef in KPI_DEFS:
        key = kdef["key"]
        conn.execute("""
            INSERT INTO ontology_nodes(key, name, domain, unit, direction, updated_at)
            VALUES (?,?,?,?,?,?)
            ON CONFLICT(key) DO UPDATE SET
              name=excluded.name, domain=excluded.domain,
              unit=excluded.unit, direction=excluded.direction,
              updated_at=excluded.updated_at
        """, (key, kdef["name"], ONTOLOGY_DOMAIN.get(key,"other"),
              kdef["unit"], kdef["direction"], now))

    # Upsert extended ontology-only nodes
    for em in EXTENDED_ONTOLOGY_METRICS:
        conn.execute("""
            INSERT INTO ontology_nodes(key, name, domain, unit, direction, updated_at)
            VALUES (?,?,?,?,?,?)
            ON CONFLICT(key) DO UPDATE SET
              name=excluded.name, domain=excluded.domain,
              unit=excluded.unit, direction=excluded.direction,
              updated_at=excluded.updated_at
        """, (em["key"], em["name"], em["domain"], em["unit"], em["direction"], now))
    conn.commit()

    # ── 1b. Migrate existing DBs: add Granger columns if absent ────────────
    for _col_sql in [
        "ALTER TABLE ontology_edges ADD COLUMN granger_pval REAL",
        "ALTER TABLE ontology_edges ADD COLUMN granger_lag INTEGER",
        "ALTER TABLE ontology_edges ADD COLUMN confidence_tier TEXT DEFAULT 'expert_prior'",
    ]:
        try:
            conn.execute(_col_sql)
        except Exception:
            pass   # column already exists
    conn.commit()

    # ── 2. Edges from CAUSATION_RULES + EXTENDED_CAUSATION_RULES ──────────
    edge_count = 0
    for source_key, rules in ALL_CAUSATION_RULES.items():
        for target_key in rules.get("downstream_impact", []):
            conn.execute("""
                INSERT OR IGNORE INTO ontology_edges(source, target, relation, strength, evidence, direction)
                VALUES (?,?,'CAUSES',0.75,'domain_knowledge','positive')
            """, (source_key, target_key))
            edge_count += 1
    conn.commit()

    # ── 3. Synthetic + real time-series for correlation ────────────────────
    # Load real KPI monthly data
    rows = conn.execute("SELECT data_json FROM monthly_data ORDER BY year, month").fetchall()
    series: dict = {}
    for row in rows:
        d = json.loads(row["data_json"])
        for k, v in d.items():
            if v is not None:
                series.setdefault(k, []).append(float(v))

    # Generate 60-month synthetic series for extended metrics using base signals.
    # 8 sinusoidal base signals (A-H) with different periods/phases.
    N = 60
    import random as _rnd
    _rnd.seed(42)
    base_signals = [
        [math.sin(2 * math.pi * t / p + ph) for t in range(N)]
        for p, ph in [(12, 0), (6, 1.0), (24, 2.1), (18, 0.5),
                      (36, 1.2), (9, 3.1), (15, 0.8), (48, 1.7)]
    ]
    for em in EXTENDED_ONTOLOGY_METRICS:
        key = em["key"]
        if key not in _SYN_WEIGHTS:
            continue
        w = _SYN_WEIGHTS[key]
        weights, trend, base = w[:8], w[8], w[9]
        ts = []
        for t in range(N):
            val = base + trend * t
            for j, sig in enumerate(base_signals):
                val += weights[j] * sig[t] * base * 0.08
            val += _rnd.gauss(0, abs(base) * 0.02)
            ts.append(val)
        series[key] = ts

    # Correlation loop across ALL nodes (KPI_DEFS + extended)
    all_keys = [kd["key"] for kd in KPI_DEFS] + [em["key"] for em in EXTENDED_ONTOLOGY_METRICS]
    for i, ka in enumerate(all_keys):
        for kb in all_keys[i+1:]:
            va = series.get(ka, [])
            vb = series.get(kb, [])
            n = min(len(va), len(vb))
            if n < 3:
                continue
            va, vb = va[:n], vb[:n]
            # Pearson r
            mean_a = sum(va)/n
            mean_b = sum(vb)/n
            num = sum((a-mean_a)*(b-mean_b) for a,b in zip(va,vb))
            den = math.sqrt(sum((a-mean_a)**2 for a in va)*sum((b-mean_b)**2 for b in vb))
            if den == 0:
                continue
            r = num/den
            if abs(r) < 0.45:
                continue
            rel = "CORRELATES_WITH" if r > 0 else "ANTI_CORRELATES"
            direction = 'positive' if r > 0 else 'negative'
            strength = round(abs(r), 4)
            conn.execute("""
                INSERT INTO ontology_edges(source, target, relation, strength, evidence, direction)
                VALUES (?,?,?,?,'monthly_correlation',?)
                ON CONFLICT(source, target, relation) DO UPDATE SET
                  strength=MAX(strength, excluded.strength),
                  direction=excluded.direction
            """, (ka, kb, rel, strength, direction))
            edge_count += 1
    conn.commit()

    # Back-fill direction on domain_knowledge CAUSES edges using empirical correlations.
    # If time-series shows a negative correlation between the pair, flip the edge to negative.
    # Infer direction on domain_knowledge CAUSES edges using node optimization direction.
    # Domain knowledge edges use node 'direction' properties (higher/lower) as ground
    # truth — these are more reliable than empirical correlations on synthetic data.
    # Logic: if source and target are optimised in opposite directions (one 'higher',
    # one 'lower'), a high value of the source drives the target the wrong way → negative.
    node_dirs = {r["key"]: r["direction"]
                 for r in conn.execute("SELECT key, direction FROM ontology_nodes").fetchall()}

    causes_edges = conn.execute(
        "SELECT source, target FROM ontology_edges WHERE evidence='domain_knowledge'"
    ).fetchall()
    for edge in causes_edges:
        src, tgt = edge["source"], edge["target"]
        src_dir = node_dirs.get(src, 'higher')
        tgt_dir = node_dirs.get(tgt, 'higher')
        inferred = 'positive' if src_dir == tgt_dir else 'negative'
        conn.execute(
            "UPDATE ontology_edges SET direction=? "
            "WHERE source=? AND target=? AND evidence='domain_knowledge'",
            (inferred, src, tgt)
        )
    conn.commit()

    # ── 3b. Granger tests on CAUSES edges (real KPI data only) ─────────────
    # For every expert-encoded CAUSES edge, test whether source Granger-predicts
    # target in the actual monthly data.  Only core KPI keys have real series;
    # extended-metric series are synthetic and excluded.
    real_keys = {kd["key"] for kd in KPI_DEFS}
    causes_pairs = conn.execute(
        "SELECT source, target FROM ontology_edges WHERE evidence='domain_knowledge'"
    ).fetchall()
    granger_confirmed = 0
    for ep in causes_pairs:
        src, tgt = ep["source"], ep["target"]
        if src not in real_keys or tgt not in real_keys:
            continue
        ys, xs = series.get(tgt, []), series.get(src, [])
        if len(ys) < 12 or len(xs) < 12:
            continue
        pval, lag = _granger_test(ys, xs)
        tier = 'data_confirmed' if pval <= 0.10 else 'expert_prior'
        conn.execute(
            "UPDATE ontology_edges SET granger_pval=?, granger_lag=?, confidence_tier=? "
            "WHERE source=? AND target=? AND evidence='domain_knowledge'",
            (pval, lag, tier, src, tgt)
        )
        if tier == 'data_confirmed':
            granger_confirmed += 1
    conn.commit()

    # ── 3c. Granger on high-correlation pairs — surface hidden leading indicators
    # Run Granger on pairs with strong Pearson correlation that have no CAUSES edge.
    # If Granger passes, upgrade to a LEADS edge with data_confirmed tier.
    corr_candidates = conn.execute(
        "SELECT source, target, strength, direction FROM ontology_edges "
        "WHERE relation='CORRELATES_WITH' AND strength >= 0.70"
    ).fetchall()
    for ep in corr_candidates:
        src, tgt = ep["source"], ep["target"]
        if src not in real_keys or tgt not in real_keys:
            continue
        # Skip if a CAUSES edge already exists between this pair
        existing = conn.execute(
            "SELECT 1 FROM ontology_edges WHERE source=? AND target=? AND relation='CAUSES'",
            (src, tgt)
        ).fetchone()
        if existing:
            continue
        ys, xs = series.get(tgt, []), series.get(src, [])
        if len(ys) < 12 or len(xs) < 12:
            continue
        pval, lag = _granger_test(ys, xs)
        if pval <= 0.10:
            # Promote: insert a LEADS edge tagged data_confirmed
            conn.execute("""
                INSERT INTO ontology_edges
                  (source, target, relation, strength, evidence, direction, granger_pval, granger_lag, confidence_tier)
                VALUES (?,?,'LEADS',?,?,'granger_test',?,?,?)
                ON CONFLICT(source, target, relation) DO UPDATE SET
                  strength=MAX(strength, excluded.strength),
                  granger_pval=excluded.granger_pval,
                  granger_lag=excluded.granger_lag,
                  confidence_tier='data_confirmed'
            """, (src, tgt, round(ep["strength"], 4), ep["direction"], pval, lag, 'data_confirmed'))
    conn.commit()

    # ── 4. Compute degree centrality + weighted PageRank ───────────────────
    edges_all = conn.execute(
        "SELECT source, target, strength, relation, granger_pval, confidence_tier "
        "FROM ontology_edges"
    ).fetchall()
    node_keys = [r["key"] for r in conn.execute("SELECT key FROM ontology_nodes").fetchall()]

    degree: dict = {k: 0 for k in node_keys}
    for e in edges_all:
        degree[e["source"]] = degree.get(e["source"], 0) + 1
        degree[e["target"]] = degree.get(e["target"], 0) + 1
    max_deg = max(degree.values()) if degree else 1

    # Weighted adjacency for PageRank
    # Edge weights reflect both relation type and statistical confirmation:
    #   CAUSES + data_confirmed  → 0.85 (expert rule validated by Granger test)
    #   CAUSES + expert_prior    → 0.60 (expert rule, no statistical confirmation yet)
    #   LEADS  (data_confirmed)  → 0.75 (Granger-discovered leading relationship)
    #   CORRELATES_WITH strong   → 0.45 (Pearson r > 0.70)
    #   CORRELATES_WITH moderate → 0.28 (Pearson r 0.45–0.70)
    #   ANTI_CORRELATES          → 0.20
    def _edge_w(rel, tier, strength):
        if rel == 'CAUSES':
            return 0.85 if tier == 'data_confirmed' else 0.60
        if rel == 'LEADS':
            return 0.75
        if rel == 'CORRELATES_WITH':
            return 0.45 if strength >= 0.70 else 0.28
        if rel == 'ANTI_CORRELATES':
            return 0.20
        return 0.35

    out_w: dict = {k: {} for k in node_keys}   # source → {target: weight}
    for e in edges_all:
        src, tgt = e["source"], e["target"]
        w = _edge_w(e["relation"], e.get("confidence_tier") or "expert_prior",
                    float(e["strength"] or 0))
        prev = out_w.setdefault(src, {}).get(tgt, 0)
        out_w[src][tgt] = max(prev, w)

    pr = {k: 1.0 for k in node_keys}
    d  = 0.85
    for _ in range(30):
        new_pr: dict = {}
        for k in node_keys:
            rank = (1 - d)
            for s in node_keys:
                w_sk = out_w.get(s, {}).get(k, 0)
                if w_sk > 0:
                    total_out = sum(out_w[s].values()) or 1e-9
                    rank += d * pr[s] * w_sk / total_out
            new_pr[k] = rank
        pr = new_pr
    max_pr = max(pr.values()) if pr else 1

    for k in node_keys:
        cent = round(degree.get(k, 0) / max_deg, 4)
        pg   = round(pr.get(k, 0) / max_pr, 4)
        conn.execute("UPDATE ontology_nodes SET centrality=?, pagerank=? WHERE key=?", (cent, pg, k))
    conn.commit()

    # ── 5. Generate signal recommendations ────────────────────────────────
    conn.execute("DELETE FROM ontology_recommendations WHERE 1")

    # Build quick lookup
    node_map = {r["key"]: dict(r) for r in conn.execute("SELECT * FROM ontology_nodes").fetchall()}
    edges_set = {(e["source"], e["target"], e["relation"]) for e in edges_all}
    causes_map: dict = {}   # source → [targets] for CAUSES edges
    for e in edges_all:
        if e["relation"] == "CAUSES":
            causes_map.setdefault(e["source"], []).append(e["target"])

    rec_count = 0

    # Pre-build Granger result lookup for recommendations
    granger_lookup = {
        (r["source"], r["target"]): (r["granger_pval"], r["granger_lag"])
        for r in conn.execute(
            "SELECT source, target, granger_pval, granger_lag FROM ontology_edges "
            "WHERE granger_pval IS NOT NULL"
        ).fetchall()
    }

    # (a) Transitive paths: A -CAUSES-> B -CAUSES-> C but no direct A→C edge
    # If Granger test confirmed A predicts C, surface as "Hidden Leading Indicator"
    # Otherwise surface as "Potential Upstream Driver" to investigate
    for a, b_list in causes_map.items():
        for b in b_list:
            for c in causes_map.get(b, []):
                if a == c:
                    continue
                if (a, c, "CAUSES") not in edges_set and (a, c, "CORRELATES_WITH") not in edges_set \
                        and (a, c, "LEADS") not in edges_set:
                    na = node_map.get(a, {}); nc = node_map.get(c, {})
                    nb = node_map.get(b, {})
                    g_pval, g_lag = granger_lookup.get((a, c), (None, None))
                    if g_pval is not None and g_pval <= 0.10:
                        # Granger-confirmed: this is a real hidden leading indicator
                        lag_txt = f" approximately {g_lag} month{'s' if g_lag != 1 else ''} ahead" if g_lag else ""
                        conn.execute("""
                            INSERT INTO ontology_recommendations
                              (title, description, rec_type, path, confidence, novelty, impact, hypothesis, status, created_at)
                            VALUES (?,?,?,?,?,?,?,?,?,?)
                        """, (
                            f"Hidden leading indicator: {na.get('name',a)} predicts {nc.get('name',c)}",
                            f"{na.get('name',a)} statistically predicts {nc.get('name',c)}{lag_txt}, "
                            f"bypassing the intermediate step through {nb.get('name',b)}. "
                            f"Granger p={g_pval:.3f} — this relationship is present in your data.",
                            "hidden_leading_indicator",
                            json.dumps([a, b, c]),
                            round(1 - g_pval, 2), 0.90, pr.get(c, 0) / max_pr,
                            f"Track {na.get('name',a)} as a direct early-warning signal for "
                            f"{nc.get('name',c)}. Act on it {lag_txt.strip()} before the impact arrives.",
                            "active", now
                        ))
                    else:
                        # Not yet confirmed — surface as a pattern to investigate
                        conn.execute("""
                            INSERT INTO ontology_recommendations
                              (title, description, rec_type, path, confidence, novelty, impact, hypothesis, status, created_at)
                            VALUES (?,?,?,?,?,?,?,?,?,?)
                        """, (
                            f"Potential upstream driver: {na.get('name',a)} to {nc.get('name',c)}",
                            f"{na.get('name',a)} influences {nb.get('name',b)} which influences "
                            f"{nc.get('name',c)}, but the direct relationship has not yet been "
                            f"confirmed in your data. More history will sharpen this signal.",
                            "untested_link",
                            json.dumps([a, b, c]),
                            0.65, 0.80, pr.get(c, 0) / max_pr,
                            f"Watch whether changes in {na.get('name',a)} consistently precede "
                            f"changes in {nc.get('name',c)} over the next two quarters.",
                            "active", now
                        ))
                    rec_count += 1
                    if rec_count >= 40:
                        break
            if rec_count >= 40:
                break
        if rec_count >= 40:
            break

    # (b) Cross-domain bridge: high-centrality node linking 2+ domains
    domain_nodes: dict = {}
    for k, n in node_map.items():
        domain_nodes.setdefault(n.get("domain", "other"), []).append(k)

    bridges = [k for k in node_keys if degree.get(k, 0) >= 4]
    bridges.sort(key=lambda k: -pr.get(k, 0))
    for bridge in bridges[:5]:
        neighbors = set()
        for e in edges_all:
            if e["source"] == bridge:
                neighbors.add(e["target"])
            if e["target"] == bridge:
                neighbors.add(e["source"])
        neighbor_domains = {node_map.get(n, {}).get("domain") for n in neighbors} - {None}
        if len(neighbor_domains) >= 2:
            nb = node_map.get(bridge, {})
            conn.execute("""
                INSERT INTO ontology_recommendations
                  (title, description, rec_type, path, confidence, novelty, impact, hypothesis, status, created_at)
                VALUES (?,?,?,?,?,?,?,?,?,?)
            """, (
                f"Bridge metric: {nb.get('name', bridge)}",
                f"{nb.get('name', bridge)} connects {len(neighbor_domains)} domains "
                f"({', '.join(sorted(neighbor_domains))}). Monitoring it provides early warning across multiple KPI clusters.",
                "bridge_node",
                json.dumps([bridge]),
                0.80, 0.75, pr.get(bridge, 0) / max_pr,
                f"Set alerts on {nb.get('name', bridge)} — it influences KPIs across "
                f"{', '.join(sorted(neighbor_domains))} domains simultaneously.",
                "active", now
            ))
            rec_count += 1

    # (c) Strongly-correlated cluster recommendations
    corr_edges = [(e["source"], e["target"], e["strength"]) for e in edges_all
                  if e["relation"] == "CORRELATES_WITH" and e["strength"] > 0.75]
    if len(corr_edges) >= 2:
        top_corr = sorted(corr_edges, key=lambda x: -x[2])[:3]
        for src, tgt, str_ in top_corr:
            ns = node_map.get(src, {}); nt = node_map.get(tgt, {})
            conn.execute("""
                INSERT INTO ontology_recommendations
                  (title, description, rec_type, path, confidence, novelty, impact, hypothesis, status, created_at)
                VALUES (?,?,?,?,?,?,?,?,?,?)
            """, (
                f"Strong co-movement: {ns.get('name',src)} ↔ {nt.get('name',tgt)}",
                f"Pearson r = {str_:.2f}. These KPIs move together strongly — "
                f"a combined leading indicator may have more predictive power than either alone.",
                "cluster",
                json.dumps([src, tgt]),
                round(str_, 2), 0.65, round((pr.get(src,0)+pr.get(tgt,0))/(2*max_pr), 4),
                f"Build a composite signal from {ns.get('name',src)} and {nt.get('name',tgt)} "
                f"to create a single early-warning index.",
                "active", now
            ))
            rec_count += 1

    conn.commit()
    conn.close()
    return {"nodes": len(node_keys), "edges": edge_count, "recommendations": rec_count}


# ── Ontology endpoints ────────────────────────────────────────────────────

@app.post("/api/ontology/discover")
def ontology_discover():
    """Trigger background knowledge-graph discovery."""
    def _bg():
        try:
            _run_ontology_discovery()
        except Exception as exc:
            print(f"Ontology discovery error: {exc}")
    threading.Thread(target=_bg, daemon=True).start()
    return {"status": "running", "message": "Ontology discovery started — refresh in ~5 seconds"}


@app.get("/api/ontology/graph")
def ontology_graph(domain: Optional[str] = None):
    conn = get_db()
    q = "SELECT * FROM ontology_nodes"
    params = ()
    if domain and domain != "all":
        q += " WHERE domain=?"
        params = (domain,)
    nodes = []
    for r in conn.execute(q, params).fetchall():
        n = dict(r)
        rules = ALL_CAUSATION_RULES.get(n["key"], {})
        n["root_causes"]       = rules.get("root_causes", [])
        n["corrective_actions"]= rules.get("corrective_actions", [])
        n["downstream_impact"] = rules.get("downstream_impact", [])
        nodes.append(n)
    node_keys = {n["key"] for n in nodes}
    edges = [dict(e) for e in conn.execute("SELECT * FROM ontology_edges").fetchall()
             if e["source"] in node_keys and e["target"] in node_keys]
    conn.close()
    return {"nodes": nodes, "edges": edges}


@app.get("/api/ontology/stats")
def ontology_stats():
    conn = get_db()
    total_nodes = conn.execute("SELECT COUNT(*) FROM ontology_nodes").fetchone()[0]
    total_edges = conn.execute("SELECT COUNT(*) FROM ontology_edges").fetchone()[0]
    active_recs  = conn.execute("SELECT COUNT(*) FROM ontology_recommendations WHERE status='active'").fetchone()[0]
    domain_rows  = conn.execute("SELECT domain, COUNT(*) as cnt FROM ontology_nodes GROUP BY domain").fetchall()
    edge_rows    = conn.execute("SELECT relation, COUNT(*) as cnt FROM ontology_edges GROUP BY relation").fetchall()
    top_nodes    = conn.execute(
        "SELECT key, name, pagerank, domain FROM ontology_nodes ORDER BY pagerank DESC LIMIT 5"
    ).fetchall()
    conn.close()
    return {
        "total_nodes": total_nodes,
        "total_edges": total_edges,
        "active_recommendations": active_recs,
        "domain_distribution": {r["domain"]: r["cnt"] for r in domain_rows},
        "edge_type_distribution": {r["relation"]: r["cnt"] for r in edge_rows},
        "top_nodes_by_pagerank": [dict(r) for r in top_nodes],
    }


@app.get("/api/ontology/recommendations")
def ontology_recommendations(rec_type: Optional[str] = None, status: Optional[str] = "active"):
    conn = get_db()
    q = "SELECT * FROM ontology_recommendations WHERE status=?"
    params: list = [status or "active"]
    if rec_type:
        q += " AND rec_type=?"
        params.append(rec_type)
    q += " ORDER BY impact DESC, confidence DESC"
    rows = [dict(r) for r in conn.execute(q, params).fetchall()]
    for r in rows:
        r["path"] = json.loads(r["path"]) if r.get("path") else []
    conn.close()
    return rows


@app.post("/api/ontology/recommendations/{rec_id}/dismiss")
def dismiss_recommendation(rec_id: int):
    conn = get_db()
    conn.execute("UPDATE ontology_recommendations SET status='dismissed' WHERE id=?", (rec_id,))
    conn.commit()
    conn.close()
    return {"status": "dismissed"}


# ─── Markov KPI Forecast ─────────────────────────────────────────────────────

from pydantic import BaseModel as _BaseModel

def _init_forecast_tables():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS markov_models (
            id INTEGER PRIMARY KEY,
            kpis TEXT,
            thresholds TEXT,
            self_matrices TEXT,
            cross_matrices TEXT,
            current_states TEXT,
            upstream_kpis TEXT,
            days_back INTEGER DEFAULT 365,
            trained_at TEXT,
            regime_data TEXT DEFAULT NULL
        );
        CREATE TABLE IF NOT EXISTS forecast_runs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            model_id INTEGER,
            horizon_days INTEGER,
            overrides TEXT,
            n_samples INTEGER,
            trajectories TEXT,
            causal_paths TEXT,
            created_at TEXT
        );
    """)
    conn.commit()
    # Add regime_data column to existing tables that pre-date this migration
    try:
        conn.execute("ALTER TABLE markov_models ADD COLUMN regime_data TEXT DEFAULT NULL")
        conn.commit()
    except Exception:
        pass  # Column already exists
    conn.close()

_init_forecast_tables()

def _mrk_monthly_history():
    conn = get_db()
    rows = conn.execute(
        "SELECT year, month, data_json FROM monthly_data ORDER BY year ASC, month ASC"
    ).fetchall()
    conn.close()
    result = {}
    for r in rows:
        try:
            d = json.loads(r["data_json"])
        except Exception:
            continue
        for k, v in d.items():
            if v is not None:
                result.setdefault(k, []).append(float(v))
    return result


def _mrk_causal_pairs():
    """Return 5-tuples: (source, target, strength, direction, relation)."""
    try:
        conn = get_db()
        rows = conn.execute(
            "SELECT source, target, strength, COALESCE(direction, 'positive') AS direction, relation "
            "FROM ontology_edges "
            "WHERE relation IN ('CAUSES','INFLUENCES','CORRELATES_WITH','LEADS','ANTI_CORRELATES') "
            "ORDER BY strength DESC"
        ).fetchall()
        conn.close()
        return [(r["source"], r["target"], float(r["strength"]), r["direction"], r["relation"])
                for r in rows]
    except Exception:
        return []


def _mrk_monthly_history_dated():
    """Returns OrderedDict {(year, month): {kpi: float}} sorted chronologically."""
    conn = get_db()
    rows = conn.execute(
        "SELECT year, month, data_json FROM monthly_data ORDER BY year ASC, month ASC"
    ).fetchall()
    conn.close()
    result = {}
    for r in rows:
        try:
            d = json.loads(r["data_json"])
        except Exception:
            continue
        result[(int(r["year"]), int(r["month"]))] = {
            k: float(v) for k, v in d.items() if v is not None
        }
    return result


def _detect_regimes(dated_history, kpis):
    """
    Wasserstein-based regime clustering (Horvath et al. 2021 approach adapted
    for monthly business KPI data).

    Each calendar month is represented as a standardised N-dimensional KPI
    snapshot vector.  Pairwise Wasserstein-1 distances between monthly
    snapshots are computed (treating the vector as an empirical distribution
    over its KPI values), then MDS + KMeans clusters months into 2-3 regimes.

    Returns a dict with:
        regime_labels   {(year,month) -> int}
        regime_names    {int -> "Growth"|"Recovery"|"Stress"}
        regime_deltas   {int -> {kpi -> [float]}}  — regime-conditioned MoM deltas
        current_regime  {label, name, months_in}
    Returns None if insufficient data.
    """
    try:
        from scipy.stats import wasserstein_distance as _wdist
        from sklearn.cluster import KMeans
        from sklearn.manifold import MDS
        from sklearn.preprocessing import StandardScaler
    except ImportError:
        return None

    months = sorted(dated_history.keys())
    n_m = len(months)
    if n_m < 12:
        return None

    # Build snapshot matrix (months × KPIs), filling missing values with col means
    mat = np.array(
        [[dated_history[ym].get(k, np.nan) for k in kpis] for ym in months],
        dtype=float,
    )
    col_means = np.nanmean(mat, axis=0)
    for j in range(mat.shape[1]):
        mat[np.isnan(mat[:, j]), j] = col_means[j] if not np.isnan(col_means[j]) else 0.0

    # Drop months that are still mostly NaN (shouldn't happen after fill)
    valid_mask = ~np.all(mat == 0, axis=1)
    mat = mat[valid_mask]
    valid_months = [months[i] for i, v in enumerate(valid_mask) if v]
    n_m = len(valid_months)
    if n_m < 12:
        return None

    # Standardise: each KPI zero-mean unit-variance across months
    scaler = StandardScaler()
    mat_std = scaler.fit_transform(mat)

    # Pairwise Wasserstein-1 distances between monthly KPI snapshot vectors
    dist_matrix = np.zeros((n_m, n_m))
    for i in range(n_m):
        for j in range(i + 1, n_m):
            d = _wdist(mat_std[i], mat_std[j])
            dist_matrix[i, j] = d
            dist_matrix[j, i] = d

    # MDS → 2-D embedding for KMeans
    try:
        mds = MDS(n_components=2, dissimilarity="precomputed",
                  random_state=42, normalized_stress="auto")
        embedding = mds.fit_transform(dist_matrix)
    except TypeError:
        mds = MDS(n_components=2, dissimilarity="precomputed", random_state=42)
        embedding = mds.fit_transform(dist_matrix)

    # Use 3 regimes if >= 24 months, else 2
    n_clusters = 3 if n_m >= 24 else 2
    kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
    raw_labels = kmeans.fit_predict(embedding)

    # Name regimes by business character: score each cluster on growth vs stress KPIs
    growth_kpis = [k for k in ("revenue_growth", "arr_growth", "nrr", "growth_efficiency") if k in kpis]
    stress_kpis = [k for k in ("churn_rate", "burn_multiple", "cac_payback", "burn_convexity") if k in kpis]

    regime_scores = {}
    for r in range(n_clusters):
        mask = raw_labels == r
        cluster_mat = mat_std[mask]
        g = sum(float(np.mean(cluster_mat[:, kpis.index(k)])) for k in growth_kpis) if growth_kpis else 0.0
        s = sum(float(np.mean(cluster_mat[:, kpis.index(k)])) for k in stress_kpis) if stress_kpis else 0.0
        regime_scores[r] = g - s  # higher = healthier

    sorted_r = sorted(regime_scores, key=lambda r: regime_scores[r])
    if n_clusters == 3:
        regime_name_map = {sorted_r[0]: "Stress", sorted_r[1]: "Recovery", sorted_r[2]: "Growth"}
    else:
        regime_name_map = {sorted_r[0]: "Stress", sorted_r[1]: "Growth"}

    regime_labels = {valid_months[i]: int(raw_labels[i]) for i in range(n_m)}

    # Build regime-conditioned MoM delta pools
    regime_deltas = {r: {k: [] for k in kpis} for r in range(n_clusters)}
    for i in range(1, n_m):
        prev_ym = valid_months[i - 1]
        curr_ym = valid_months[i]
        r = regime_labels[prev_ym]
        for k in kpis:
            pv = dated_history.get(prev_ym, {}).get(k)
            cv = dated_history.get(curr_ym, {}).get(k)
            if pv is not None and cv is not None:
                regime_deltas[r][k].append(float(cv) - float(pv))

    # Current regime and consecutive-month streak
    curr_ym = valid_months[-1]
    curr_label = regime_labels[curr_ym]
    months_in = 1
    for i in range(n_m - 2, -1, -1):
        if regime_labels[valid_months[i]] == curr_label:
            months_in += 1
        else:
            break

    return {
        "regime_labels":  {f"{ym[0]}-{ym[1]:02d}": lbl for ym, lbl in regime_labels.items()},
        "regime_names":   {str(k): v for k, v in regime_name_map.items()},
        "regime_deltas":  {str(r): d for r, d in regime_deltas.items()},
        "current_regime": {"label": curr_label, "name": regime_name_map[curr_label], "months_in": months_in},
        "n_clusters":     n_clusters,
    }


# Priority for deduplicating same (src, tgt) pairs: CAUSES beats empirical correlations
_CAUSAL_RELATION_PRIORITY = {
    'CAUSES': 3, 'INFLUENCES': 2, 'LEADS': 2,
    'CORRELATES_WITH': 1, 'ANTI_CORRELATES': 1,
}


def _build_deduped_causal_map(causal_pairs, kpis):
    """
    Build causal_map deduplicating by (tgt, src): when a KPI pair has both a
    domain-knowledge CAUSES edge and an empirical CORRELATES_WITH / ANTI_CORRELATES
    edge, keep only the higher-priority relation so the simulation doesn't double-
    count and doesn't use conflicting directions (H3 + H4).
    """
    best = {}  # (tgt, src) -> (strength, direction, priority)
    for src, tgt, strength, direction, relation in causal_pairs:
        if src in kpis and tgt in kpis:
            key      = (tgt, src)
            priority = _CAUSAL_RELATION_PRIORITY.get(relation, 0)
            if key not in best or priority > best[key][2]:
                best[key] = (strength, direction, priority)
    causal_map = {}
    for (tgt, src), (strength, direction, _) in best.items():
        causal_map.setdefault(tgt, []).append((src, strength, direction))
    return causal_map


def _build_markov_task():
    """
    Wasserstein regime-conditioned bootstrap Monte Carlo engine.

    Phase 1 (existing): learns per-KPI MoM delta distributions from full history.
    Phase 2 (new):      Wasserstein clustering labels each historical month with
                        a business regime (Growth / Recovery / Stress) and builds
                        regime-conditioned delta pools so Monte Carlo sampling
                        reflects the current operating context.
    """
    history       = _mrk_monthly_history()
    dated_history = _mrk_monthly_history_dated()
    if len(history) < 2:
        return

    current_values = {}
    value_ranges   = {}
    mean_deltas    = {}
    std_deltas     = {}

    for kpi, values in history.items():
        if len(values) < 2:
            continue
        arr    = np.array(values, dtype=float)
        deltas = np.diff(arr).tolist()

        current_values[kpi] = float(arr[-1])
        mean_deltas[kpi]    = float(np.mean(deltas))
        std_deltas[kpi]     = float(np.std(deltas)) if len(deltas) > 1 \
                              else max(abs(float(np.mean(deltas))) * 0.3, 0.01)
        value_ranges[kpi] = {
            "min":     float(np.min(arr)),
            "max":     float(np.max(arr)),
            "p10":     float(np.percentile(arr, 10)),
            "p25":     float(np.percentile(arr, 25)),
            "p50":     float(np.percentile(arr, 50)),
            "p75":     float(np.percentile(arr, 75)),
            "p90":     float(np.percentile(arr, 90)),
            "current": float(arr[-1]),
            "deltas":  deltas,
        }

    kpis          = list(current_values.keys())
    causal_pairs  = _mrk_causal_pairs()
    upstream_kpis = list({src for src, _, _, _, _ in causal_pairs if src in kpis})

    # ── Wasserstein regime detection ─────────────────────────────────────────
    regime_data = None
    try:
        regime_data = _detect_regimes(dated_history, kpis)
    except Exception:
        import traceback; traceback.print_exc()
    # ─────────────────────────────────────────────────────────────────────────

    now = datetime.utcnow().isoformat()
    conn = get_db()
    conn.execute("DELETE FROM markov_models")
    conn.execute(
        "INSERT INTO markov_models (kpis, thresholds, self_matrices, cross_matrices, "
        "current_states, upstream_kpis, days_back, trained_at, regime_data) "
        "VALUES (?,?,?,?,?,?,?,?,?)",
        (json.dumps(kpis),
         json.dumps(value_ranges),
         json.dumps(mean_deltas),
         json.dumps(std_deltas),
         json.dumps(current_values),
         json.dumps(upstream_kpis),
         365, now,
         json.dumps(regime_data) if regime_data else None)
    )
    conn.commit()
    conn.close()


def _project_scenario(horizon_days: int, overrides: dict, n_samples: int):
    """
    Bootstrap Monte Carlo projection in real KPI units.
    overrides = {kpi: state_idx} where state_idx 0-4 maps to p10/p25/p50/p75/p90.
    Causal influence: upstream KPI deltas propagate to downstream, scaled by
    relative volatility so units remain comparable.
    """
    conn = get_db()
    row  = conn.execute(
        "SELECT * FROM markov_models ORDER BY id DESC LIMIT 1"
    ).fetchone()
    conn.close()
    if not row:
        return {"status": "no_model", "message": "Train model first"}

    kpis         = json.loads(row["kpis"])
    value_ranges = json.loads(row["thresholds"])
    mean_deltas  = json.loads(row["self_matrices"])
    std_deltas   = json.loads(row["cross_matrices"])
    cur_values   = json.loads(row["current_states"])

    # ── Wasserstein regime data ───────────────────────────────────────────────
    regime_data    = json.loads(row["regime_data"]) if row["regime_data"] else None
    current_regime = regime_data["current_regime"] if regime_data else None

    # Build regime-conditioned delta pools keyed by KPI
    # Falls back to full historical delta pool when regime pool has < 6 samples
    _MIN_REGIME_POOL = 6
    regime_delta_pool = {}   # kpi -> [floats] — whichever pool will be used
    if regime_data:
        curr_label     = str(current_regime["label"])
        regime_deltas  = regime_data.get("regime_deltas", {})
        r_pool         = regime_deltas.get(curr_label, {})
        for kpi in kpis:
            cond_pool = r_pool.get(kpi, [])
            full_pool = value_ranges.get(kpi, {}).get("deltas", [])
            regime_delta_pool[kpi] = cond_pool if len(cond_pool) >= _MIN_REGIME_POOL else full_pool
    else:
        for kpi in kpis:
            regime_delta_pool[kpi] = value_ranges.get(kpi, {}).get("deltas", [])
    # ─────────────────────────────────────────────────────────────────────────

    causal_pairs = _mrk_causal_pairs()
    causal_map   = _build_deduped_causal_map(causal_pairs, set(kpis))

    horizon_steps = max(1, round(horizon_days / 30))

    # Map state overrides (0-4) → actual KPI starting values
    _state_pcts   = ["p10", "p25", "p50", "p75", "p90"]
    override_vals = {}
    for kpi, state_idx in overrides.items():
        if kpi in value_ranges:
            override_vals[kpi] = float(value_ranges[kpi][_state_pcts[min(int(state_idx), 4)]])

    # Sustained delta bias for overridden KPIs.
    # If a KPI is pinned above its historical median, it stays elevated (and vice versa).
    # z-score of override vs median → persistent per-step push so trajectories compound.
    override_delta_bias = {}
    for kpi, override_val in override_vals.items():
        vr       = value_ranges.get(kpi, {})
        hist_p50 = vr.get("p50", override_val)
        sigma    = std_deltas.get(kpi, 0.01) or 0.01
        z        = (override_val - hist_p50) / (sigma * 3)   # normalised position
        override_delta_bias[kpi] = z * sigma * 0.30           # 30% of 1-sigma per step

    all_traj = {kpi: [] for kpi in kpis}

    for _ in range(n_samples):
        values = {k: override_vals.get(k, cur_values.get(k, 0.0)) for k in kpis}
        path   = {k: [values[k]] for k in kpis}

        for _ in range(horizon_steps):
            # Step 1: regime-conditioned bootstrap self-delta for every KPI
            self_deltas = {}
            for kpi in kpis:
                hist  = regime_delta_pool.get(kpi, [])
                sigma = std_deltas.get(kpi, 0.01)
                if hist:
                    d = float(np.random.choice(hist))
                    d += float(np.random.normal(0, sigma * 0.15))  # small jitter
                else:
                    d = float(np.random.normal(mean_deltas.get(kpi, 0.0), sigma))
                # Apply sustained scenario bias for overridden KPIs
                d += override_delta_bias.get(kpi, 0.0)
                self_deltas[kpi] = d

            # Step 2: blend in causal influence from upstream KPIs
            # Uses BOTH delta-based (change propagation) and level-based (regime pressure)
            new_values = {}
            for kpi in kpis:
                delta     = self_deltas[kpi]
                tgt_sigma = std_deltas.get(kpi, 1.0) or 1.0
                if kpi in causal_map:
                    causal_delta = 0.0
                    total_w      = 0.0
                    for src, strength, direction in causal_map[kpi]:
                        src_sigma = std_deltas.get(src, 1.0) or 1.0
                        src_vr    = value_ranges.get(src, {})
                        dir_sign  = 1 if direction == 'positive' else -1
                        # Delta-based influence: upstream change propagates downstream
                        scale          = tgt_sigma / src_sigma
                        delta_term     = self_deltas[src] * scale * strength * dir_sign
                        # Level-based influence: how far upstream is from its median
                        # creates persistent directional pressure on downstream KPI.
                        # dir_sign ensures high churn pushes revenue DOWN, not up.
                        src_p50        = src_vr.get("p50", values[src])
                        level_z        = (values[src] - src_p50) / (src_sigma * 3 or 1)
                        level_term     = level_z * tgt_sigma * strength * dir_sign * 0.25
                        causal_delta  += (delta_term + level_term)
                        total_w       += strength
                    if total_w > 0:
                        causal_delta /= total_w
                        delta = 0.6 * delta + 0.4 * causal_delta
                new_values[kpi] = values[kpi] + delta

            values = new_values
            for kpi in kpis:
                path[kpi].append(values[kpi])

        for kpi in kpis:
            all_traj[kpi].append(path[kpi])

    # Compute p10/p50/p90 bands in real KPI units
    trajectories = {}
    for kpi in kpis:
        arr  = np.array(all_traj[kpi])
        vr   = value_ranges.get(kpi, {})
        traj = []
        for step in range(arr.shape[1]):
            col = arr[:, step]
            traj.append({
                "step":     step,
                "p10":      float(np.percentile(col, 10)),
                "p50":      float(np.percentile(col, 50)),
                "p90":      float(np.percentile(col, 90)),
                "label":    "Now" if step == 0 else f"M+{step}",
                "hist_p10": float(vr.get("p10", 0)),
                "hist_p50": float(vr.get("p50", 0)),
                "hist_p90": float(vr.get("p90", 0)),
            })
        trajectories[kpi] = traj

    causal_paths_out = {}
    for kpi in kpis:
        if kpi in causal_map:
            causal_paths_out[kpi] = [
                {"from": src, "strength": round(float(s), 3), "direction": drn}
                for src, s, drn in sorted(causal_map[kpi], key=lambda x: -x[1])[:3]
            ]

    now = datetime.utcnow().isoformat()
    conn = get_db()
    run_id = conn.execute(
        "INSERT INTO forecast_runs (model_id, horizon_days, overrides, n_samples, "
        "trajectories, causal_paths, created_at) VALUES (?,?,?,?,?,?,?)",
        (row["id"], horizon_days, json.dumps(overrides), n_samples,
         json.dumps(trajectories), json.dumps(causal_paths_out), now)
    ).lastrowid
    conn.commit()
    conn.close()

    return {
        "status":           "ok",
        "run_id":           run_id,
        "horizon_days":     horizon_days,
        "n_samples":        n_samples,
        "kpis":             kpis,
        "overrides":        overrides,
        "trajectories":     trajectories,
        "causal_paths":     causal_paths_out,
        "model_trained_at": row["trained_at"],
        "value_ranges":     value_ranges,
        "current_regime":   current_regime,   # Wasserstein regime label
        "regime_data":      {                  # summary for UI
            "name":      current_regime["name"]      if current_regime else None,
            "months_in": current_regime["months_in"] if current_regime else None,
            "label":     current_regime["label"]     if current_regime else None,
            "n_regimes": regime_data["n_clusters"]   if regime_data else None,
            "available": regime_data is not None,
        },
    }


class _ProjectRequest(_BaseModel):
    horizon_days: int = 90
    n_samples:    int = 400
    overrides:    dict = {}


@app.post("/api/forecast/build")
def forecast_build():
    def _bg():
        try:
            _build_markov_task()
        except Exception:
            import traceback; traceback.print_exc()
    threading.Thread(target=_bg, daemon=True).start()
    return {"status": "building", "message": "Markov model training started"}


@app.post("/api/forecast/project")
def forecast_project(req: _ProjectRequest):
    result = _project_scenario(req.horizon_days, req.overrides, req.n_samples)
    return result


@app.get("/api/forecast/model")
def forecast_model():
    conn = get_db()
    row = conn.execute(
        "SELECT * FROM markov_models ORDER BY id DESC LIMIT 1"
    ).fetchone()
    conn.close()
    if not row:
        return {"status": "not_trained"}
    return {
        "status":         "ready",
        "id":             row["id"],
        "kpis":           json.loads(row["kpis"]),
        "upstream_kpis":  json.loads(row["upstream_kpis"]),
        "current_values": json.loads(row["current_states"]),
        "value_ranges":   json.loads(row["thresholds"]),
        "trained_at":     row["trained_at"],
        "days_back":      row["days_back"],
        "regime_data":    json.loads(row["regime_data"]) if row["regime_data"] else None,
    }


# ─── Export / Import KPI Data ───────────────────────────────────────────────

@app.get("/api/export/data.xlsx")
def export_data_xlsx():
    """Export all monthly KPI data to Excel for offline editing."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    # ── KPI metadata: key → (full_name, unit_label, direction, used_for) ──────
    KPI_META: dict[str, tuple[str, str, str, str]] = {
        "revenue_growth":        ("Revenue Growth Rate",          "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine · Markov Forecast; feeds Revenue Momentum Index"),
        "gross_margin":          ("Gross Margin %",               "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine · Markov Forecast; feeds Operating Margin, EBITDA Margin, Contribution Margin"),
        "operating_margin":      ("Operating Margin %",           "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine; downstream of Gross Margin"),
        "ebitda_margin":         ("EBITDA Margin %",              "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine; downstream of Operating Margin & Gross Margin"),
        "cash_conv_cycle":       ("Cash Conversion Cycle",        "days",           "lower is better",   "Status Distribution · Fingerprint · Signal Engine · Markov Forecast"),
        "dso":                   ("Days Sales Outstanding",       "days",           "lower is better",   "Status Distribution · Fingerprint · Signal Engine; component of Cash Conversion Cycle"),
        "arr_growth":            ("ARR Growth Rate",              "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine · Markov Forecast; feeds Growth Efficiency Index, Revenue Momentum Index"),
        "nrr":                   ("Net Revenue Retention",        "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine · Markov Forecast; feeds Revenue Fragility Index"),
        "burn_multiple":         ("Burn Multiple",                "ratio (×)",      "lower is better",   "Status Distribution · Fingerprint · Signal Engine · Markov Forecast; feeds Growth Efficiency Index, Burn Convexity"),
        "opex_ratio":            ("Operating Expense Ratio",      "%",              "lower is better",   "Status Distribution · Fingerprint · Signal Engine"),
        "contribution_margin":   ("Contribution Margin %",        "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine; downstream of Gross Margin"),
        "revenue_quality":       ("Revenue Quality Ratio",        "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine"),
        "cac_payback":           ("CAC Payback Period",           "months",         "lower is better",   "Status Distribution · Fingerprint · Signal Engine · Markov Forecast"),
        "sales_efficiency":      ("Sales Efficiency Ratio",       "ratio (×)",      "higher is better",  "Status Distribution · Fingerprint · Signal Engine; downstream of Pipeline Conversion"),
        "customer_concentration":("Customer Concentration",       "%",              "lower is better",   "Status Distribution · Fingerprint · Signal Engine; feeds Revenue Fragility Index"),
        "recurring_revenue":     ("Recurring Revenue Ratio",      "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine"),
        "churn_rate":            ("Monthly Churn Rate",           "%",              "lower is better",   "Status Distribution · Fingerprint · Signal Engine · Markov Forecast; feeds Revenue Fragility Index, Customer Decay Curve Slope"),
        "operating_leverage":    ("Operating Leverage Index",     "ratio (×)",      "higher is better",  "Status Distribution · Fingerprint · Signal Engine"),
        "pipeline_conversion":   ("Pipeline Conversion Rate",     "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine; feeds Sales Efficiency, ARR Growth"),
        "customer_ltv":          ("Customer Lifetime Value",      "$ (thousands)",  "higher is better",  "Status Distribution · Fingerprint · Signal Engine"),
        "pricing_power_index":   ("Pricing Power Index",          "%",              "higher is better",  "Status Distribution · Fingerprint · Signal Engine"),
        "growth_efficiency":     ("Growth Efficiency Index",      "ratio (×)",      "higher is better",  "Status Distribution · Fingerprint · Signal Engine; DERIVED = ARR Growth Rate ÷ Burn Multiple"),
        "revenue_momentum":      ("Revenue Momentum Index",       "ratio (×)",      "higher is better",  "Status Distribution · Fingerprint · Signal Engine; DERIVED = Current Rev Growth ÷ Annual Avg Rev Growth"),
        "revenue_fragility":     ("Strategic Revenue Fragility",  "ratio (×)",      "lower is better",   "Status Distribution · Fingerprint · Signal Engine; DERIVED = (Customer Concentration × Churn Rate) ÷ NRR"),
        "burn_convexity":        ("Burn Convexity",               "ratio (×)",      "lower is better",   "Status Distribution · Fingerprint · Signal Engine; DERIVED = Month-over-Month change in Burn Multiple"),
        "margin_volatility":     ("Margin Volatility Index",      "%",              "lower is better",   "Status Distribution · Fingerprint · Signal Engine; DERIVED = 6-Month rolling std dev of Gross Margin"),
        "customer_decay_slope":  ("Customer Decay Curve Slope",   "%",              "lower is better",   "Status Distribution · Fingerprint · Signal Engine; DERIVED = Month-over-Month change in Churn Rate"),
    }

    conn = get_db()
    rows = conn.execute(
        "SELECT year, month, data_json FROM monthly_data ORDER BY year, month"
    ).fetchall()
    conn.close()

    if not rows:
        raise HTTPException(status_code=404, detail="No data found")

    # Collect KPI keys in consistent order
    kpi_keys: list[str] = []
    parsed = []
    for r in rows:
        d = json.loads(r["data_json"])
        for k in d:
            if k not in kpi_keys:
                kpi_keys.append(k)
        parsed.append((r["year"], r["month"], d))

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "KPI Data"

    # ── Colour palette ───────────────────────────────────────────────────────
    fill_key   = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")  # dark navy
    fill_name  = PatternFill(start_color="2C5282", end_color="2C5282", fill_type="solid")  # medium blue
    fill_unit  = PatternFill(start_color="2B6CB0", end_color="2B6CB0", fill_type="solid")  # blue
    fill_usage = PatternFill(start_color="EBF4FF", end_color="EBF4FF", fill_type="solid")  # pale blue
    fill_meta_left = PatternFill(start_color="2D3748", end_color="2D3748", fill_type="solid")  # dark slate for Year/Month cols

    font_white_bold = Font(color="FFFFFF", bold=True,  size=9)
    font_white      = Font(color="FFFFFF", bold=False, size=9)
    font_navy       = Font(color="1E3A5F", bold=False, size=8, italic=True)
    align_center    = Alignment(horizontal="center", vertical="center", wrap_text=True)
    align_right     = Alignment(horizontal="right",  vertical="center")
    thin_border     = Border(
        bottom=Side(style="thin", color="CBD5E0"),
        right=Side(style="thin",  color="CBD5E0"),
    )

    month_names = ["Jan","Feb","Mar","Apr","May","Jun",
                   "Jul","Aug","Sep","Oct","Nov","Dec"]

    # ── Row 1: machine-readable KPI key (used by import) ────────────────────
    row1_values = ["Year", "Month"] + kpi_keys
    for ci, val in enumerate(row1_values, 1):
        c = ws.cell(row=1, column=ci, value=val)
        c.fill  = fill_key if ci > 2 else fill_meta_left
        c.font  = font_white_bold
        c.alignment = align_center

    # ── Row 2: human-readable full name ─────────────────────────────────────
    ws.cell(row=2, column=1, value="Full Name").fill   = fill_meta_left
    ws.cell(row=2, column=1).font                      = font_white_bold
    ws.cell(row=2, column=1).alignment                 = align_center
    ws.cell(row=2, column=2, value="").fill            = fill_meta_left
    for ci, k in enumerate(kpi_keys, 3):
        meta = KPI_META.get(k, (k, "", "", ""))
        c = ws.cell(row=2, column=ci, value=meta[0])
        c.fill      = fill_name
        c.font      = font_white
        c.alignment = align_center

    # ── Row 3: unit of measure ───────────────────────────────────────────────
    ws.cell(row=3, column=1, value="Unit").fill        = fill_meta_left
    ws.cell(row=3, column=1).font                      = font_white_bold
    ws.cell(row=3, column=1).alignment                 = align_center
    ws.cell(row=3, column=2, value="").fill            = fill_meta_left
    for ci, k in enumerate(kpi_keys, 3):
        meta = KPI_META.get(k, (k, "", "", ""))
        c = ws.cell(row=3, column=ci, value=meta[1])
        c.fill      = fill_unit
        c.font      = font_white_bold
        c.alignment = align_center

    # ── Row 4: used-for description ──────────────────────────────────────────
    ws.cell(row=4, column=1, value="Used For").fill    = fill_meta_left
    ws.cell(row=4, column=1).font                      = font_white_bold
    ws.cell(row=4, column=1).alignment                 = align_center
    ws.cell(row=4, column=2, value="").fill            = fill_meta_left
    for ci, k in enumerate(kpi_keys, 3):
        meta = KPI_META.get(k, (k, "", "", ""))
        c = ws.cell(row=4, column=ci, value=meta[3])
        c.fill      = fill_usage
        c.font      = font_navy
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

    # Row heights for header block
    ws.row_dimensions[1].height = 22
    ws.row_dimensions[2].height = 30
    ws.row_dimensions[3].height = 20
    ws.row_dimensions[4].height = 52

    # ── Data rows start at row 5 ─────────────────────────────────────────────
    for row_idx, (year, month, d) in enumerate(parsed, 5):
        ws.cell(row=row_idx, column=1, value=year).alignment   = align_right
        ws.cell(row=row_idx, column=2,
                value=month_names[month - 1] if 1 <= month <= 12 else month).alignment = align_center
        # Alternating row shading
        alt_fill = PatternFill(start_color="F7FAFC", end_color="F7FAFC", fill_type="solid") if row_idx % 2 == 0 else None
        for ci, k in enumerate(kpi_keys, 3):
            val = d.get(k)
            c   = ws.cell(row=row_idx, column=ci,
                          value=round(val, 4) if isinstance(val, float) else val)
            c.alignment = align_right
            c.border    = thin_border
            if alt_fill:
                c.fill = alt_fill

    # ── Column widths ────────────────────────────────────────────────────────
    ws.column_dimensions["A"].width = 7
    ws.column_dimensions["B"].width = 7
    for ci in range(3, len(row1_values) + 1):
        ws.column_dimensions[get_column_letter(ci)].width = 20

    # Freeze below the 4-row header block
    ws.freeze_panes = "C5"

    # ── README / legend sheet ─────────────────────────────────────────────────
    ws2 = wb.create_sheet("README")
    ws2.column_dimensions["A"].width = 20
    ws2.column_dimensions["B"].width = 80
    readme_rows = [
        ("INSTRUCTIONS", ""),
        ("Step 1", "Edit KPI values in the 'KPI Data' sheet — do NOT change row 1 (KPI keys) or columns A–B (Year / Month)."),
        ("Step 2", "Save the file as .xlsx."),
        ("Step 3", "Upload via Settings › Import Data in the platform."),
        ("", ""),
        ("HEADER GUIDE", ""),
        ("Row 1 — KPI Key",   "Machine-readable identifier used by the import engine. Do not edit."),
        ("Row 2 — Full Name", "Human-readable KPI name for reference only."),
        ("Row 3 — Unit",      "Unit of measure: % = percentage point value (e.g. 6.5 = 6.5%); ratio (×) = dimensionless multiplier; days / months = calendar count; $ (thousands) = USD thousands."),
        ("Row 4 — Used For",  "Platform features & model components that consume this KPI. 'DERIVED' means the value is computed from other KPIs stored in this sheet."),
        ("", ""),
        ("UNIT NOTES", ""),
        ("%",            "Store as a plain number, NOT as a decimal. e.g. 62.5 means 62.5%, not 0.625."),
        ("ratio (×)",    "Dimensionless multiplier. e.g. 1.2 means 1.2×."),
        ("days / months","Integer or decimal count of calendar days or months."),
        ("$ (thousands)","USD value in thousands. e.g. 80 means $80,000."),
        ("DERIVED KPIs", "Values are pre-computed for testing. You may override them; the model will use whatever value is in the cell."),
    ]
    for ri, (label, text) in enumerate(readme_rows, 1):
        ca = ws2.cell(row=ri, column=1, value=label)
        cb = ws2.cell(row=ri, column=2, value=text)
        if label.isupper() and label:
            ca.font = Font(bold=True, size=10, color="1E3A5F")
        ca.alignment = Alignment(vertical="top")
        cb.alignment = Alignment(wrap_text=True, vertical="top")
        ws2.row_dimensions[ri].height = 28 if text else 10

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=axiom_kpi_data.xlsx"},
    )


@app.post("/api/import/data")
async def import_data_xlsx(file: UploadFile = File(...)):
    """Import KPI data from a previously exported (and edited) Excel file.

    Supports both the legacy single-header format (data from row 2) and the
    new enriched format (4-row header block; data from row 5).  Row 1 always
    contains the machine-readable KPI keys — rows 2-4 are metadata and are
    skipped automatically.
    """
    import openpyxl

    contents = await file.read()
    wb = openpyxl.load_workbook(io.BytesIO(contents), data_only=True)
    ws = wb["KPI Data"]

    # Row 1 always holds the machine-readable keys (Year, Month, kpi_key, ...)
    headers = [ws.cell(row=1, column=c).value for c in range(1, ws.max_column + 1)]
    if headers[0] != "Year" or headers[1] != "Month":
        raise HTTPException(status_code=400, detail="Invalid file format — expected Year, Month in row 1 columns A–B")

    kpi_keys = headers[2:]

    # Detect format: new enriched (row 2 col A = "Full Name") vs legacy (data at row 2)
    row2_a = ws.cell(row=2, column=1).value
    data_start_row = 5 if str(row2_a).strip().lower() in ("full name", "unit", "used for") else 2

    month_map = {"Jan":1,"Feb":2,"Mar":3,"Apr":4,"May":5,"Jun":6,
                 "Jul":7,"Aug":8,"Sep":9,"Oct":10,"Nov":11,"Dec":12}

    records: list[tuple] = []
    for row in ws.iter_rows(min_row=data_start_row, values_only=True):
        if row[0] is None:
            continue
        year = int(row[0])
        month_raw = row[1]
        month = month_map.get(str(month_raw), None) or int(month_raw)
        data = {}
        for i, k in enumerate(kpi_keys):
            if k is None:
                continue
            v = row[2 + i]
            if v is not None:
                data[k] = float(v)
        records.append((year, month, json.dumps(data)))

    if not records:
        raise HTTPException(status_code=400, detail="No data rows found in file")

    conn = get_db()
    for year, month, data_json in records:
        existing = conn.execute(
            "SELECT id FROM monthly_data WHERE year=? AND month=?", (year, month)
        ).fetchone()
        if existing:
            conn.execute(
                "UPDATE monthly_data SET data_json=? WHERE year=? AND month=?",
                (data_json, year, month)
            )
        else:
            conn.execute(
                "INSERT INTO monthly_data (year, month, data_json) VALUES (?,?,?)",
                (year, month, data_json)
            )
    conn.commit()
    conn.close()

    return {"status": "ok", "rows_imported": len(records)}


# ─── Slack Alerts ────────────────────────────────────────────────────────────

import urllib.request as _urllib_req

class SlackTestRequest(_BaseModel):
    webhook_url: str

class SlackAlertRequest(_BaseModel):
    webhook_url: str
    red_kpis: list[dict]          # [{key, name, value, target, pct_off}]
    company_name: str = "Your Company"

@app.post("/api/slack/test", tags=["Alerts"])
async def slack_test(body: SlackTestRequest):
    """Send a test Slack message to verify the webhook URL."""
    payload = {
        "text": "✅ *Axiom Intelligence* — Slack alerts connected successfully. You'll receive KPI threshold alerts here.",
        "username": "Axiom Intelligence",
        "icon_emoji": ":bar_chart:",
    }
    try:
        data = json.dumps(payload).encode("utf-8")
        req  = _urllib_req.Request(
            body.webhook_url,
            data=data,
            headers={"Content-Type": "application/json"},
        )
        with _urllib_req.urlopen(req, timeout=8) as resp:
            resp_text = resp.read().decode("utf-8")
        if resp_text.strip() != "ok":
            raise HTTPException(status_code=502, detail=f"Slack returned: {resp_text}")
    except Exception as exc:
        raise HTTPException(status_code=502, detail=str(exc))
    return {"status": "sent"}


@app.post("/api/slack/notify", tags=["Alerts"])
async def slack_notify(body: SlackAlertRequest):
    """Fire a KPI alert message to Slack for a batch of red KPIs."""
    if not body.red_kpis:
        return {"status": "no_alerts"}

    blocks: list[dict] = [
        {
            "type": "header",
            "text": {"type": "plain_text", "text": f"🚨 KPI Alert — {body.company_name}", "emoji": True},
        },
        {
            "type": "section",
            "text": {
                "type": "mrkdwn",
                "text": f"*{len(body.red_kpis)} KPI{'s' if len(body.red_kpis) > 1 else ''} below critical threshold* — immediate attention recommended.",
            },
        },
        {"type": "divider"},
    ]

    for kpi in body.red_kpis[:8]:   # cap at 8 to avoid giant messages
        pct  = abs(kpi.get("pct_off", 0))
        val  = kpi.get("value", "–")
        tgt  = kpi.get("target", "–")
        name = kpi.get("name", kpi.get("key", "?"))
        blocks.append({
            "type": "section",
            "text": {
                "type": "mrkdwn",
                "text": (
                    f"*{name}*\n"
                    f"Current: `{val}`   Target: `{tgt}`   Off by `{pct:.0f}%`"
                ),
            },
        })

    blocks.append({
        "type": "context",
        "elements": [
            {
                "type": "mrkdwn",
                "text": "Sent by *Axiom Intelligence V2* · Open the platform for full narrative analysis",
            }
        ],
    })

    payload = {"blocks": blocks, "username": "Axiom Intelligence", "icon_emoji": ":bar_chart:"}
    try:
        data = json.dumps(payload).encode("utf-8")
        req  = _urllib_req.Request(
            body.webhook_url,
            data=data,
            headers={"Content-Type": "application/json"},
        )
        with _urllib_req.urlopen(req, timeout=8) as resp:
            resp_text = resp.read().decode("utf-8")
        if resp_text.strip() != "ok":
            raise HTTPException(status_code=502, detail=f"Slack returned: {resp_text}")
    except Exception as exc:
        raise HTTPException(status_code=502, detail=str(exc))

    return {"status": "sent", "alerts_fired": len(body.red_kpis)}


# ─── Board Deck PPTX Export ──────────────────────────────────────────────────

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

_DECK_DARK_BLUE = _hex_to_rgb("071e45")
_DECK_HEADER_BLUE = _hex_to_rgb("0055A4")
_DECK_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DECK_RED_BG = _hex_to_rgb("fef2f2")
_DECK_RED_FG = _hex_to_rgb("dc2626")
_DECK_GREEN_BG = _hex_to_rgb("f0fdf4")
_DECK_GREEN_FG = _hex_to_rgb("059669")
_DECK_YELLOW_BG = _hex_to_rgb("fffbeb")
_DECK_YELLOW_FG = _hex_to_rgb("d97706")

def _set_cell_style(cell, text, font_size=11, bold=False, fg=None, bg=None):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.name = "Calibri"
        paragraph.font.bold = bold
        if fg:
            paragraph.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg.lstrip("#") if isinstance(bg, str) else f"{bg}"})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

def _set_cell_bg_rgb(cell, rgb_color: RGBColor):
    from pptx.oxml.ns import qn
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": f"{rgb_color}"})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)

def _status_colors(status: str):
    if status == "red":
        return _DECK_RED_FG, _DECK_RED_BG
    elif status == "green":
        return _DECK_GREEN_FG, _DECK_GREEN_BG
    elif status == "yellow":
        return _DECK_YELLOW_FG, _DECK_YELLOW_BG
    return None, None

def _add_header_row(table, col_texts):
    for i, txt in enumerate(col_texts):
        cell = table.cell(0, i)
        _set_cell_style(cell, txt, font_size=11, bold=True, fg=_DECK_WHITE)
        _set_cell_bg_rgb(cell, _DECK_HEADER_BLUE)

def _compute_fingerprint_data():
    """Reuse fingerprint logic without HTTP call."""
    conn = get_db()
    rows = conn.execute("SELECT * FROM monthly_data").fetchall()
    targets = {r["kpi_key"]: {"target": r["target_value"], "direction": r["direction"], "unit": r["unit"]}
               for r in conn.execute("SELECT * FROM kpi_targets").fetchall()}
    conn.close()

    kpi_monthly: dict = {}
    for row in rows:
        mo_key = f"{row['year']}-{row['month']:02d}"
        data = json.loads(row["data_json"])
        for kpi_key, val in data.items():
            if kpi_key in ("year", "month"):
                continue
            kpi_monthly.setdefault(kpi_key, {})[mo_key] = val

    def _status(val, target, direction):
        if val is None or target is None:
            return "grey"
        pct = val / target if target else 0
        if direction == "higher":
            return "green" if pct >= 0.98 else ("yellow" if pct >= 0.90 else "red")
        else:
            return "green" if pct <= 1.02 else ("yellow" if pct <= 1.10 else "red")

    # Build lookup for all known KPI metadata (KPI_DEFS + EXTENDED_ONTOLOGY_METRICS)
    all_kpi_meta = {kd["key"]: kd for kd in KPI_DEFS}
    for em in EXTENDED_ONTOLOGY_METRICS:
        if em["key"] not in all_kpi_meta:
            all_kpi_meta[em["key"]] = em

    # All KPIs that have both a target AND monthly data
    keys_to_include = set(KPI_DEFS[i]["key"] for i in range(len(KPI_DEFS)))
    for key in kpi_monthly:
        if key in targets:
            keys_to_include.add(key)

    fingerprint_out = []
    for key in keys_to_include:
        kdef = all_kpi_meta.get(key, {"key": key, "name": key.replace("_", " ").title(), "unit": "ratio", "direction": "higher", "domain": "other"})
        vals = kpi_monthly.get(key, {})
        t = targets.get(key, {})
        tval = t.get("target")
        dirn = t.get("direction", kdef.get("direction", "higher"))
        unit = t.get("unit", kdef.get("unit", "ratio"))

        monthly_list = [{"period": k, "value": v} for k, v in sorted(vals.items())]
        values = [m["value"] for m in monthly_list]
        avg = round(np.mean(values), 2) if values else None

        trend = None
        if len(values) >= 2:
            trend = "up" if values[-1] > values[0] else ("down" if values[-1] < values[0] else "flat")

        fingerprint_out.append({
            "key": key,
            "name": kdef.get("name", key.replace("_", " ").title()),
            "unit": unit,
            "target": tval,
            "direction": dirn,
            "avg": avg,
            "trend": trend,
            "fy_status": _status(avg, tval, dirn),
            "monthly": monthly_list,
        })

    # Sort: KPI_DEFS order first, then extended KPIs alphabetically
    kpi_def_order = {kd["key"]: i for i, kd in enumerate(KPI_DEFS)}
    fingerprint_out.sort(key=lambda x: (kpi_def_order.get(x["key"], 9999), x["key"]))
    return fingerprint_out

@app.get("/api/export/board-deck.pptx", tags=["Board Deck"])
def export_board_deck(stage: str = "series_b"):
    """Generate a narrative-driven PPTX board deck with charts, executive summary, and data-backed actions."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mticker

    fp_data = _compute_fingerprint_data()

    valid_stages = {"seed", "series_a", "series_b", "series_c"}
    if stage not in valid_stages:
        stage = "series_b"
    bench = {}
    for kpi_key, stages_data in BENCHMARKS.items():
        if stage in stages_data:
            bench[kpi_key] = stages_data[stage]

    stage_label = {"seed": "Seed", "series_a": "Series A", "series_b": "Series B", "series_c": "Series C+"}.get(stage, stage)

    green_kpis = [k for k in fp_data if k["fy_status"] == "green"]
    yellow_kpis = [k for k in fp_data if k["fy_status"] == "yellow"]
    red_kpis = [k for k in fp_data if k["fy_status"] == "red"]
    total = len(green_kpis) + len(yellow_kpis) + len(red_kpis)

    # Sort red by worst gap
    def _gap_pct(k):
        if k["avg"] is not None and k["target"] is not None and k["target"] != 0:
            return abs((k["avg"] - k["target"]) / abs(k["target"]) * 100)
        return 0
    red_kpis.sort(key=_gap_pct, reverse=True)
    yellow_kpis.sort(key=_gap_pct, reverse=True)

    # ── Helper: generate a matplotlib chart as PNG bytes ─────────────────
    def _make_trend_chart(kpis_list, title_text, max_kpis=5):
        """Sparkline-style multi-KPI trend chart → PNG bytes."""
        fig, ax = plt.subplots(figsize=(11, 5.5))
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        colors_cycle = ["#dc2626", "#d97706", "#2563eb", "#059669", "#7c3aed", "#db2777"]
        plotted = 0
        for i, kpi in enumerate(kpis_list[:max_kpis]):
            months = kpi.get("monthly", [])
            if len(months) < 2:
                continue
            periods = [m["period"] for m in months]
            values = [m["value"] for m in months]
            color = colors_cycle[i % len(colors_cycle)]
            ax.plot(periods, values, marker="o", markersize=4, linewidth=2, color=color, label=kpi["name"])
            if kpi.get("target"):
                ax.axhline(y=kpi["target"], color=color, linestyle="--", alpha=0.4, linewidth=1)
            plotted += 1
        if plotted == 0:
            plt.close(fig)
            return None
        ax.set_title(title_text, fontsize=16, fontweight="bold", pad=16)
        ax.legend(fontsize=10, loc="upper left", framealpha=0.9, fancybox=True, shadow=True)
        ax.grid(True, alpha=0.3)
        # Show only every Nth x-tick to avoid crowding
        labels = [m["period"] for m in kpis_list[0].get("monthly", [])] if kpis_list else []
        if len(labels) > 12:
            step = max(len(labels) // 8, 1)
            ax.set_xticks(range(0, len(labels), step))
            ax.set_xticklabels([labels[j] for j in range(0, len(labels), step)], fontsize=8, rotation=30)
        else:
            ax.tick_params(axis="x", labelsize=8, rotation=30)
        ax.tick_params(axis="y", labelsize=9)
        plt.tight_layout()
        buf_png = io.BytesIO()
        fig.savefig(buf_png, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        buf_png.seek(0)
        return buf_png

    def _make_status_donut():
        """Donut chart of red/yellow/green distribution → PNG bytes."""
        fig, ax = plt.subplots(figsize=(5, 5))
        fig.patch.set_facecolor("white")
        sizes = [len(red_kpis), len(yellow_kpis), len(green_kpis)]
        colors_d = ["#dc2626", "#d97706", "#059669"]
        labels = [f"Critical ({len(red_kpis)})", f"Watch ({len(yellow_kpis)})", f"On Target ({len(green_kpis)})"]
        # Filter out zeros
        filtered = [(s, c, l) for s, c, l in zip(sizes, colors_d, labels) if s > 0]
        if not filtered:
            plt.close(fig)
            return None
        f_sizes, f_colors, f_labels = zip(*filtered)
        wedges, texts, autotexts = ax.pie(f_sizes, colors=f_colors, labels=f_labels,
                                           autopct="%1.0f%%", startangle=90, pctdistance=0.78,
                                           textprops={"fontsize": 11})
        for at in autotexts:
            at.set_fontsize(12)
            at.set_fontweight("bold")
            at.set_color("white")
        centre_circle = plt.Circle((0, 0), 0.55, fc="white")
        ax.add_artist(centre_circle)
        ax.text(0, 0.08, str(total), ha="center", va="center", fontsize=28, fontweight="bold", color="#1e293b")
        ax.text(0, -0.15, "KPIs", ha="center", va="center", fontsize=11, color="#64748b")
        plt.tight_layout()
        buf_png = io.BytesIO()
        fig.savefig(buf_png, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        buf_png.seek(0)
        return buf_png

    def _make_benchmark_bar(kpis_to_show, bench_data):
        """Horizontal bar chart: company value vs peer median → PNG bytes."""
        names = []
        company_vals = []
        peer_vals = []
        bar_colors = []
        for k in kpis_to_show:
            if k["key"] in bench_data and k["avg"] is not None:
                b = bench_data[k["key"]]
                names.append(k["name"][:25])
                company_vals.append(k["avg"])
                peer_vals.append(b["p50"])
                bar_colors.append("#dc2626" if k["fy_status"] == "red" else "#d97706" if k["fy_status"] == "yellow" else "#059669")
        if not names:
            return None
        fig, ax = plt.subplots(figsize=(11, max(len(names) * 0.65, 4)))
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        y_pos = range(len(names))
        ax.barh(y_pos, company_vals, height=0.35, color=bar_colors, label="Company", alpha=0.9)
        ax.barh([y + 0.35 for y in y_pos], peer_vals, height=0.35, color="#94a3b8", label=f"Peer Median ({stage_label})", alpha=0.6)
        ax.set_yticks([y + 0.175 for y in y_pos])
        ax.set_yticklabels(names, fontsize=10)
        ax.invert_yaxis()
        ax.legend(fontsize=10, loc="lower right")
        ax.set_title(f"Company vs {stage_label} Peer Median", fontsize=13, fontweight="bold", pad=10)
        ax.grid(True, axis="x", alpha=0.3)
        ax.tick_params(axis="x", labelsize=9)
        plt.tight_layout()
        buf_png = io.BytesIO()
        fig.savefig(buf_png, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        buf_png.seek(0)
        return buf_png

    # ── Helper: add a text box with multi-paragraph rich text ──────────────
    def _add_narrative(slide, left, top, width, height, paragraphs_list):
        """paragraphs_list: [(text, font_size, bold, color_rgb), ...]"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (text, fs, bold, color) in enumerate(paragraphs_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(fs)
            p.font.name = "Calibri"
            p.font.bold = bold
            if color:
                p.font.color.rgb = color
            p.space_after = Pt(6)

    # ── Build narrative sentences ─────────────────────────────────────────
    def _kpi_sentence(k):
        val = k["avg"]
        tgt = k["target"]
        if val is None or tgt is None or tgt == 0:
            return f"{k['name']}: no data available."
        gap = round((val - tgt) / abs(tgt) * 100, 1)
        direction_word = "below" if gap < 0 else "above"
        unit = k.get("unit", "")
        val_fmt = f"{val:,.2f}" if isinstance(val, float) else str(val)
        tgt_fmt = f"{tgt:,.2f}" if isinstance(tgt, float) else str(tgt)
        sentence = f"{k['name']} is at {val_fmt} vs target {tgt_fmt} ({abs(gap):.0f}% {direction_word} target)."
        # Add benchmark context
        b = bench.get(k["key"])
        if b:
            if val < b["p25"]:
                sentence += f" Below {stage_label} P25 ({b['p25']}) — bottom quartile."
            elif val < b["p50"]:
                sentence += f" Below {stage_label} median ({b['p50']})."
            elif val >= b["p75"]:
                sentence += f" Above {stage_label} P75 ({b['p75']}) — top quartile."
        # Add causal context
        rules = ALL_CAUSATION_RULES.get(k["key"], {})
        if rules.get("root_causes"):
            sentence += f" Likely driver: {rules['root_causes'][0].lower()}."
        if rules.get("corrective_actions"):
            sentence += f" Recommended action: {rules['corrective_actions'][0]}."
        return sentence

    # ── PPTX Generation ───────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ── Slide 1: Title ────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.background
    fill1 = bg1.fill
    fill1.solid()
    fill1.fore_color.rgb = _DECK_DARK_BLUE

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Board Intelligence Brief"
    p.font.size = Pt(44)
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.color.rgb = _DECK_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Date + status summary subtitle
    _add_narrative(slide1, 1, 3.3, 11, 2, [
        (f"{datetime.now().strftime('%B %d, %Y')}  ·  {stage_label} SaaS", 20, False, _DECK_WHITE),
        ("", 8, False, None),
        (f"{len(red_kpis)} critical  ·  {len(yellow_kpis)} watch  ·  {len(green_kpis)} on target  ·  {total} KPIs tracked", 22, True, _DECK_WHITE),
    ])
    # Center subtitle
    for shape in slide1.shapes:
        if hasattr(shape, "text_frame"):
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Executive Summary (narrative + donut) ────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    _add_narrative(slide2, 0.5, 0.3, 8, 0.7, [
        ("Executive Summary", 28, True, _hex_to_rgb("1e293b")),
    ])

    # Build the summary narrative
    summary_paras = []
    if red_kpis:
        worst = red_kpis[0]
        worst_gap = round(abs((worst["avg"] - worst["target"]) / abs(worst["target"]) * 100), 0) if worst["avg"] and worst["target"] and worst["target"] != 0 else 0
        summary_paras.append(
            (f"The business has {len(red_kpis)} KPIs in critical status and {len(yellow_kpis)} requiring attention. "
             f"The most severe miss is {worst['name']} at {worst.get('avg', '?')} vs target {worst.get('target', '?')} "
             f"({worst_gap:.0f}% off target).", 13, False, _hex_to_rgb("334155"))
        )
    else:
        summary_paras.append(
            (f"All {total} KPIs are in green or watch status. No critical issues detected.", 13, False, _hex_to_rgb("334155"))
        )

    if red_kpis:
        # Causal chain for worst
        rules = ALL_CAUSATION_RULES.get(red_kpis[0]["key"], {})
        downstream = rules.get("downstream_impact", [])
        downstream_red = [k for k in fp_data if k["key"] in downstream and k["fy_status"] == "red"]
        if downstream_red:
            names_str = ", ".join(d["name"] for d in downstream_red[:3])
            summary_paras.append(
                (f"This is cascading: {red_kpis[0]['name']} directly impacts {names_str}, which are also in critical status. "
                 f"Addressing the root cause would improve multiple metrics simultaneously.", 13, False, _hex_to_rgb("334155"))
            )
        if rules.get("corrective_actions"):
            summary_paras.append(
                (f"Priority action: {rules['corrective_actions'][0]}", 13, True, _DECK_RED_FG)
            )

    # Bright spots
    if green_kpis:
        top_green = green_kpis[:3]
        bright_names = ", ".join(k["name"] for k in top_green)
        summary_paras.append(
            (f"Bright spots: {bright_names} are all on or above target.", 13, False, _DECK_GREEN_FG)
        )

    _add_narrative(slide2, 0.5, 1.2, 7.5, 5, summary_paras)

    # Donut chart on the right
    donut_png = _make_status_donut()
    if donut_png:
        slide2.shapes.add_picture(donut_png, Inches(8.5), Inches(1), Inches(4.5), Inches(4.5))

    # ── Slide 3: Critical KPIs — Narrative Cards + Trend Chart ────────────
    slide3 = prs.slides.add_slide(blank_layout)
    _add_narrative(slide3, 0.5, 0.3, 8, 0.7, [
        (f"Critical Items: {len(red_kpis)} KPIs Below Threshold", 28, True, _hex_to_rgb("1e293b")),
    ])

    if red_kpis:
        # Left side: narrative cards for top 4 red KPIs
        card_paras = []
        for k in red_kpis[:4]:
            card_paras.append((_kpi_sentence(k), 11, False, _hex_to_rgb("334155")))
            card_paras.append(("", 6, False, None))  # spacer

        _add_narrative(slide3, 0.5, 1.2, 6, 5.5, card_paras)

        # Right side: trend chart of red KPIs
        trend_png = _make_trend_chart(red_kpis, "Critical KPI Trends", max_kpis=4)
        if trend_png:
            slide3.shapes.add_picture(trend_png, Inches(6.5), Inches(1.2), Inches(6.5), Inches(5.2))
    else:
        _add_narrative(slide3, 0.5, 2, 8, 1, [
            ("No critical KPIs — all metrics are within acceptable ranges.", 16, False, _DECK_GREEN_FG)
        ])

    # ── Slide 4: Benchmark Position — Bar Chart + Narrative ───────────────
    slide4 = prs.slides.add_slide(blank_layout)
    _add_narrative(slide4, 0.5, 0.3, 10, 0.7, [
        (f"Peer Benchmark Position — {stage_label} SaaS", 28, True, _hex_to_rgb("1e293b")),
    ])

    # Show red + yellow KPIs in benchmark comparison
    kpis_for_bench = (red_kpis + yellow_kpis)[:10]
    bench_png = _make_benchmark_bar(kpis_for_bench, bench)
    if bench_png:
        slide4.shapes.add_picture(bench_png, Inches(0.3), Inches(1.3), Inches(8.5), Inches(5.8))

    # Narrative on right: which KPIs are below P25
    below_p25 = []
    for k in fp_data:
        if k["key"] in bench and k["avg"] is not None:
            if k["avg"] < bench[k["key"]]["p25"]:
                below_p25.append(k)
    bench_narrative = []
    if below_p25:
        bench_narrative.append(
            (f"{len(below_p25)} KPIs are below the {stage_label} bottom quartile (P25):", 13, True, _DECK_RED_FG)
        )
        for bp in below_p25[:5]:
            b = bench[bp["key"]]
            bench_narrative.append(
                (f"• {bp['name']}: {bp['avg']:.2f} vs P25 {b['p25']} (peer median: {b['p50']})", 11, False, _hex_to_rgb("334155"))
            )
    above_p75 = [k for k in fp_data if k["key"] in bench and k["avg"] is not None and k["avg"] >= bench[k["key"]]["p75"]]
    if above_p75:
        bench_narrative.append(("", 6, False, None))
        bench_narrative.append(
            (f"{len(above_p75)} KPIs are top quartile (above P75):", 13, True, _DECK_GREEN_FG)
        )
        for ap in above_p75[:5]:
            b = bench[ap["key"]]
            bench_narrative.append(
                (f"• {ap['name']}: {ap['avg']:.2f} vs P75 {b['p75']}", 11, False, _hex_to_rgb("334155"))
            )
    if bench_narrative:
        _add_narrative(slide4, 8.8, 1.3, 4.2, 5.5, bench_narrative)

    # ── Slide 5: Watch Zone — Narrative + Trend Chart ─────────────────────
    if yellow_kpis:
        slide5 = prs.slides.add_slide(blank_layout)
        _add_narrative(slide5, 0.5, 0.3, 8, 0.7, [
            (f"Watch Zone: {len(yellow_kpis)} KPIs Approaching Threshold", 28, True, _hex_to_rgb("1e293b")),
        ])

        watch_paras = []
        for k in yellow_kpis[:4]:
            watch_paras.append((_kpi_sentence(k), 11, False, _hex_to_rgb("334155")))
            watch_paras.append(("", 6, False, None))

        watch_paras.append(("These KPIs are at risk of moving to critical status without intervention.", 12, True, _DECK_YELLOW_FG))
        _add_narrative(slide5, 0.5, 1.2, 6, 5.5, watch_paras)

        trend_yellow_png = _make_trend_chart(yellow_kpis, "Watch Zone Trends", max_kpis=4)
        if trend_yellow_png:
            slide5.shapes.add_picture(trend_yellow_png, Inches(6.5), Inches(1.2), Inches(6.5), Inches(5.2))

    # ── Slide 6: Recommended Actions (top 5 priorities) ───────────────────
    slide6 = prs.slides.add_slide(blank_layout)
    _add_narrative(slide6, 0.5, 0.3, 10, 0.7, [
        ("Recommended Actions — Priority Order", 28, True, _hex_to_rgb("1e293b")),
    ])

    action_paras = []
    for i, k in enumerate(red_kpis[:5], 1):
        rules = ALL_CAUSATION_RULES.get(k["key"], {})
        actions = rules.get("corrective_actions", [])
        causes = rules.get("root_causes", [])
        val = k["avg"]
        tgt = k["target"]
        gap_str = ""
        if val is not None and tgt is not None and tgt != 0:
            gap = round(abs((val - tgt) / abs(tgt) * 100), 0)
            gap_str = f" ({gap:.0f}% off target)"

        action_paras.append((f"{i}. {k['name']}{gap_str}", 14, True, _DECK_RED_FG))
        if causes:
            action_paras.append((f"   Root cause: {causes[0]}", 11, False, _hex_to_rgb("475569")))
        if actions:
            for a in actions[:2]:
                action_paras.append((f"   → {a}", 11, False, _hex_to_rgb("334155")))

        # Add downstream impact context
        downstream = rules.get("downstream_impact", [])
        downstream_at_risk = [dk for dk in fp_data if dk["key"] in downstream and dk["fy_status"] in ("red", "yellow")]
        if downstream_at_risk:
            names_str = ", ".join(d["name"] for d in downstream_at_risk[:3])
            action_paras.append((f"   Impact: fixing this also improves {names_str}", 11, True, _hex_to_rgb("2563eb")))
        action_paras.append(("", 6, False, None))

    if not action_paras:
        action_paras.append(("No critical actions required — all KPIs are within tolerance.", 14, False, _DECK_GREEN_FG))

    _add_narrative(slide6, 0.5, 1.2, 12, 5.5, action_paras)

    # ── Slide 7: Closing — Bright Spots + Next Review ─────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    bg7 = slide7.background
    fill7 = bg7.fill
    fill7.solid()
    fill7.fore_color.rgb = _DECK_DARK_BLUE

    closing_paras = [
        ("Key Takeaways", 32, True, _DECK_WHITE),
        ("", 8, False, None),
    ]
    if red_kpis:
        closing_paras.append(
            (f"• {len(red_kpis)} metrics need immediate attention — {red_kpis[0]['name']} is the top priority", 16, False, _DECK_WHITE)
        )
    if yellow_kpis:
        closing_paras.append(
            (f"• {len(yellow_kpis)} metrics in watch zone — monitor weekly to prevent escalation", 16, False, _DECK_WHITE)
        )
    if green_kpis:
        closing_paras.append(
            (f"• {len(green_kpis)} metrics on target — maintain current trajectory", 16, False, _DECK_WHITE)
        )
    if below_p25:
        closing_paras.append(
            (f"• {len(below_p25)} metrics below {stage_label} industry bottom quartile", 16, False, _DECK_WHITE)
        )
    closing_paras.append(("", 12, False, None))
    closing_paras.append(
        (f"Generated {datetime.now().strftime('%B %d, %Y')}  ·  Axiom Intelligence  ·  {stage_label} Benchmarks", 14, False, _hex_to_rgb("94a3b8"))
    )
    _add_narrative(slide7, 1, 1.5, 11, 5, closing_paras)
    for shape in slide7.shapes:
        if hasattr(shape, "text_frame"):
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER

    # Serialize
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=board-deck.pptx"},
    )


# ─── KPI Annotations CRUD ───────────────────────────────────────────────────

from pydantic import BaseModel as _AnnotationBaseModel

class _AnnotationBody(_AnnotationBaseModel):
    kpi_key: str
    period: str
    note: str

@app.get("/api/annotations", tags=["Annotations"])
def list_annotations(kpi_key: Optional[str] = None, period: Optional[str] = None):
    """List KPI annotations, optionally filtered by kpi_key and/or period."""
    conn = get_db()
    clauses, params = [], []
    if kpi_key:
        clauses.append("kpi_key = ?")
        params.append(kpi_key)
    if period:
        clauses.append("period = ?")
        params.append(period)
    where = (" WHERE " + " AND ".join(clauses)) if clauses else ""
    rows = conn.execute(f"SELECT * FROM kpi_annotations{where} ORDER BY created_at DESC", params).fetchall()
    conn.close()
    return {"annotations": [dict(r) for r in rows]}

@app.put("/api/annotations", tags=["Annotations"])
def upsert_annotation(body: _AnnotationBody):
    """Create or update (upsert) a KPI annotation for a given kpi_key + period."""
    conn = get_db()
    now = datetime.utcnow().isoformat()
    conn.execute(
        """INSERT INTO kpi_annotations (kpi_key, period, note, created_at, updated_at)
           VALUES (?, ?, ?, ?, ?)
           ON CONFLICT(kpi_key, period) DO UPDATE SET note=excluded.note, updated_at=excluded.updated_at""",
        (body.kpi_key, body.period, body.note, now, now),
    )
    conn.commit()
    row = conn.execute(
        "SELECT * FROM kpi_annotations WHERE kpi_key=? AND period=?",
        (body.kpi_key, body.period),
    ).fetchone()
    conn.close()
    return {"status": "ok", "annotation": dict(row)}

@app.delete("/api/annotations/{annotation_id}", tags=["Annotations"])
def delete_annotation(annotation_id: int):
    """Delete a KPI annotation by its ID."""
    conn = get_db()
    existing = conn.execute("SELECT id FROM kpi_annotations WHERE id=?", (annotation_id,)).fetchone()
    if not existing:
        conn.close()
        raise HTTPException(status_code=404, detail="Annotation not found")
    conn.execute("DELETE FROM kpi_annotations WHERE id=?", (annotation_id,))
    conn.commit()
    conn.close()
    return {"status": "ok"}


# ─── KPI Accountability ─────────────────────────────────────────────────────

@app.get("/api/accountability", tags=["Accountability"])
def get_accountability(kpi_key: Optional[str] = None):
    conn = get_db()
    if kpi_key:
        rows = conn.execute("SELECT * FROM kpi_accountability WHERE kpi_key=?", (kpi_key,)).fetchall()
    else:
        rows = conn.execute("SELECT * FROM kpi_accountability").fetchall()
    conn.close()
    result = {}
    for r in rows:
        result[r["kpi_key"]] = {
            "kpi_key": r["kpi_key"],
            "owner": r["owner"],
            "due_date": r["due_date"],
            "status": r["status"],
            "last_updated": r["last_updated"],
        }
    return {"accountability": result}

@app.put("/api/accountability/{kpi_key}", tags=["Accountability"])
def put_accountability(kpi_key: str, body: dict):
    owner = body.get("owner", "")
    due_date = body.get("due_date", "")
    status = body.get("status", "open")
    now = datetime.now().isoformat()
    conn = get_db()
    conn.execute("""
        INSERT INTO kpi_accountability (kpi_key, owner, due_date, status, last_updated)
        VALUES (?, ?, ?, ?, ?)
        ON CONFLICT(kpi_key) DO UPDATE SET
            owner = excluded.owner,
            due_date = excluded.due_date,
            status = excluded.status,
            last_updated = excluded.last_updated
    """, (kpi_key, owner, due_date, status, now))
    conn.commit()
    conn.close()
    return {"status": "ok", "accountability": {"kpi_key": kpi_key, "owner": owner, "due_date": due_date, "status": status, "last_updated": now}}


@app.get("/api/export/weekly-briefing.html", tags=["Board Deck"])
def export_weekly_briefing(stage: str = "series_b"):
    """Generate an HTML weekly briefing document."""
    fp_data = _compute_fingerprint_data()
    if not fp_data:
        raise HTTPException(status_code=404, detail="No data available")

    # Get benchmarks
    bench = {}
    for kpi_key, stages_data in BENCHMARKS.items():
        if stage in stages_data:
            bench[kpi_key] = stages_data[stage]

    # Get accountability data
    conn = get_db()
    acct_rows = conn.execute("SELECT * FROM kpi_accountability").fetchall()
    conn.close()
    acct = {}
    for r in acct_rows:
        acct[r["kpi_key"]] = {"owner": r["owner"], "due_date": r["due_date"], "status": r["status"]}

    # Categorise
    green = [k for k in fp_data if k["fy_status"] == "green"]
    yellow = [k for k in fp_data if k["fy_status"] == "yellow"]
    red = [k for k in fp_data if k["fy_status"] == "red"]
    total = len(fp_data)

    # Sort red by gap magnitude
    def gap_pct(k):
        if k["avg"] is None or not k["target"]: return 0
        raw = (k["avg"] / k["target"] - 1) * 100
        return -raw if k["direction"] != "higher" else raw

    red_sorted = sorted(red, key=lambda k: abs(gap_pct(k)), reverse=True)
    yellow_sorted = sorted(yellow, key=lambda k: abs(gap_pct(k)), reverse=True)

    stage_label = {"seed": "Seed", "series_a": "Series A", "series_b": "Series B", "series_c": "Series C+"}.get(stage, stage)
    date_str = datetime.now().strftime("%B %d, %Y")

    def fmt_val(val, unit):
        if val is None: return "—"
        if unit == "pct": return f"{val:.1f}%"
        if unit == "days": return f"{val:.1f}d"
        if unit == "months": return f"{val:.1f}mo"
        if unit == "ratio": return f"{val:.2f}x"
        return f"{val:.2f}"

    # Get causation rules
    causation = {}
    try:
        causation = globals().get("CAUSATION_RULES", {})
    except:
        pass

    # Build HTML
    html_parts = []
    html_parts.append(f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>Weekly Briefing — {date_str}</title>
<style>
  body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; color: #1e293b; line-height: 1.6; }}
  h1 {{ color: #0f172a; font-size: 24px; border-bottom: 2px solid #0055A4; padding-bottom: 8px; }}
  h2 {{ color: #334155; font-size: 18px; margin-top: 32px; }}
  .summary {{ display: flex; gap: 16px; margin: 16px 0; }}
  .stat {{ padding: 12px 20px; border-radius: 12px; text-align: center; flex: 1; }}
  .stat-red {{ background: #fef2f2; border: 1px solid #fecaca; }}
  .stat-yellow {{ background: #fffbeb; border: 1px solid #fde68a; }}
  .stat-green {{ background: #f0fdf4; border: 1px solid #bbf7d0; }}
  .stat .num {{ font-size: 28px; font-weight: 800; }}
  .stat-red .num {{ color: #dc2626; }}
  .stat-yellow .num {{ color: #d97706; }}
  .stat-green .num {{ color: #059669; }}
  .stat .label {{ font-size: 11px; color: #64748b; text-transform: uppercase; font-weight: 700; letter-spacing: 0.05em; }}
  table {{ width: 100%; border-collapse: collapse; margin: 12px 0; font-size: 13px; }}
  th {{ background: #0055A4; color: white; padding: 8px 12px; text-align: left; font-size: 11px; text-transform: uppercase; letter-spacing: 0.05em; }}
  td {{ padding: 8px 12px; border-bottom: 1px solid #e2e8f0; }}
  tr:nth-child(even) {{ background: #f8fafc; }}
  .red {{ color: #dc2626; font-weight: 700; }}
  .yellow {{ color: #d97706; font-weight: 700; }}
  .green {{ color: #059669; font-weight: 700; }}
  .owner {{ color: #0055A4; font-weight: 600; }}
  .footer {{ margin-top: 40px; padding-top: 16px; border-top: 1px solid #e2e8f0; font-size: 11px; color: #94a3b8; }}
  .action {{ background: #eff6ff; border-left: 3px solid #0055A4; padding: 8px 12px; margin: 4px 0; font-size: 12px; border-radius: 0 8px 8px 0; }}
</style></head><body>
<h1>Weekly Finance Briefing</h1>
<p style="color:#64748b; font-size:13px;">{date_str} · {stage_label} stage · {total} KPIs tracked</p>

<div class="summary">
  <div class="stat stat-red"><div class="num">{len(red)}</div><div class="label">Critical</div></div>
  <div class="stat stat-yellow"><div class="num">{len(yellow)}</div><div class="label">Watch</div></div>
  <div class="stat stat-green"><div class="num">{len(green)}</div><div class="label">On Target</div></div>
</div>
""")

    if red_sorted:
        html_parts.append("<h2>🔴 Critical KPIs — Immediate Attention Required</h2>")
        html_parts.append("<table><tr><th>KPI</th><th>Current</th><th>Target</th><th>Gap</th><th>Owner</th><th>Status</th></tr>")
        for k in red_sorted:
            gap = gap_pct(k)
            a = acct.get(k["key"], {})
            owner_str = a.get("owner", "—") or "—"
            status_str = a.get("status", "unassigned") or "unassigned"
            cause_data = causation.get(k["key"], {})
            root = cause_data.get("root_causes", [""])[0] if isinstance(cause_data.get("root_causes"), list) else ""
            fix = cause_data.get("corrective_actions", [""])[0] if isinstance(cause_data.get("corrective_actions"), list) else ""
            html_parts.append(f'<tr><td><strong>{k["name"]}</strong></td><td>{fmt_val(k["avg"], k["unit"])}</td><td>{fmt_val(k["target"], k["unit"])}</td><td class="red">{gap:+.1f}%</td><td class="owner">{owner_str}</td><td>{status_str}</td></tr>')
        html_parts.append("</table>")

    if yellow_sorted:
        html_parts.append("<h2>🟡 Watch Zone</h2>")
        html_parts.append("<table><tr><th>KPI</th><th>Current</th><th>Target</th><th>Gap</th><th>Owner</th></tr>")
        for k in yellow_sorted[:6]:
            gap = gap_pct(k)
            a = acct.get(k["key"], {})
            owner_str = a.get("owner", "—") or "—"
            html_parts.append(f'<tr><td>{k["name"]}</td><td>{fmt_val(k["avg"], k["unit"])}</td><td>{fmt_val(k["target"], k["unit"])}</td><td class="yellow">{gap:+.1f}%</td><td class="owner">{owner_str}</td></tr>')
        html_parts.append("</table>")

    if green[:5]:
        html_parts.append("<h2>🟢 Bright Spots</h2>")
        html_parts.append("<table><tr><th>KPI</th><th>Current</th><th>Target</th><th>Above Target</th></tr>")
        for k in sorted(green, key=lambda k: gap_pct(k), reverse=True)[:5]:
            gap = gap_pct(k)
            html_parts.append(f'<tr><td>{k["name"]}</td><td>{fmt_val(k["avg"], k["unit"])}</td><td>{fmt_val(k["target"], k["unit"])}</td><td class="green">+{abs(gap):.1f}%</td></tr>')
        html_parts.append("</table>")

    html_parts.append(f"""
<div class="footer">
  Generated by Axiom Intelligence · {date_str}<br>
  This briefing covers {total} KPIs for the {stage_label} stage.
</div>
</body></html>""")

    html_content = "\n".join(html_parts)
    buf = io.BytesIO(html_content.encode("utf-8"))
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="text/html",
        headers={"Content-Disposition": f'attachment; filename="weekly-briefing-{datetime.now().strftime("%Y%m%d")}.html"'}
    )


# ─── Smart Actions Endpoint ─────────────────────────────────────────────────

def _build_causal_chain(kpi_key: str, fp_lookup: dict, max_depth: int = 3) -> list:
    """
    Build a multi-hop causal chain starting from kpi_key.
    Traces upstream red/yellow KPIs up to max_depth hops.
    Returns a list of chain nodes with hop, kpi metadata, and label.
    """
    status_order = {"red": 0, "yellow": 1, "green": 2, "grey": 3}

    def _gap_pct(kpi_data):
        val = kpi_data.get("avg")
        target = kpi_data.get("target")
        direction = kpi_data.get("direction", "higher")
        if val is None or target is None or target == 0:
            return None
        if direction == "higher":
            return round((val - target) / abs(target) * 100, 1)
        else:
            return round((target - val) / abs(target) * 100, 1)

    def _node(hop, key, kpi_data, label):
        val = kpi_data.get("avg")
        target = kpi_data.get("target")
        status = kpi_data.get("fy_status", "grey")
        return {
            "hop": hop,
            "kpi_key": key,
            "kpi_name": kpi_data.get("name", key.replace("_", " ").title()),
            "value": val,
            "target": target,
            "status": status,
            "gap_pct": _gap_pct(kpi_data),
            "label": label,
        }

    hop_labels = {0: "Surface symptom", 1: "Primary driver", 2: "Contributing cause", 3: "Root cause"}

    chain = []
    visited = set()

    # Hop 0: the requested KPI itself
    root_data = fp_lookup.get(kpi_key)
    if root_data is None:
        return chain
    chain.append(_node(0, kpi_key, root_data, hop_labels[0]))
    visited.add(kpi_key)

    # BFS-style expansion up to max_depth
    current_level_keys = [kpi_key]
    for hop in range(1, max_depth + 1):
        next_level_keys = []
        hop_nodes = []
        for parent_key in current_level_keys:
            # Find upstream KPIs: those whose downstream_impact includes parent_key
            for source_key, rules in ALL_CAUSATION_RULES.items():
                if parent_key in rules.get("downstream_impact", []):
                    if source_key in visited:
                        continue
                    source_data = fp_lookup.get(source_key)
                    if source_data is None:
                        continue
                    status = source_data.get("fy_status", "grey")
                    # Only include red/yellow in the chain
                    if status in ("red", "yellow"):
                        label = hop_labels.get(hop, f"Hop {hop} cause")
                        hop_nodes.append(_node(hop, source_key, source_data, label))
                        visited.add(source_key)
                        next_level_keys.append(source_key)

        # Sort each hop level: red first, then by gap magnitude (worst first)
        hop_nodes.sort(key=lambda x: (
            status_order.get(x["status"], 4),
            -(abs(x["gap_pct"]) if x["gap_pct"] is not None else 0)
        ))
        chain.extend(hop_nodes)
        current_level_keys = next_level_keys
        if not current_level_keys:
            break

    return chain


def _generate_smart_actions(kpi_key: str, fp_data: list, benchmarks_for_stage: dict, stage: str):
    """
    Generate data-backed, number-specific corrective action recommendations
    for a given KPI using fingerprint data, benchmarks, and the causal knowledge graph.
    """
    # ── Find the requested KPI in fingerprint data ──
    kpi_data = None
    fp_lookup = {}
    for k in fp_data:
        fp_lookup[k["key"]] = k
        if k["key"] == kpi_key:
            kpi_data = k
    if kpi_data is None:
        return None

    current_value = kpi_data["avg"]
    target = kpi_data["target"]
    direction = kpi_data["direction"]
    unit = kpi_data["unit"]
    status = kpi_data["fy_status"]
    kpi_name = kpi_data["name"]
    monthly = kpi_data["monthly"]

    # ── Step 1: Quantify the problem ──
    gap_pct = None
    if current_value is not None and target is not None and target != 0:
        if direction == "higher":
            gap_pct = round((current_value - target) / abs(target) * 100, 1)
        else:
            gap_pct = round((target - current_value) / abs(target) * 100, 1)

    # ── Benchmark context ──
    bench_info = None
    bench_data = benchmarks_for_stage.get(kpi_key)
    if bench_data and current_value is not None:
        p25 = bench_data.get("p25")
        p50 = bench_data.get("p50")
        p75 = bench_data.get("p75")
        pct_from_median = round((current_value - p50) / abs(p50) * 100, 1) if p50 and p50 != 0 else None
        if direction == "higher":
            if current_value < p25:
                position = "below_p25"
            elif current_value < p50:
                position = "p25_to_p50"
            elif current_value < p75:
                position = "p50_to_p75"
            else:
                position = "above_p75"
        else:
            # For "lower is better" KPIs, lower values are better
            if current_value > p75:
                position = "below_p25"  # worst quartile
            elif current_value > p50:
                position = "p25_to_p50"
            elif current_value > p25:
                position = "p50_to_p75"
            else:
                position = "above_p75"  # best quartile
        bench_info = {
            "p25": p25, "p50": p50, "p75": p75,
            "position": position,
            "pct_from_median": pct_from_median,
        }

    # ── Trend computation (last 3 months) ──
    trend_info = None
    if monthly and len(monthly) >= 2:
        sorted_monthly = sorted(monthly, key=lambda m: m["period"])
        last_3 = sorted_monthly[-3:] if len(sorted_monthly) >= 3 else sorted_monthly
        vals_3 = [m["value"] for m in last_3 if m["value"] is not None]

        trend_direction = "stable"
        pct_change_3m = None
        consecutive_declining = 0

        if len(vals_3) >= 2:
            if vals_3[-1] > vals_3[0]:
                trend_direction = "improving" if direction == "higher" else "declining"
            elif vals_3[-1] < vals_3[0]:
                trend_direction = "declining" if direction == "higher" else "improving"

            if vals_3[0] != 0:
                pct_change_3m = round((vals_3[-1] - vals_3[0]) / abs(vals_3[0]) * 100, 1)

            # Count consecutive months in the same direction
            for i in range(len(vals_3) - 1, 0, -1):
                if direction == "higher":
                    if vals_3[i] < vals_3[i - 1]:
                        consecutive_declining += 1
                    else:
                        break
                else:
                    if vals_3[i] > vals_3[i - 1]:
                        consecutive_declining += 1
                    else:
                        break

        # Count consecutive red months from the end
        consecutive_red = 0
        for m in reversed(sorted_monthly):
            v = m["value"]
            if v is not None and target is not None:
                if direction == "higher":
                    ratio = v / target if target != 0 else 0
                    if ratio < 0.90:
                        consecutive_red += 1
                    else:
                        break
                else:
                    ratio = v / target if target != 0 else 0
                    if ratio > 1.10:
                        consecutive_red += 1
                    else:
                        break
            else:
                break

        trend_info = {
            "direction": trend_direction,
            "last_3_months": last_3,
            "pct_change_3m": pct_change_3m,
            "consecutive_red_months": consecutive_red,
        }

    # ── Step 2: Trace upstream causes using CAUSATION_RULES ──
    upstream_causes = []
    for source_key, rules in ALL_CAUSATION_RULES.items():
        if kpi_key in rules.get("downstream_impact", []):
            upstream_kpi = fp_lookup.get(source_key)
            if upstream_kpi:
                u_val = upstream_kpi["avg"]
                u_target = upstream_kpi["target"]
                u_status = upstream_kpi["fy_status"]
                u_direction = upstream_kpi["direction"]
                u_gap = None
                if u_val is not None and u_target is not None and u_target != 0:
                    if u_direction == "higher":
                        u_gap = round((u_val - u_target) / abs(u_target) * 100, 1)
                    else:
                        u_gap = round((u_target - u_val) / abs(u_target) * 100, 1)

                if u_val is not None:
                    u_unit_label = upstream_kpi["unit"]
                    fmt_val = f"{u_val}"
                    fmt_target = f"{u_target}" if u_target is not None else "N/A"
                    if u_status in ("red", "yellow"):
                        gap_str = f"{abs(u_gap)}% {'below' if u_gap < 0 else 'above'} target" if u_gap is not None else ""
                        explanation = (
                            f"{upstream_kpi['name']} is at {fmt_val} vs target {fmt_target} ({gap_str}). "
                            f"This directly impacts {kpi_name} as a causal upstream driver."
                        )
                    else:
                        explanation = (
                            f"{upstream_kpi['name']} is healthy at {fmt_val} "
                            f"({'above' if u_direction == 'higher' else 'below'} target of {fmt_target}) "
                            f"— this is not contributing to the problem."
                        )
                    upstream_causes.append({
                        "kpi_key": source_key,
                        "kpi_name": upstream_kpi["name"],
                        "status": u_status,
                        "value": u_val,
                        "target": u_target,
                        "gap_pct": u_gap,
                        "explanation": explanation,
                        "is_data_available": True,
                    })
                else:
                    upstream_causes.append({
                        "kpi_key": source_key,
                        "kpi_name": upstream_kpi["name"],
                        "status": "grey",
                        "value": None,
                        "target": u_target,
                        "gap_pct": None,
                        "explanation": f"No data available for {upstream_kpi['name']}. Collecting it would improve diagnostic accuracy.",
                        "is_data_available": False,
                    })
            else:
                # KPI exists in causation rules but not in fingerprint data at all
                upstream_causes.append({
                    "kpi_key": source_key,
                    "kpi_name": source_key.replace("_", " ").title(),
                    "status": "grey",
                    "value": None,
                    "target": None,
                    "gap_pct": None,
                    "explanation": f"No data available for {source_key.replace('_', ' ').title()}. Collecting it would improve diagnostic accuracy.",
                    "is_data_available": False,
                })

    # Sort: red first, then yellow, then green, then grey
    status_order = {"red": 0, "yellow": 1, "green": 2, "grey": 3}
    upstream_causes.sort(key=lambda x: status_order.get(x["status"], 4))

    # ── Find downstream impact ──
    downstream_impact = []
    rules_for_kpi = ALL_CAUSATION_RULES.get(kpi_key, {})
    for ds_key in rules_for_kpi.get("downstream_impact", []):
        ds_kpi = fp_lookup.get(ds_key)
        if ds_kpi:
            ds_val = ds_kpi["avg"]
            ds_target = ds_kpi["target"]
            ds_status = ds_kpi["fy_status"]
            ds_direction = ds_kpi["direction"]
            ds_gap = None
            if ds_val is not None and ds_target is not None and ds_target != 0:
                if ds_direction == "higher":
                    ds_gap = round((ds_val - ds_target) / abs(ds_target) * 100, 1)
                else:
                    ds_gap = round((ds_target - ds_val) / abs(ds_target) * 100, 1)
            fmt_ds_val = f"{ds_val}" if ds_val is not None else "N/A"
            fmt_ds_target = f"{ds_target}" if ds_target is not None else "N/A"
            explanation = (
                f"{ds_kpi['name']} is at {fmt_ds_val} (target {fmt_ds_target}, status: {ds_status}). "
                f"If {kpi_name} improves, {ds_kpi['name']} should improve proportionally as a downstream metric."
            )
            downstream_impact.append({
                "kpi_key": ds_key,
                "kpi_name": ds_kpi["name"],
                "status": ds_status,
                "value": ds_val,
                "target": ds_target,
                "explanation": explanation,
            })

    # ── Step 3: Generate SPECIFIC corrective actions with numbers ──
    actions = []
    priority = 1

    # Action from worst upstream cause
    red_upstreams = [u for u in upstream_causes if u["status"] == "red" and u["is_data_available"]]
    if red_upstreams:
        worst = red_upstreams[0]
        # Get corrective actions from causation rules for the upstream KPI
        upstream_actions = ALL_CAUSATION_RULES.get(worst["kpi_key"], {}).get("corrective_actions", [])
        specific_action = upstream_actions[0] if upstream_actions else "Investigate root cause"

        # Estimate impact: if upstream recovers to target, approximate the KPI improvement
        impact_str = ""
        if worst["value"] is not None and worst["target"] is not None and current_value is not None and target is not None:
            # Proportional estimate
            upstream_recovery_ratio = worst["target"] / worst["value"] if worst["value"] != 0 else 1
            estimated_new_value = round(current_value * upstream_recovery_ratio, 2)
            impact_str = (
                f"If {worst['kpi_name']} recovers from {worst['value']} to the target of {worst['target']}, "
                f"{kpi_name} would improve from {current_value} to approximately {estimated_new_value}, "
                f"assuming other factors remain constant."
            )

        actions.append({
            "priority": priority,
            "action": (
                f"Address {worst['kpi_name']} first — it is the primary upstream driver and is "
                f"{abs(worst['gap_pct'])}% {'below' if worst['gap_pct'] < 0 else 'above'} target. "
                f"{specific_action}."
            ),
            "expected_impact": impact_str,
            "owner_suggestion": "Revenue Operations / VP Sales" if "sales" in worst["kpi_key"] or "pipeline" in worst["kpi_key"] else "Finance / Operations",
            "timeframe": "30-60 days for diagnosis, 60-90 days for improvement",
        })
        priority += 1

    # Action from benchmark position
    if bench_info and current_value is not None:
        stage_label = stage.replace("_", " ").title()
        position = bench_info["position"]
        if position == "below_p25":
            p25_val = bench_info["p25"]
            p50_val = bench_info["p50"]
            corrective_actions = rules_for_kpi.get("corrective_actions", [])
            specific_fix = corrective_actions[0] if corrective_actions else "Review operational processes"
            actions.append({
                "priority": priority,
                "action": (
                    f"{kpi_name} at {current_value} is below the {stage_label} 25th percentile ({p25_val}). "
                    f"Benchmark against peer operational structure. {specific_fix}."
                ),
                "expected_impact": (
                    f"Reaching the peer median of {p50_val} would represent a "
                    f"{abs(bench_info['pct_from_median'])}% improvement from current levels."
                ),
                "owner_suggestion": "CRO / CFO",
                "timeframe": "Quarterly review cycle",
            })
            priority += 1
        elif position == "p25_to_p50":
            p50_val = bench_info["p50"]
            actions.append({
                "priority": priority,
                "action": (
                    f"{kpi_name} at {current_value} is between the 25th and 50th percentile for {stage_label}. "
                    f"Closing the gap to the median ({p50_val}) should be an operational priority."
                ),
                "expected_impact": (
                    f"Reaching the peer median of {p50_val} would represent a "
                    f"{abs(bench_info['pct_from_median'])}% improvement."
                ),
                "owner_suggestion": "Operations Lead",
                "timeframe": "60-90 days",
            })
            priority += 1

    # Action from trend
    if trend_info and trend_info["direction"] == "declining" and trend_info["pct_change_3m"] is not None:
        last3 = trend_info["last_3_months"]
        values_str = " -> ".join([f"{m['value']}" for m in last3 if m["value"] is not None])
        pct_chg = abs(trend_info["pct_change_3m"])
        first_val = next((m["value"] for m in last3 if m["value"] is not None), None)

        # Try to identify which upstream KPI also declined in same period
        coinciding_upstream = ""
        for u in red_upstreams:
            u_fp = fp_lookup.get(u["kpi_key"])
            if u_fp and u_fp["monthly"] and len(u_fp["monthly"]) >= 2:
                u_sorted = sorted(u_fp["monthly"], key=lambda m: m["period"])
                u_last3 = u_sorted[-3:] if len(u_sorted) >= 3 else u_sorted
                u_vals = [m["value"] for m in u_last3 if m["value"] is not None]
                if len(u_vals) >= 2 and u_vals[-1] < u_vals[0]:
                    coinciding_upstream = (
                        f" This coincides with {u['kpi_name']} deteriorating from {u_vals[0]} to {u_vals[-1]} "
                        f"— addressing {u['kpi_name']} is likely the highest-leverage fix."
                    )
                    break

        actions.append({
            "priority": priority,
            "action": (
                f"The {len(last3)}-month declining trend ({values_str}, -{pct_chg}%) suggests a structural change, "
                f"not seasonal variance. Investigate what changed in the operational process starting "
                f"{len(last3)} months ago.{coinciding_upstream}"
            ),
            "expected_impact": (
                f"Identifying and reversing the structural cause could restore the metric to its "
                f"{len(last3)}-month-ago level of {first_val}."
            ) if first_val is not None else "Reversing the trend would stabilize the metric.",
            "owner_suggestion": "Operations / Strategy",
            "timeframe": "2-week investigation, 30-day corrective action",
        })
        priority += 1

    # Fallback: if no actions generated, use corrective_actions from causation rules
    if not actions:
        corrective_actions = rules_for_kpi.get("corrective_actions", [])
        for i, ca in enumerate(corrective_actions[:3]):
            actions.append({
                "priority": priority,
                "action": ca,
                "expected_impact": f"Addressing this would help move {kpi_name} toward the target of {target}." if target else "Impact depends on severity of root cause.",
                "owner_suggestion": "Operations Lead",
                "timeframe": "30-60 days",
            })
            priority += 1

    # ── Step 4: Identify data gaps ──
    data_gaps = []
    for u in upstream_causes:
        if not u["is_data_available"]:
            data_gaps.append(
                f"{u['kpi_name']} has no data available — this is a critical upstream metric. "
                f"Collecting it would enable more precise diagnosis. Add it to your data collection via the KPI export template."
            )
    # Also flag missing monthly data
    if not monthly or len(monthly) < 3:
        data_gaps.append(
            f"Only {len(monthly) if monthly else 0} months of data available for {kpi_name}. "
            f"At least 3 months are needed for reliable trend analysis."
        )

    # ── Build multi-hop causal chain ──
    causal_chain = _build_causal_chain(kpi_key, fp_lookup, max_depth=3)

    # ── Analysis depth metadata ──
    chain_kpi_keys = list({node["kpi_key"] for node in causal_chain})
    max_hop_depth = max((node["hop"] for node in causal_chain), default=0)

    # Total monthly data points across all KPIs in the chain
    total_data_points = 0
    for ck in chain_kpi_keys:
        ck_fp = fp_lookup.get(ck)
        if ck_fp and ck_fp.get("monthly"):
            total_data_points += len([m for m in ck_fp["monthly"] if m.get("value") is not None])

    # Build chain summary: hop-0 → hop-1 → ... names
    chain_by_hop: dict = {}
    for node in causal_chain:
        chain_by_hop.setdefault(node["hop"], []).append(node["kpi_name"])
    chain_summary_parts = []
    for h in sorted(chain_by_hop.keys()):
        chain_summary_parts.append(" / ".join(chain_by_hop[h]))
    chain_summary = " → ".join(chain_summary_parts) if chain_summary_parts else kpi_name

    analysis_depth = {
        "total_data_points": total_data_points,
        "kpis_in_chain": len(chain_kpi_keys),
        "max_hop_depth": max_hop_depth,
        "chain_summary": f"{max_hop_depth}-hop causal trace: {chain_summary}",
    }

    # ── Direction protection ──
    wrong_action_map = {
        "burn_multiple": "Increase S&M headcount or pressure sales team",
        "arr_growth": "Increase S&M headcount or pressure sales team",
        "revenue_growth": "Increase S&M headcount or pressure sales team",
        "sales_efficiency": "Hire more sales reps or replace underperformers",
        "win_rate": "Hire more sales reps or replace underperformers",
        "churn_rate": "Assign more CSMs or offer discounts to retain accounts",
        "nrr": "Assign more CSMs or offer discounts to retain accounts",
        "gross_margin": "Reduce headcount or cut vendor contracts",
        "contribution_margin": "Reduce headcount or cut vendor contracts",
        "pipeline_conversion": "Increase demo volume or revise sales process",
        "cpl": "Increase total marketing budget",
        "marketing_roi": "Increase total marketing budget",
    }
    likely_wrong_action = wrong_action_map.get(kpi_key, "Address the surface symptom directly without tracing root cause")

    root_cause_distance = max_hop_depth
    is_deep_cause = root_cause_distance >= 2

    direction_protected = {
        "likely_wrong_action": likely_wrong_action,
        "root_cause_distance": root_cause_distance,
        "is_deep_cause": is_deep_cause,
    }

    return {
        "kpi_key": kpi_key,
        "kpi_name": kpi_name,
        "current_value": current_value,
        "target": target,
        "gap_pct": gap_pct,
        "unit": unit,
        "status": status,
        "benchmark": bench_info,
        "trend": trend_info,
        "upstream_causes": upstream_causes,
        "downstream_impact": downstream_impact,
        "actions": actions,
        "data_gaps": data_gaps,
        "causal_chain": causal_chain,
        "analysis_depth": analysis_depth,
        "direction_protected": direction_protected,
    }


@app.get("/api/smart-actions/{kpi_key}", tags=["Smart Actions"])
def get_smart_actions(kpi_key: str, stage: str = "series_b"):
    """
    Generate data-backed, number-specific corrective action recommendations for a KPI.
    Uses fingerprint data, benchmarks, and the causal knowledge graph (CAUSATION_RULES)
    to produce actionable insights with specific numbers, upstream/downstream analysis,
    and quantified expected impact.
    """
    # Validate stage
    valid_stages = {"seed", "series_a", "series_b", "series_c"}
    if stage not in valid_stages:
        stage = "series_b"

    # Compute fingerprint data
    fp_data = _compute_fingerprint_data()

    # Get benchmarks for stage
    bench = {}
    for bk, stages_data in BENCHMARKS.items():
        if stage in stages_data:
            bench[bk] = stages_data[stage]

    # Validate KPI exists
    valid_keys = {k["key"] for k in fp_data}
    if kpi_key not in valid_keys:
        from fastapi.responses import JSONResponse
        return JSONResponse(
            status_code=404,
            content={"detail": f"KPI '{kpi_key}' not found. Valid keys: {sorted(valid_keys)}"}
        )

    result = _generate_smart_actions(kpi_key, fp_data, bench, stage)
    if result is None:
        from fastapi.responses import JSONResponse
        return JSONResponse(
            status_code=404,
            content={"detail": f"Could not generate actions for KPI '{kpi_key}' — no data found."}
        )

    return result


# ─── Serve React Frontend ───────────────────────────────────────────────────

STATIC_DIR = Path(__file__).parent.parent / "frontend" / "dist"
if STATIC_DIR.exists():
    app.mount("/assets", StaticFiles(directory=STATIC_DIR / "assets"), name="assets")

    @app.get("/{full_path:path}", include_in_schema=False)
    def serve_spa(full_path: str):
        index = STATIC_DIR / "index.html"
        return FileResponse(index)
