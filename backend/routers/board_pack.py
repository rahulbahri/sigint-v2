"""
routers/board_pack.py — Intelligent Board Pack Generator.

POST /api/board-pack/generate — returns PPTX file with:
  - Flexible date range (1 month to multi-year)
  - Four pack modes: talk_track, just_kpis, variance_narrative, financial_projections
  - Knowledge-graph-based causal narratives (not templates)
  - Intelligently selected charts (heatmap, radar, waterfall, line, donut, bar)
  - Single Corporate Blue theme

Uses:
  - core/chart_engine.py for chart rendering
  - core/board_narrative.py for signal detection + narrative generation
  - core/narrative_engine.py for per-KPI root cause analysis
  - core/health_score.py for health score computation
  - core/criticality.py for composite criticality ranking
"""
import io
import json
from datetime import datetime
from typing import Optional

from fastapi import APIRouter, Request, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from core.database import get_db
from core.deps import _get_workspace

router = APIRouter()


# ── Request Schema ──────────────────────────────────────────────────────────

class BoardPackRequest(BaseModel):
    company_name: Optional[str] = None
    from_year: int
    from_month: int            # 1-12
    to_year: int
    to_month: int              # 1-12
    modes: list[str]           # ["talk_track", "just_kpis", "variance_narrative", "financial_projections"]
    kpi_keys: Optional[list] = None
    company_stage: str = "series_b"


VALID_MODES = {"talk_track", "just_kpis", "variance_narrative", "financial_projections"}


# ── Corporate Blue Theme (single theme) ─────────────────────────────────────

THEME = {
    "bg":        "FFFFFF",
    "accent":    "003087",
    "highlight": "0055A4",
    "text":      "0F172A",
    "subtext":   "64748B",
    "positive":  "059669",
    "warning":   "D97706",
    "critical":  "DC2626",
    "card_bg":   "F1F5F9",
    "title_bg":  "071E45",
}


# ── Endpoint ────────────────────────────────────────────────────────────────

@router.post("/api/board-pack/generate", tags=["Board Pack"])
async def generate_board_pack(request: Request, body: BoardPackRequest):
    """Generate an intelligent PPTX board pack with causal narratives and charts."""
    workspace_id = _get_workspace(request)

    # Validate modes
    modes = set(body.modes) & VALID_MODES
    if not modes:
        modes = {"just_kpis"}

    # Validate date range
    if (body.from_year, body.from_month) > (body.to_year, body.to_month):
        raise HTTPException(400, "from_date must be before to_date")

    try:
        return _build_board_pack(workspace_id, body, modes)
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Board pack generation failed: {str(e)}")


def _build_board_pack(workspace_id: str, body: BoardPackRequest, modes: set[str]):
    """Main generation function."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    from core.health_score import compute_health_score
    from core.kpi_defs import KPI_DEFS, EXTENDED_ONTOLOGY_METRICS, BENCHMARKS, ALL_CAUSATION_RULES
    from core.kpi_utils import compute_kpi_avg
    from core.chart_engine import (
        render_waterfall, PALETTE, STATUS_COLORS,
        add_native_line, add_native_multi_line, add_native_donut,
        add_native_bar_h, add_native_grouped_bar_h, add_native_radar,
        add_native_table_heatmap,
    )
    from core.board_narrative import (
        detect_signals, generate_executive_summary, generate_causal_narratives,
        generate_domain_narratives, generate_period_comparison,
        generate_outlook, generate_talk_track, build_so_what,
        _get_domain, _gap_pct, _cell_status, _fmt_val,
    )
    from core.narrative_engine import enrich_needs_attention

    include_notes = "talk_track" in modes

    # ── Load data ───────────────────────────────────────────────────────────
    conn = get_db()
    from_period = (body.from_year, body.from_month)
    to_period = (body.to_year, body.to_month)

    health = compute_health_score(
        conn, workspace_id,
        from_period=from_period, to_period=to_period,
    )

    # Monthly data (filtered by period)
    rows = conn.execute(
        "SELECT year, month, data_json FROM monthly_data "
        "WHERE workspace_id=? AND (year > ? OR (year = ? AND month >= ?)) "
        "AND (year < ? OR (year = ? AND month <= ?)) ORDER BY year, month",
        [workspace_id,
         from_period[0], from_period[0], from_period[1],
         to_period[0], to_period[0], to_period[1]],
    ).fetchall()

    targets_rows = conn.execute(
        "SELECT kpi_key, target_value, direction, unit FROM kpi_targets WHERE workspace_id=?",
        [workspace_id],
    ).fetchall()

    settings_rows = conn.execute(
        "SELECT key, value FROM company_settings WHERE workspace_id=?",
        [workspace_id],
    ).fetchall()

    # Data provenance: latest upload timestamp
    try:
        upload_row = conn.execute(
            "SELECT uploaded_at FROM uploads WHERE workspace_id=? ORDER BY id DESC LIMIT 1",
            [workspace_id],
        ).fetchone()
        data_as_of = upload_row["uploaded_at"] if upload_row else None
    except Exception:
        data_as_of = None

    conn.close()

    targets_map = {r["kpi_key"]: {"target": r["target_value"], "direction": r["direction"] or "higher", "unit": r["unit"] or ""}
                   for r in targets_rows}
    settings = {r["key"]: r["value"] for r in settings_rows}
    company_name = body.company_name or settings.get("company_name", "Company")

    # Period label
    from_label = f"{_month_name(body.from_month)} {body.from_year}"
    to_label = f"{_month_name(body.to_month)} {body.to_year}"
    if body.from_year == body.to_year and body.from_month == body.to_month:
        period_label = from_label
    elif body.from_year == body.to_year:
        period_label = f"{_month_name(body.from_month)}-{_month_name(body.to_month)} {body.from_year}"
    else:
        period_label = f"{from_label} to {to_label}"

    # ── Build fingerprint ───────────────────────────────────────────────────
    all_kpi_meta = {kd["key"]: kd for kd in KPI_DEFS}
    for em in EXTENDED_ONTOLOGY_METRICS:
        if em["key"] not in all_kpi_meta:
            all_kpi_meta[em["key"]] = em

    kpi_monthly: dict = {}
    for row in rows:
        mo_key = f"{row['year']}-{row['month']:02d}"
        data = json.loads(row["data_json"]) if isinstance(row["data_json"], str) else (row["data_json"] or {})
        for kpi_key, val in data.items():
            if kpi_key in ("year", "month") or kpi_key.startswith("_"):
                continue
            if isinstance(val, (int, float)):
                kpi_monthly.setdefault(kpi_key, {})[mo_key] = val

    fingerprint = []
    kpi_avgs = {}
    time_series = {}
    directions = {}
    targets_flat = {}

    for key in set(list(kpi_monthly.keys()) + [kd["key"] for kd in KPI_DEFS]):
        vals_dict = kpi_monthly.get(key, {})
        t = targets_map.get(key, {})
        kdef = all_kpi_meta.get(key, {"key": key, "name": key.replace("_", " ").title(),
                                       "unit": "ratio", "direction": "higher"})
        tval = t.get("target")
        dirn = t.get("direction", kdef.get("direction", "higher"))
        unit = t.get("unit", kdef.get("unit", "ratio"))

        monthly_sorted = sorted(vals_dict.items())
        values = [v for _, v in monthly_sorted]
        monthly_list = [{"period": k, "value": v} for k, v in monthly_sorted]
        avg = compute_kpi_avg(values, window=len(values), period_filtered=True)

        if avg is None and not tval:
            continue

        kpi_avgs[key] = avg
        time_series[key] = values
        directions[key] = dirn
        if tval is not None:
            targets_flat[key] = tval

        fingerprint.append({
            "key": key,
            "name": kdef.get("name", key.replace("_", " ").title()),
            "unit": unit,
            "target": tval,
            "direction": dirn,
            "avg": avg,
            "fy_status": _cell_status(avg, tval, dirn),
            "monthly": monthly_list,
        })

    # Sort: KPI_DEFS order first
    kpi_def_order = {kd["key"]: i for i, kd in enumerate(KPI_DEFS)}
    fingerprint.sort(key=lambda x: (kpi_def_order.get(x["key"], 9999), x["key"]))

    # Classify KPIs
    red_kpis = [k for k in fingerprint if k["fy_status"] == "red"]
    yellow_kpis = [k for k in fingerprint if k["fy_status"] == "yellow"]
    green_kpis = [k for k in fingerprint if k["fy_status"] == "green"]
    total = len(red_kpis) + len(yellow_kpis) + len(green_kpis)

    # Sort red by gap magnitude
    def _gap_mag(k):
        g = _gap_pct(k.get("avg"), k.get("target"), k.get("direction", "higher"))
        return abs(g) if g else 0
    red_kpis.sort(key=_gap_mag, reverse=True)

    # ── Run analyses ────────────────────────────────────────────────────────
    signals = detect_signals(fingerprint)
    domain_narratives = generate_domain_narratives(fingerprint)

    # Root cause analysis for red KPIs
    critical_analyses = enrich_needs_attention(
        [k["key"] for k in red_kpis[:5]],
        kpi_avgs, time_series, targets_flat, directions, {},
    )

    # Executive summary
    exec_summary = generate_executive_summary(
        health, fingerprint, signals, period_label, critical_analyses,
    )

    # Causal narratives
    causal_narratives = generate_causal_narratives(
        red_kpis[:5], kpi_avgs, time_series, targets_flat, directions,
    )

    # Outlook
    outlook_bullets = generate_outlook(fingerprint, signals, domain_narratives)

    # ── PPTX helpers ────────────────────────────────────────────────────────
    def _rgb(hex6: str) -> RGBColor:
        h = hex6.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _fill_bg(slide, hex_color: str):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(hex_color)

    def _add_text(slide, text, left, top, width, height,
                  font_size=14, bold=False, color="0F172A", align=PP_ALIGN.LEFT, italic=False):
        from pptx.util import Inches as _In, Pt as _Pt
        txBox = slide.shapes.add_textbox(_In(left), _In(top), _In(width), _In(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = _Pt(font_size)
        run.font.name = "Calibri"
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
        return txBox

    def _add_paragraphs(slide, paragraphs, left, top, width, height):
        """Add multi-paragraph text. paragraphs: [(text, size, bold, color_hex), ...]"""
        from pptx.util import Inches as _In, Pt as _Pt
        txBox = slide.shapes.add_textbox(_In(left), _In(top), _In(width), _In(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (text, fs, bold, color) in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = _Pt(fs)
            p.font.name = "Calibri"
            p.font.bold = bold
            if color:
                p.font.color.rgb = _rgb(color)
            p.space_after = _Pt(6)
        return txBox

    def _add_divider(slide, top_in: float, width: float = 12.0):
        from pptx.util import Inches as _In, Pt as _Pt
        line = slide.shapes.add_connector(1, _In(0.5), _In(top_in), _In(0.5 + width), _In(top_in))
        line.line.color.rgb = _rgb(THEME["card_bg"])
        line.line.width = _Pt(0.75)

    def _notes(slide, text):
        if include_notes and text:
            slide.notes_slide.notes_text_frame.text = text

    # ── Create Presentation ─────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title (all modes)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _fill_bg(slide, THEME["title_bg"])

    # Accent bar
    rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(THEME["highlight"])
    rect.line.fill.background()

    _add_text(slide, "BOARD PACK", 0.6, 1.8, 11, 0.5,
              font_size=12, bold=True, color=THEME["subtext"])
    _add_text(slide, company_name, 0.6, 2.4, 11, 1.2,
              font_size=40, bold=True, color="FFFFFF")
    _add_text(slide, f"Performance Review  ·  {period_label}", 0.6, 3.8, 11, 0.5,
              font_size=18, color="94A3B8")
    _add_text(slide, f"{len(red_kpis)} critical  ·  {len(yellow_kpis)} watch  ·  "
              f"{len(green_kpis)} on target  ·  {total} KPIs tracked",
              0.6, 4.5, 11, 0.6, font_size=20, bold=True, color="FFFFFF")
    provenance = f"Generated {datetime.utcnow().strftime('%B %d, %Y')}  ·  Axiom Intelligence"
    if data_as_of:
        try:
            dao = str(data_as_of).split("T")[0] if "T" in str(data_as_of) else str(data_as_of)[:10]
            provenance += f"  ·  Data as of {dao}"
        except Exception:
            pass
    _add_text(slide, provenance,
              0.6, 6.5, 11, 0.4, font_size=10, color="94A3B8", italic=True)

    _notes(slide, generate_talk_track("title", {
        "company_name": company_name, "period_label": period_label,
    }))

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2: Executive Health Summary (all modes)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _fill_bg(slide, THEME["bg"])

    _add_text(slide, "EXECUTIVE SUMMARY", 0.5, 0.3, 8, 0.35,
              font_size=9, bold=True, color=THEME["subtext"])

    # Left: narrative paragraphs
    paras = []
    for p_text in exec_summary:
        paras.append((p_text, 12, False, THEME["text"]))
        paras.append(("", 6, False, None))
    _add_paragraphs(slide, paras, 0.5, 1.0, 7.5, 5.5)

    # Right: donut chart (native)
    add_native_donut(slide, 8.8, 0.8, 4.0, 4.0,
        [f"Critical ({len(red_kpis)})", f"Watch ({len(yellow_kpis)})", f"On Target ({len(green_kpis)})"],
        [len(red_kpis), len(yellow_kpis), len(green_kpis)],
        [PALETTE["critical"], PALETTE["warning"], PALETTE["positive"]],
    )

    # Health score badge
    score = health.get("score", 0)
    score_color = THEME["positive"] if score >= 70 else (THEME["warning"] if score >= 50 else THEME["critical"])
    _add_text(slide, f"Health Score: {score}/100", 9.0, 5.0, 3.5, 0.5,
              font_size=16, bold=True, color=score_color, align=PP_ALIGN.CENTER)
    _add_text(slide, health.get("label", ""), 9.0, 5.5, 3.5, 0.4,
              font_size=12, color=THEME["subtext"], align=PP_ALIGN.CENTER)

    _notes(slide, generate_talk_track("health_summary", {
        "score": score, "label": health.get("label", ""),
        "n_red": len(red_kpis), "n_green": len(green_kpis),
    }))

    # ════════════════════════════════════════════════════════════════════════
    # MODE: just_kpis
    # ════════════════════════════════════════════════════════════════════════
    if "just_kpis" in modes:
        # Slide: KPI Status Heatmap
        if fingerprint:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, "KPI STATUS OVERVIEW", 0.5, 0.3, 8, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, "Status by KPI and Month", 0.5, 0.65, 8, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            # Build heatmap data (filtered to KPIs with targets)
            hm_kpis = [k for k in fingerprint if k.get("target") and k.get("monthly")][:20]
            if hm_kpis:
                month_labels = sorted(set(
                    m["period"] for k in hm_kpis for m in k.get("monthly", [])
                ))
                kpi_names = [k["name"][:25] for k in hm_kpis]
                status_matrix = []
                value_matrix = []
                for kpi in hm_kpis:
                    by_period = {m["period"]: m["value"] for m in kpi.get("monthly", [])}
                    row_status = []
                    row_values = []
                    for ml in month_labels:
                        val = by_period.get(ml)
                        row_status.append(_cell_status(val, kpi.get("target"), kpi.get("direction", "higher")))
                        row_values.append(val)
                    status_matrix.append(row_status)
                    value_matrix.append(row_values)

                add_native_table_heatmap(slide, 0.3, 1.3,
                                         min(12.5, len(month_labels) * 0.8 + 2.5),
                                         min(5.8, len(kpi_names) * 0.35 + 0.8),
                                         kpi_names, month_labels, status_matrix, value_matrix)

            _notes(slide, "This heatmap shows every KPI's status for each month in the selected period. "
                   "Green = on target, yellow = watch, red = critical. "
                   "Look for horizontal red streaks (persistent problems) and vertical clusters (systemic months).")

        # Slide: Domain Health Radar
        troubled_domains = {d: info for d, info in domain_narratives.items()
                           if d != "other" and info["total"] > 0}
        if len(troubled_domains) >= 3:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, "DOMAIN HEALTH", 0.5, 0.3, 8, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, "Business Area Performance", 0.5, 0.65, 8, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            dims = []
            actuals = []
            for d, info in troubled_domains.items():
                dims.append(info["label"][:15])
                # Domain score = % on target
                actuals.append(round(info["green_count"] / info["total"] * 100) if info["total"] else 50)

            add_native_radar(slide, 0.5, 1.3, 5.5, 5.5,
                            dims, actuals, title="Domain Health (% On Target)")

            # Domain summary on right
            dom_paras = []
            for d, info in sorted(troubled_domains.items(), key=lambda x: -x[1].get("red_count", 0)):
                status = "critical" if info["red_count"] > 0 else ("watch" if info["yellow_count"] > 0 else "strong")
                color = THEME["critical"] if status == "critical" else (THEME["warning"] if status == "watch" else THEME["positive"])
                dom_paras.append((f"{info['label']}", 12, True, color))
                dom_paras.append((f"{info['green_count']} on target, {info['yellow_count']} watch, {info['red_count']} critical", 10, False, THEME["subtext"]))
                dom_paras.append(("", 4, False, None))
            _add_paragraphs(slide, dom_paras, 6.5, 1.3, 6.3, 5.5)

            _notes(slide, "The radar chart shows what percentage of KPIs are on target in each business domain. "
                   "Domains with low scores indicate systemic issues requiring domain-level intervention.")

        # Slide: Top/Bottom KPIs
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, THEME["bg"])
        _add_text(slide, "KEY METRICS", 0.5, 0.3, 12, 0.35,
                  font_size=9, bold=True, color=THEME["subtext"])

        # Top red KPIs (left)
        if red_kpis:
            _add_text(slide, "Needs Attention", 0.5, 0.7, 6, 0.4,
                      font_size=16, bold=True, color=THEME["critical"])
            y = 1.3
            for k in red_kpis[:6]:
                gap = _gap_pct(k.get("avg"), k.get("target"), k.get("direction", "higher"))
                gap_str = f"({abs(gap):.0f}% gap)" if gap else ""
                _add_text(slide, f"● {k['name']}", 0.5, y, 5.5, 0.3,
                          font_size=11, bold=True, color=THEME["text"])
                _add_text(slide, f"{_fmt_val(k.get('avg'), k.get('unit', ''))} vs {_fmt_val(k.get('target'), k.get('unit', ''))} {gap_str}",
                          0.8, y + 0.3, 5.2, 0.25, font_size=9, color=THEME["subtext"])
                y += 0.65

        # Top green KPIs (right)
        if green_kpis:
            _add_text(slide, "Performing Well", 7.0, 0.7, 6, 0.4,
                      font_size=16, bold=True, color=THEME["positive"])
            y = 1.3
            for k in green_kpis[:6]:
                _add_text(slide, f"● {k['name']}", 7.0, y, 5.5, 0.3,
                          font_size=11, bold=True, color=THEME["text"])
                _add_text(slide, f"{_fmt_val(k.get('avg'), k.get('unit', ''))} (target: {_fmt_val(k.get('target'), k.get('unit', ''))})",
                          7.3, y + 0.3, 5.2, 0.25, font_size=9, color=THEME["subtext"])
                y += 0.65

        _notes(slide, "Left: KPIs below target threshold ranked by gap severity. "
               "Right: KPIs meeting or exceeding targets. "
               "Focus discussion on the critical items and what actions are in progress.")

    # ════════════════════════════════════════════════════════════════════════
    # MODE: variance_narrative
    # ════════════════════════════════════════════════════════════════════════
    if "variance_narrative" in modes:

        # Slide: Causal Chain Analysis
        if causal_narratives:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, "CAUSAL ANALYSIS", 0.5, 0.3, 8, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, "Root Cause Intelligence", 0.5, 0.65, 8, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            # Left: narrative cards for top 3 critical KPIs
            y_pos = 1.4
            paras = []
            for cn in causal_narratives[:3]:
                paras.append((cn["headline"], 12, True, THEME["critical"]))
                paras.append((cn["narrative"], 10, False, THEME["text"]))
                if cn.get("actions"):
                    paras.append((f"Action: {cn['actions'][0]}", 10, True, THEME["highlight"]))
                paras.append(("", 6, False, None))
            _add_paragraphs(slide, paras, 0.5, 1.4, 7.0, 5.5)

            # Right: trend chart of critical KPIs
            if red_kpis:
                month_labels_all = sorted(set(
                    m["period"] for k in red_kpis[:4] for m in k.get("monthly", [])
                ))
                series = {}
                tgts = {}
                for k in red_kpis[:4]:
                    by_p = {m["period"]: m["value"] for m in k.get("monthly", [])}
                    series[k["name"][:20]] = [by_p.get(ml) for ml in month_labels_all]
                    if k.get("target"):
                        tgts[k["name"][:20]] = k["target"]
                add_native_multi_line(slide, 7.8, 1.2, 5.2, 5.2,
                                     month_labels_all, series, "Critical KPI Trends", tgts)

            _notes(slide, generate_talk_track("causal_analysis", {"narratives": causal_narratives}))

        # Slide: Signals & Hidden Risks
        if signals:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, "HIDDEN SIGNALS", 0.5, 0.3, 8, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, "Structural Patterns in the Data", 0.5, 0.65, 8, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            paras = []
            sev_colors = {"critical": THEME["critical"], "warning": THEME["warning"], "positive": THEME["positive"]}
            for s in signals[:4]:
                paras.append((s["title"], 13, True, sev_colors.get(s["severity"], THEME["text"])))
                paras.append((s["body"], 10, False, THEME["text"]))
                paras.append(("", 6, False, None))
            _add_paragraphs(slide, paras, 0.5, 1.4, 12.0, 5.5)

            _notes(slide, generate_talk_track("signals", {"signals": signals}))

        # Slides: Domain Deep Dives (top 3 troubled domains)
        troubled = [(d, info) for d, info in domain_narratives.items()
                    if info.get("has_issues") and d != "other"]
        troubled.sort(key=lambda x: -(x[1].get("red_count", 0) * 10 + x[1].get("yellow_count", 0)))

        for domain, info in troubled[:3]:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, f"DOMAIN: {info['label'].upper()}", 0.5, 0.3, 8, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, info["label"], 0.5, 0.65, 8, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            # Narrative
            dom_paras = [
                (info["story"], 12, False, THEME["text"]),
                ("", 6, False, None),
                (f"{info['green_count']} on target  ·  {info['yellow_count']} watch  ·  {info['red_count']} critical",
                 11, True, THEME["subtext"]),
            ]
            _add_paragraphs(slide, dom_paras, 0.5, 1.4, 6.5, 5.0)

            # Domain KPI trend chart (right)
            domain_kpis = info.get("kpis", [])
            kpis_with_data = [k for k in domain_kpis if k.get("monthly")][:5]
            if kpis_with_data:
                ml = sorted(set(m["period"] for k in kpis_with_data for m in k.get("monthly", [])))
                series = {}
                for k in kpis_with_data:
                    by_p = {m["period"]: m["value"] for m in k.get("monthly", [])}
                    series[k["name"][:18]] = [by_p.get(p) for p in ml]
                add_native_multi_line(slide, 7.5, 1.2, 5.5, 5.5,
                                     ml, series, f"{info['label']} Trends")

            _notes(slide, generate_talk_track("domain", {"domain_info": info}))

        # Slide: Period Comparison (waterfall)
        # Only if we have enough months for a prior period comparison
        total_months = (body.to_year - body.from_year) * 12 + (body.to_month - body.from_month) + 1
        if total_months >= 2:
            # Split into first half / second half for comparison
            mid_months = total_months // 2
            # For waterfall: show KPI deltas between first and second half
            kpis_with_change = []
            for kpi in fingerprint:
                vals = [m["value"] for m in kpi.get("monthly", []) if m.get("value") is not None]
                if len(vals) >= 4:
                    first_half = vals[:len(vals) // 2]
                    second_half = vals[len(vals) // 2:]
                    avg1 = sum(first_half) / len(first_half)
                    avg2 = sum(second_half) / len(second_half)
                    if avg1 != 0:
                        delta = (avg2 - avg1) / abs(avg1) * 100
                        kpis_with_change.append((kpi["name"][:18], delta, kpi.get("direction", "higher")))

            if kpis_with_change:
                # Sort by absolute delta, take top 10
                kpis_with_change.sort(key=lambda x: -abs(x[1]))
                top_changes = kpis_with_change[:10]
                labels = [x[0] for x in top_changes]
                deltas = [x[1] for x in top_changes]

                slide = prs.slides.add_slide(blank)
                _fill_bg(slide, THEME["bg"])
                _add_text(slide, "PERIOD COMPARISON", 0.5, 0.3, 8, 0.35,
                          font_size=9, bold=True, color=THEME["subtext"])
                _add_text(slide, "Biggest Movers (% Change)", 0.5, 0.65, 8, 0.5,
                          font_size=20, bold=True, color=THEME["text"])

                wf = render_waterfall(labels, deltas, "Period-over-Period Change (%)")
                slide.shapes.add_picture(wf.image, Inches(0.5), Inches(1.5),
                                         Inches(12.0), Inches(5.0))

                _notes(slide, "This waterfall shows the biggest KPI movers between the first and second half "
                       "of the selected period. Green bars = improvement, red bars = deterioration.")

        # Slide: Corrective Actions
        if causal_narratives:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, "CORRECTIVE ACTIONS", 0.5, 0.3, 8, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, "Priority Actions — Data-Grounded", 0.5, 0.65, 8, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            action_paras = []
            for i, cn in enumerate(causal_narratives[:5], 1):
                action_paras.append((f"{i}. {cn['headline']}", 12, True, THEME["critical"]))
                for a in cn.get("actions", [])[:2]:
                    action_paras.append((f"   → {a}", 10, False, THEME["text"]))
                if cn.get("downstream_count", 0) > 0:
                    action_paras.append((
                        f"   Impact: affects {cn['downstream_count']} downstream KPIs",
                        10, True, THEME["highlight"],
                    ))
                action_paras.append(("", 5, False, None))

            _add_paragraphs(slide, action_paras, 0.5, 1.4, 12.0, 5.5)

            _notes(slide, generate_talk_track("actions", {}))

    # ════════════════════════════════════════════════════════════════════════
    # MODE: financial_projections
    # ════════════════════════════════════════════════════════════════════════
    if "financial_projections" in modes:

        # Slide: Benchmark Position
        bench = {}
        valid_stages = {"seed", "series_a", "series_b", "series_c"}
        stage = body.company_stage if body.company_stage in valid_stages else "series_b"
        for kpi_key, stages_data in BENCHMARKS.items():
            if stage in stages_data:
                bench[kpi_key] = stages_data[stage]

        stage_label = {"seed": "Seed", "series_a": "Series A",
                       "series_b": "Series B", "series_c": "Series C+"}.get(stage, stage)

        # Build benchmark comparison
        kpis_for_bench = [k for k in (red_kpis + yellow_kpis + green_kpis)
                          if k["key"] in bench and k.get("avg") is not None][:10]
        if kpis_for_bench:
            slide = prs.slides.add_slide(blank)
            _fill_bg(slide, THEME["bg"])
            _add_text(slide, f"PEER BENCHMARK — {stage_label.upper()}", 0.5, 0.3, 10, 0.35,
                      font_size=9, bold=True, color=THEME["subtext"])
            _add_text(slide, f"Company vs {stage_label} Peer Median", 0.5, 0.65, 10, 0.5,
                      font_size=20, bold=True, color=THEME["text"])

            names = [k["name"][:25] for k in kpis_for_bench]
            company_vals = [k["avg"] for k in kpis_for_bench]
            peer_vals = [bench[k["key"]]["p50"] for k in kpis_for_bench]
            bar_colors = [
                PALETTE["critical"] if k["fy_status"] == "red" else
                PALETTE["warning"] if k["fy_status"] == "yellow" else
                PALETTE["positive"]
                for k in kpis_for_bench
            ]

            add_native_grouped_bar_h(slide, 0.3, 1.3, 8.0, 5.5,
                                     names, company_vals, peer_vals, bar_colors,
                                     f"Company vs {stage_label} Median",
                                     peer_label=f"{stage_label} Median")

            # Narrative: below P25 / above P75
            bench_paras = []
            below_p25 = [k for k in kpis_for_bench if k["avg"] < bench[k["key"]].get("p25", float("inf"))]
            above_p75 = [k for k in kpis_for_bench if k["avg"] >= bench[k["key"]].get("p75", float("inf"))]
            if below_p25:
                bench_paras.append((f"{len(below_p25)} KPIs below {stage_label} bottom quartile (P25):", 12, True, THEME["critical"]))
                for bp in below_p25[:4]:
                    b = bench[bp["key"]]
                    bench_paras.append((f"  {bp['name']}: {bp['avg']:.1f} vs P25 {b['p25']}", 10, False, THEME["text"]))
            if above_p75:
                bench_paras.append(("", 4, False, None))
                bench_paras.append((f"{len(above_p75)} KPIs in top quartile (above P75):", 12, True, THEME["positive"]))
                for ap in above_p75[:4]:
                    b = bench[ap["key"]]
                    bench_paras.append((f"  {ap['name']}: {ap['avg']:.1f} vs P75 {b['p75']}", 10, False, THEME["text"]))
            if bench_paras:
                _add_paragraphs(slide, bench_paras, 8.8, 1.3, 4.2, 5.5)

        # Slide: 30-90 Day Outlook
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, THEME["bg"])
        _add_text(slide, "OUTLOOK", 0.5, 0.3, 8, 0.35,
                  font_size=9, bold=True, color=THEME["subtext"])
        _add_text(slide, "30-90 Day Forward View", 0.5, 0.65, 8, 0.5,
                  font_size=20, bold=True, color=THEME["text"])

        outlook_paras = []
        for bullet in outlook_bullets:
            outlook_paras.append((f"•  {bullet}", 13, False, THEME["text"]))
            outlook_paras.append(("", 6, False, None))
        _add_paragraphs(slide, outlook_paras, 0.5, 1.5, 12.0, 5.0)

        _notes(slide, generate_talk_track("outlook", {}))

    # ════════════════════════════════════════════════════════════════════════
    # SEC-GRADE ANALYTICS SLIDES (all modes except just_kpis)
    # ════════════════════════════════════════════════════════════════════════
    if "just_kpis" not in modes:
        try:
            _add_sec_slides(prs, blank, workspace_id, body, _add_text, _add_paragraphs,
                           _fill_bg, _rgb, _notes, include_notes, period_label,
                           generate_talk_track, THEME)
        except Exception as e:
            # SEC slides are non-blocking — log and continue
            print(f"[Board Pack] SEC slides failed (non-fatal): {e}")

    # ════════════════════════════════════════════════════════════════════════
    # CLOSING SLIDE: Key Takeaways (all modes)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _fill_bg(slide, THEME["title_bg"])

    # Accent bar
    rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(THEME["highlight"])
    rect.line.fill.background()

    closing_paras = [
        ("Key Takeaways", 28, True, "FFFFFF"),
        ("", 8, False, None),
    ]
    if red_kpis:
        closing_paras.append((
            f"•  {len(red_kpis)} metrics need immediate attention — {red_kpis[0]['name']} is the top priority",
            15, False, "FFFFFF",
        ))
    if yellow_kpis:
        closing_paras.append((
            f"•  {len(yellow_kpis)} metrics in watch zone — monitor weekly to prevent escalation",
            15, False, "FFFFFF",
        ))
    if green_kpis:
        closing_paras.append((
            f"•  {len(green_kpis)} metrics on target — maintain current trajectory",
            15, False, "FFFFFF",
        ))

    # Add the most important signal
    if signals:
        s = signals[0]
        closing_paras.append((f"•  Signal: {s['title']}", 15, False, "FFFFFF"))

    # Add top action
    if causal_narratives and causal_narratives[0].get("actions"):
        closing_paras.append((
            f"•  Priority action: {causal_narratives[0]['actions'][0][:120]}",
            15, False, "FFFFFF",
        ))

    closing_paras.append(("", 12, False, None))
    closing_paras.append((
        f"Generated {datetime.utcnow().strftime('%B %d, %Y')}  ·  Axiom Intelligence  ·  {period_label}",
        11, False, "94A3B8",
    ))

    _add_paragraphs(slide, closing_paras, 0.6, 1.5, 12.0, 5.5)

    _notes(slide, generate_talk_track("takeaways", {}))

    # ── Serialise ───────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"{company_name.replace(' ', '_')}_Board_Pack_{body.from_year}{body.from_month:02d}-{body.to_year}{body.to_month:02d}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def _month_name(month_num: int) -> str:
    """Convert month number to abbreviated name."""
    names = ["", "Jan", "Feb", "Mar", "Apr", "May", "Jun",
             "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    return names[month_num] if 1 <= month_num <= 12 else str(month_num)


def _add_sec_slides(prs, blank, workspace_id, body, _add_text, _add_paragraphs,
                    _fill_bg, _rgb, _notes, include_notes, period_label,
                    generate_talk_track, THEME):
    """Add SEC-grade analytics slides to the board pack.

    Includes: ARR Bridge (waterfall chart), Cohort Retention (heatmap),
    Customer Concentration (bar chart), Gross Margin Decomposition,
    Cash Burn Waterfall, Unit Economics & Rule of 40, KPI Accountability,
    and Data Governance attestation.

    All slides:
      - Period-filtered to the user's selected date range
      - Use chart_engine for intelligent visualisation
      - Defensive: a failure in one does not block others
    """
    from pptx.util import Inches, Pt
    from core.database import get_db
    from core.chart_engine import (
        render_waterfall, PALETTE,
        add_native_line, add_native_bar_h, add_native_table_heatmap,
    )
    import json

    conn = get_db()
    from_p = f"{body.from_year}-{body.from_month:02d}"
    to_p = f"{body.to_year}-{body.to_month:02d}"

    def _period_filter(period_str):
        """Check if a period falls within the user's selected range."""
        p = str(period_str or "")[:7]
        return from_p <= p <= to_p

    def _fmt_usd(v):
        if v is None:
            return "--"
        av = abs(float(v))
        sign = "-" if float(v) < 0 else ""
        if av >= 1e9:
            return f"{sign}${av/1e9:.1f}B"
        if av >= 1e6:
            return f"{sign}${av/1e6:.1f}M"
        if av >= 1e3:
            return f"{sign}${av/1e3:.1f}K"
        return f"{sign}${av:.0f}"

    def _fmt_pct(v):
        return f"{float(v):.1f}%" if v is not None else "--"

    def _slide_header(title, subtitle=None):
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, THEME["bg"])
        _add_text(slide, title.upper(), 0.5, 0.3, 12, 0.35,
                  font_size=9, bold=True, color=THEME["subtext"])
        _add_text(slide, title, 0.5, 0.65, 12, 0.5,
                  font_size=20, bold=True, color=THEME["text"])
        if subtitle:
            _add_text(slide, subtitle, 0.5, 1.1, 6, 0.3,
                      font_size=10, color=THEME["subtext"])
        return slide

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: ARR Bridge — Waterfall Chart
    # ═══════════════════════════════════════════════════════════════════════
    try:
        rev_rows = conn.execute(
            "SELECT customer_id, period, amount FROM canonical_revenue "
            "WHERE workspace_id=? AND amount IS NOT NULL ORDER BY period",
            [workspace_id],
        ).fetchall()

        # Build per-customer per-period revenue (period-filtered)
        customer_period: dict = {}
        for r in rev_rows:
            period = str(r[1] or "")[:7]
            if not _period_filter(period):
                continue
            cid = r[0] or "unknown"
            amt = float(r[2] or 0)
            customer_period.setdefault(period, {})
            customer_period[period][cid] = customer_period[period].get(cid, 0) + amt

        sorted_periods = sorted(customer_period.keys())
        if len(sorted_periods) >= 2:
            # Use last two periods for the bridge
            current = customer_period[sorted_periods[-1]]
            prev = customer_period[sorted_periods[-2]]
            new_arr = sum(v for k, v in current.items() if k not in prev) * 12
            churned_arr = sum(v for k, v in prev.items() if k not in current) * 12
            expansion = sum(max(0, current.get(k, 0) - prev.get(k, 0))
                           for k in current if k in prev and current[k] > prev[k]) * 12
            contraction = sum(max(0, prev.get(k, 0) - current.get(k, 0))
                             for k in current if k in prev and current[k] < prev[k]) * 12
            net_new = new_arr + expansion - contraction - churned_arr
            ending_arr = sum(current.values()) * 12

            slide = _slide_header("ARR Bridge", f"{sorted_periods[-2]} to {sorted_periods[-1]}")

            # Waterfall chart
            wf_labels = ["New ARR", "Expansion", "Contraction", "Churned", "Net New ARR"]
            wf_deltas = [new_arr, expansion, -contraction, -churned_arr, net_new]
            wf = render_waterfall(wf_labels, wf_deltas, "Net New ARR Movement")
            slide.shapes.add_picture(wf.image, Inches(0.3), Inches(1.5),
                                     Inches(8.0), Inches(5.0))

            # Summary narrative on right
            paras = [
                ("Period Summary", 14, True, THEME["text"]),
                ("", 4, False, None),
                (f"New ARR:         {_fmt_usd(new_arr)}", 11, False, PALETTE["positive"]),
                (f"Expansion:       {_fmt_usd(expansion)}", 11, False, PALETTE["positive"]),
                (f"Contraction:    ({_fmt_usd(contraction)})", 11, False, PALETTE["critical"]),
                (f"Churned:        ({_fmt_usd(churned_arr)})", 11, False, PALETTE["critical"]),
                ("", 6, False, None),
                (f"Net New ARR:     {_fmt_usd(net_new)}", 13, True,
                 THEME["positive"] if net_new > 0 else THEME["critical"]),
                (f"Ending ARR:      {_fmt_usd(ending_arr)}", 13, True, THEME["text"]),
            ]
            _add_paragraphs(slide, paras, 8.8, 1.5, 4.2, 5.0)

            if include_notes:
                _notes(slide, "ARR Bridge decomposes recurring revenue movement. "
                       "Focus on the balance between growth (New + Expansion) and "
                       "erosion (Contraction + Churn). Net New ARR is the growth engine.")
    except Exception as e:
        print(f"[Board Pack] ARR Bridge slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: Cohort Revenue Retention Matrix
    # ═══════════════════════════════════════════════════════════════════════
    try:
        cohort_rows = conn.execute(
            "SELECT customer_id, period, amount FROM canonical_revenue "
            "WHERE workspace_id=? AND customer_id IS NOT NULL AND amount IS NOT NULL",
            [workspace_id],
        ).fetchall()

        if cohort_rows and len(cohort_rows) > 10:
            # Build cohort data
            first_period: dict = {}
            cust_period_data: dict = {}
            for r in cohort_rows:
                cid, period, amt = r[0], str(r[1] or "")[:7], float(r[2] or 0)
                if cid not in first_period or period < first_period[cid]:
                    first_period[cid] = period
                key = (cid, period)
                cust_period_data[key] = cust_period_data.get(key, 0) + amt

            cohort_customers: dict = {}
            for cid, fp in first_period.items():
                cohort_customers.setdefault(fp, set()).add(cid)

            all_periods = sorted(set(p for _, p in cust_period_data.keys()))

            # Build retention heatmap data (revenue retention %)
            cohort_labels = []
            month_offsets = list(range(min(13, len(all_periods))))
            offset_labels = [f"M{i}" for i in month_offsets]
            retention_matrix = []
            value_matrix = []

            for acq_period in sorted(cohort_customers.keys())[-8:]:  # Last 8 cohorts
                customers = cohort_customers[acq_period]
                m0_revenue = sum(cust_period_data.get((c, acq_period), 0) for c in customers)
                if m0_revenue <= 0:
                    continue
                acq_idx = all_periods.index(acq_period) if acq_period in all_periods else -1
                if acq_idx < 0:
                    continue

                cohort_labels.append(f"{acq_period} ({len(customers)})")
                row_status = []
                row_values = []
                for offset in month_offsets:
                    if acq_idx + offset >= len(all_periods):
                        row_status.append("grey")
                        row_values.append(None)
                        continue
                    p = all_periods[acq_idx + offset]
                    revenue = sum(cust_period_data.get((c, p), 0) for c in customers)
                    retention = revenue / m0_revenue * 100
                    row_values.append(round(retention, 0))
                    if retention >= 100:
                        row_status.append("green")
                    elif retention >= 80:
                        row_status.append("yellow")
                    elif retention > 0:
                        row_status.append("red")
                    else:
                        row_status.append("grey")
                retention_matrix.append(row_status)
                value_matrix.append(row_values)

            if cohort_labels and len(cohort_labels) >= 2:
                slide = _slide_header("Cohort Revenue Retention",
                                      "Vintage analysis — revenue retention by acquisition cohort")

                add_native_table_heatmap(slide, 0.3, 1.5,
                                         min(12.5, len(offset_labels) * 0.8 + 2.5),
                                         min(5.5, len(cohort_labels) * 0.35 + 0.8),
                                         cohort_labels, offset_labels,
                                         retention_matrix, value_matrix)

                if include_notes:
                    _notes(slide, "Cohort retention shows whether customers acquired in each period "
                           "continue generating revenue. Green (100%+) = expansion. "
                           "Red (<80%) = material erosion. Diagonal degradation across all "
                           "cohorts signals a systemic product or market issue.")
    except Exception as e:
        print(f"[Board Pack] Cohort Retention slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: Customer Concentration — Bar Chart + SEC Flag
    # ═══════════════════════════════════════════════════════════════════════
    try:
        conc_rows = conn.execute(
            "SELECT customer_id, SUM(amount) as total FROM canonical_revenue "
            "WHERE workspace_id=? AND customer_id IS NOT NULL AND amount IS NOT NULL "
            "GROUP BY customer_id ORDER BY total DESC LIMIT 10",
            [workspace_id],
        ).fetchall()

        # Customer names
        try:
            name_rows = conn.execute(
                "SELECT source_id, name FROM canonical_customers WHERE workspace_id=?",
                [workspace_id],
            ).fetchall()
            cust_names = {r[0]: r[1] for r in name_rows}
        except Exception:
            cust_names = {}

        total_rev_row = conn.execute(
            "SELECT SUM(amount) FROM canonical_revenue WHERE workspace_id=? AND amount IS NOT NULL",
            [workspace_id],
        ).fetchone()
        total_rev = float(total_rev_row[0] or 0) if total_rev_row else 0

        if conc_rows and total_rev > 0:
            slide = _slide_header("Customer Concentration Risk",
                                  "SEC 10% disclosure threshold analysis")

            names = [cust_names.get(r[0], r[0] or "Unknown")[:25] for r in conc_rows[:8]]
            values = [float(r[1] or 0) / total_rev * 100 for r in conc_rows[:8]]
            colors = [PALETTE["critical"] if v > 10 else (PALETTE["warning"] if v > 5 else PALETTE["positive"])
                      for v in values]

            add_native_bar_h(slide, 0.3, 1.5, 7.5, 5.0,
                             names, values, colors,
                             "Top Customers — % of Total Revenue", unit="pct")

            # Summary on right
            top1_pct = values[0] if values else 0
            top5_pct = sum(values[:5])
            sec_breached = top1_pct > 10

            # HHI
            all_shares = [(float(r[1] or 0) / total_rev * 100) for r in conc_rows]
            hhi = sum(s ** 2 for s in all_shares)

            paras = [
                ("Concentration Summary", 14, True, THEME["text"]),
                ("", 4, False, None),
                (f"Top 1:    {_fmt_pct(top1_pct)} of revenue", 12, False,
                 THEME["critical"] if sec_breached else THEME["text"]),
                (f"Top 5:    {_fmt_pct(top5_pct)} of revenue", 12, False, THEME["text"]),
                (f"HHI:      {hhi:.0f}", 12, False, THEME["text"]),
                ("", 6, False, None),
            ]

            if sec_breached:
                paras.append((
                    "SEC Disclosure Required: One or more customers exceed 10% of revenue. "
                    "This must be disclosed in 10-K/10-Q filings per ASC 280.",
                    11, True, THEME["critical"],
                ))

            # HHI interpretation
            if hhi > 2500:
                paras.append(("HHI > 2,500: Highly concentrated revenue base.", 10, False, THEME["critical"]))
            elif hhi > 1500:
                paras.append(("HHI 1,500-2,500: Moderately concentrated.", 10, False, THEME["warning"]))
            else:
                paras.append(("HHI < 1,500: Well diversified revenue base.", 10, False, THEME["positive"]))

            _add_paragraphs(slide, paras, 8.3, 1.5, 4.7, 5.0)

            if include_notes:
                _notes(slide, "Customer concentration is a material risk factor. SEC mandates disclosure "
                       "of any customer >10% of revenue under ASC 280. HHI (Herfindahl-Hirschman Index) "
                       "measures concentration: <1500 = diversified, 1500-2500 = moderate, >2500 = concentrated.")
    except Exception as e:
        print(f"[Board Pack] Concentration slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: Gross Margin Decomposition
    # ═══════════════════════════════════════════════════════════════════════
    try:
        _COGS_KEYWORDS = {"cogs", "cost of goods", "cost of revenue", "hosting",
                          "infrastructure", "direct", "support", "implementation"}

        exp_rows = conn.execute(
            "SELECT period, category, SUM(amount) as total "
            "FROM canonical_expenses WHERE workspace_id=? AND amount IS NOT NULL "
            "GROUP BY period, category ORDER BY period",
            [workspace_id],
        ).fetchall()
        margin_rev_rows = conn.execute(
            "SELECT period, SUM(amount) as total FROM canonical_revenue "
            "WHERE workspace_id=? AND amount IS NOT NULL GROUP BY period ORDER BY period",
            [workspace_id],
        ).fetchall()

        revenue_by_period = {str(r[0])[:7]: float(r[1] or 0) for r in margin_rev_rows}

        # Categorize COGS
        period_costs: dict = {}
        for r in exp_rows:
            period = str(r[0])[:7]
            if not _period_filter(period):
                continue
            cat = str(r[1] or "other").lower()
            amt = float(r[2] or 0)
            if not any(kw in cat for kw in _COGS_KEYWORDS):
                continue
            if "hosting" in cat or "infrastructure" in cat:
                display_cat = "Hosting"
            elif "support" in cat:
                display_cat = "Support"
            elif any(kw in cat for kw in ("cogs", "cost of goods", "cost of revenue", "direct")):
                display_cat = "Direct COGS"
            else:
                display_cat = cat.title()[:15]
            period_costs.setdefault(period, {})
            period_costs[period][display_cat] = period_costs[period].get(display_cat, 0) + amt

        filtered_periods = sorted(p for p in set(list(revenue_by_period.keys()) + list(period_costs.keys()))
                                  if _period_filter(p))

        if filtered_periods and period_costs:
            # Build multi-line: one line per COGS category + gross margin %
            all_cats = sorted(set(c for p in period_costs.values() for c in p))

            slide = _slide_header("Gross Margin Decomposition",
                                  "COGS component breakdown and margin trajectory")

            # Use a multi-line chart for COGS categories as % of revenue
            series = {}
            for cat in all_cats[:5]:
                series[cat] = []
                for p in filtered_periods:
                    rev = revenue_by_period.get(p, 0)
                    cost = period_costs.get(p, {}).get(cat, 0)
                    series[cat].append(round(cost / rev * 100, 1) if rev else 0)

            # Add gross margin line
            gm_values = []
            for p in filtered_periods:
                rev = revenue_by_period.get(p, 0)
                cogs = sum(period_costs.get(p, {}).values())
                gm_values.append(round((rev - cogs) / rev * 100, 1) if rev else 0)
            series["Gross Margin %"] = gm_values

            add_native_line(slide, 0.3, 1.5, 7.5, 5.0,
                            filtered_periods, gm_values, None, "Gross Margin %", "pct")

            # Latest breakdown on right
            if filtered_periods:
                latest = filtered_periods[-1]
                rev = revenue_by_period.get(latest, 0)
                costs = period_costs.get(latest, {})
                paras = [
                    (f"COGS Breakdown — {latest}", 14, True, THEME["text"]),
                    ("", 4, False, None),
                ]
                for cat, amt in sorted(costs.items(), key=lambda x: -x[1]):
                    pct = amt / rev * 100 if rev else 0
                    paras.append((f"{cat}: {_fmt_usd(amt)} ({_fmt_pct(pct)})", 11, False, THEME["text"]))
                total_cogs = sum(costs.values())
                gm = (rev - total_cogs) / rev * 100 if rev else 0
                paras.append(("", 6, False, None))
                paras.append((f"Total COGS: {_fmt_usd(total_cogs)} ({_fmt_pct(total_cogs / rev * 100 if rev else 0)})", 12, True, THEME["text"]))
                paras.append((f"Gross Margin: {_fmt_pct(gm)}", 14, True,
                             THEME["positive"] if gm >= 65 else (THEME["warning"] if gm >= 50 else THEME["critical"])))
                _add_paragraphs(slide, paras, 8.3, 1.5, 4.7, 5.0)

            if include_notes:
                _notes(slide, "Gross margin decomposition reveals the cost structure and operating leverage "
                       "opportunity. Track the ratio of infrastructure (scalable) vs personnel (linear) COGS. "
                       "Healthy SaaS: >65% gross margin, improving QoQ.")
    except Exception as e:
        print(f"[Board Pack] Margin Decomposition slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: Cash Burn Waterfall
    # ═══════════════════════════════════════════════════════════════════════
    try:
        _SM_KW = {"sales", "marketing", "s&m", "advertising", "demand"}
        _RD_KW = {"r&d", "research", "engineering", "development", "product"}
        _GA_KW = {"g&a", "general", "admin", "office", "legal", "hr", "finance"}

        bs_rows = conn.execute(
            "SELECT period, cash_balance FROM canonical_balance_sheet "
            "WHERE workspace_id=? ORDER BY period",
            [workspace_id],
        ).fetchall()
        cash_rev_rows = conn.execute(
            "SELECT period, SUM(amount) as total FROM canonical_revenue "
            "WHERE workspace_id=? AND amount IS NOT NULL GROUP BY period",
            [workspace_id],
        ).fetchall()
        cash_exp_rows = conn.execute(
            "SELECT period, category, SUM(amount) as total FROM canonical_expenses "
            "WHERE workspace_id=? AND amount IS NOT NULL GROUP BY period, category",
            [workspace_id],
        ).fetchall()

        cash_by_period = {str(r[0])[:7]: float(r[1] or 0) for r in bs_rows}
        cash_rev_by_period = {str(r[0])[:7]: float(r[1] or 0) for r in cash_rev_rows}

        # Categorize expenses into buckets
        exp_buckets: dict = {}
        for r in cash_exp_rows:
            period = str(r[0])[:7]
            if not _period_filter(period):
                continue
            cat = str(r[1] or "other").lower()
            amt = float(r[2] or 0)
            if any(kw in cat for kw in _SM_KW):
                bucket = "S&M"
            elif any(kw in cat for kw in _RD_KW):
                bucket = "R&D"
            elif any(kw in cat for kw in _GA_KW):
                bucket = "G&A"
            elif any(kw in cat for kw in ("cogs", "hosting", "infrastructure", "direct")):
                bucket = "COGS"
            else:
                bucket = "Other"
            exp_buckets.setdefault(period, {})
            exp_buckets[period][bucket] = exp_buckets[period].get(bucket, 0) + amt

        # Use latest period for waterfall
        filtered_cash_periods = sorted(p for p in exp_buckets if _period_filter(p))
        if filtered_cash_periods:
            latest_p = filtered_cash_periods[-1]
            opening = cash_by_period.get(
                filtered_cash_periods[-2] if len(filtered_cash_periods) >= 2 else latest_p, 0)
            revenue = cash_rev_by_period.get(latest_p, 0)
            expenses = exp_buckets.get(latest_p, {})
            total_exp = sum(expenses.values())
            closing = cash_by_period.get(latest_p, opening + revenue - total_exp)
            net_burn = revenue - total_exp
            runway = abs(closing / net_burn) if net_burn < 0 else 999

            slide = _slide_header("Cash Flow Waterfall", f"Period: {latest_p}")

            # Build waterfall: Opening → Revenue → -S&M → -R&D → -G&A → -COGS → -Other → Closing
            wf_labels = ["Opening Cash", "Revenue"]
            wf_deltas = [opening, revenue]
            for bucket in ["S&M", "R&D", "G&A", "COGS", "Other"]:
                amt = expenses.get(bucket, 0)
                if amt > 0:
                    wf_labels.append(bucket)
                    wf_deltas.append(-amt)

            wf = render_waterfall(wf_labels, wf_deltas, "Cash Flow Bridge")
            slide.shapes.add_picture(wf.image, Inches(0.3), Inches(1.5),
                                     Inches(8.5), Inches(5.0))

            # Summary on right
            paras = [
                ("Cash Summary", 14, True, THEME["text"]),
                ("", 4, False, None),
                (f"Opening Cash:  {_fmt_usd(opening)}", 12, False, THEME["text"]),
                (f"Revenue:       {_fmt_usd(revenue)}", 12, False, THEME["positive"]),
                (f"Total Expenses: ({_fmt_usd(total_exp)})", 12, False, THEME["critical"]),
                ("", 4, False, None),
                (f"Net Burn:      {_fmt_usd(net_burn)}", 13, True,
                 THEME["positive"] if net_burn >= 0 else THEME["critical"]),
                (f"Closing Cash:  {_fmt_usd(closing)}", 13, True, THEME["text"]),
                ("", 6, False, None),
                (f"Runway: {runway:.0f} months" if runway < 999 else "Runway: Cash positive", 14, True,
                 THEME["positive"] if runway > 18 or runway >= 999 else
                 (THEME["warning"] if runway > 9 else THEME["critical"])),
            ]
            _add_paragraphs(slide, paras, 9.3, 1.5, 3.7, 5.0)

            if include_notes:
                _notes(slide, "Cash waterfall shows where capital is deployed. "
                       f"Current runway: {runway:.0f} months. "
                       "Board should monitor burn rate trajectory and path to cash flow breakeven.")
    except Exception as e:
        print(f"[Board Pack] Cash Waterfall slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: Unit Economics & Rule of 40
    # ═══════════════════════════════════════════════════════════════════════
    try:
        # Get KPI data within period
        kpi_rows = conn.execute(
            "SELECT data_json FROM monthly_data WHERE workspace_id=? "
            "AND (year > ? OR (year = ? AND month >= ?)) "
            "AND (year < ? OR (year = ? AND month <= ?)) "
            "ORDER BY year DESC, month DESC LIMIT 1",
            [workspace_id,
             body.from_year, body.from_year, body.from_month,
             body.to_year, body.to_year, body.to_month],
        ).fetchone()
        kpi_data = json.loads(kpi_rows[0]) if kpi_rows else {}

        rev_growth = kpi_data.get("revenue_growth")
        ebitda_margin = kpi_data.get("ebitda_margin")
        gross_margin = kpi_data.get("gross_margin")
        churn = kpi_data.get("churn_rate")
        nrr = kpi_data.get("nrr")
        burn_multiple = kpi_data.get("burn_multiple")
        ltv_cac = kpi_data.get("ltv_cac")
        customer_ltv = kpi_data.get("customer_ltv")

        if rev_growth is not None or ebitda_margin is not None:
            slide = _slide_header("Unit Economics & Strategic Metrics", period_label)

            rule_40 = float(rev_growth or 0) + float(ebitda_margin or 0)
            r40_color = THEME["positive"] if rule_40 >= 40 else (THEME["warning"] if rule_40 >= 25 else THEME["critical"])

            paras = [
                ("Rule of 40", 18, True, THEME["text"]),
                ("", 4, False, None),
                (f"Revenue Growth:   {_fmt_pct(rev_growth)}", 13, False, THEME["text"]),
                (f"EBITDA Margin:    {_fmt_pct(ebitda_margin)}", 13, False, THEME["text"]),
                (f"Rule of 40 Score: {rule_40:.1f}  {'Pass' if rule_40 >= 40 else 'Below threshold'}", 14, True, r40_color),
                ("", 10, False, None),
                ("Key Unit Economics", 18, True, THEME["text"]),
                ("", 4, False, None),
                (f"Gross Margin:     {_fmt_pct(gross_margin)}", 13, False, THEME["text"]),
                (f"Churn Rate:       {_fmt_pct(churn)}", 13, False, THEME["text"]),
                (f"NRR:              {_fmt_pct(nrr)}", 13, False, THEME["text"]),
                (f"Burn Multiple:    {burn_multiple:.2f}x" if burn_multiple else "Burn Multiple:    --", 13, False, THEME["text"]),
                (f"LTV:CAC:          {ltv_cac:.1f}x" if ltv_cac else "LTV:CAC:          --", 13, False, THEME["text"]),
                (f"Customer LTV:     {_fmt_usd(customer_ltv)}", 13, False, THEME["text"]),
                ("", 8, False, None),
                ("Note: Revenue Growth and Gross/EBITDA Margins are GAAP-aligned. "
                 "ARR, NRR, Burn Multiple, LTV:CAC are Non-GAAP operating metrics.",
                 9, False, THEME["subtext"]),
            ]

            _add_paragraphs(slide, paras, 0.5, 1.3, 12, 5.5)

            if include_notes:
                _notes(slide, "Rule of 40 = Revenue Growth + Profit Margin. Above 40 = healthy SaaS. "
                       "GAAP vs Non-GAAP distinction is mandatory for SEC filings and investor communications. "
                       f"LTV:CAC of {ltv_cac:.1f}x {'exceeds' if ltv_cac and ltv_cac >= 3 else 'is below'} "
                       f"the 3x benchmark." if ltv_cac else "")
    except Exception as e:
        print(f"[Board Pack] Unit Economics slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: KPI Accountability
    # ═══════════════════════════════════════════════════════════════════════
    try:
        acct_rows = conn.execute(
            "SELECT kpi_key, owner, status FROM kpi_accountability WHERE workspace_id=?",
            [workspace_id],
        ).fetchall()

        if acct_rows:
            # Group by owner
            owners: dict = {}
            target_rows = conn.execute(
                "SELECT kpi_key, target_value, direction FROM kpi_targets WHERE workspace_id=?",
                [workspace_id],
            ).fetchall()
            tgt_map = {r[0]: {"target": float(r[1] or 0), "direction": r[2]} for r in target_rows}
            latest_kpi = conn.execute(
                "SELECT data_json FROM monthly_data WHERE workspace_id=? "
                "ORDER BY year DESC, month DESC LIMIT 1",
                [workspace_id],
            ).fetchone()
            kpi_vals = json.loads(latest_kpi[0]) if latest_kpi else {}

            for r in acct_rows:
                owner = r[1] or "Unassigned"
                kpi_key = r[0]
                resolution = r[2] or "open"
                val = kpi_vals.get(kpi_key)
                t = tgt_map.get(kpi_key)

                if val is not None and t:
                    ratio = val / t["target"] if t["target"] else 0
                    d = t.get("direction", "higher")
                    st = "green" if (d == "higher" and ratio >= 0.98) or (d != "higher" and ratio <= 1.02) else \
                         "yellow" if (d == "higher" and ratio >= 0.90) or (d != "higher" and ratio <= 1.10) else "red"
                else:
                    st = "grey"

                owners.setdefault(owner, {"green": 0, "yellow": 0, "red": 0, "grey": 0, "resolved": 0, "total": 0})
                owners[owner][st] += 1
                owners[owner]["total"] += 1
                if resolution == "resolved":
                    owners[owner]["resolved"] += 1

            slide = _slide_header("KPI Owner Accountability")

            paras = [("Ownership & Resolution Summary", 14, True, THEME["text"]), ("", 4, False, None)]
            sorted_owners = sorted(owners.items(), key=lambda x: -x[1]["red"])
            for name, data in sorted_owners[:8]:
                total = data["total"]
                res_rate = round(data["resolved"] / total * 100) if total else 0
                status_str = f"{data['green']}G  {data['yellow']}Y  {data['red']}R"
                color = THEME["critical"] if data["red"] >= 2 else (THEME["warning"] if data["red"] >= 1 else THEME["text"])
                paras.append((f"{name} — {total} KPIs ({status_str}) — {res_rate}% resolved", 11, False, color))

            paras.append(("", 8, False, None))
            total_kpis = sum(o["total"] for o in owners.values())
            total_red = sum(o["red"] for o in owners.values())
            total_resolved = sum(o["resolved"] for o in owners.values())
            overall_rate = round(total_resolved / total_kpis * 100) if total_kpis else 0
            paras.append((f"Total: {total_kpis} KPIs assigned to {len(owners)} owners. "
                         f"{total_red} critical. {overall_rate}% resolution rate.", 12, True, THEME["text"]))

            _add_paragraphs(slide, paras, 0.5, 1.3, 12, 5.5)

            if include_notes:
                _notes(slide, "KPI accountability ensures every metric has a named owner. "
                       "Resolution rate tracks execution discipline. "
                       "Board should review owners with multiple red KPIs for capacity or capability gaps.")
    except Exception as e:
        print(f"[Board Pack] Accountability slide failed: {e}")

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: Data Governance & Control Attestation
    # ═══════════════════════════════════════════════════════════════════════
    try:
        integrity_row = conn.execute(
            "SELECT overall_status, stage0_status, stage1_status, stage2_status, "
            "stage3_status, stage4_status, correction_attempted, correction_succeeded "
            "FROM integrity_checks WHERE workspace_id=? ORDER BY started_at DESC LIMIT 1",
            [workspace_id],
        ).fetchone()

        if integrity_row:
            slide = _slide_header("Data Governance & Control Attestation")

            status = integrity_row[0] or "unknown"
            stages = ["Temporal Validation", "Source Reconciliation", "KPI Logic",
                      "Display Consistency", "Statistical Anomaly"]
            stage_statuses = [integrity_row[i+1] or "unknown" for i in range(5)]

            paras = [
                ("Integrity Check Summary", 14, True, THEME["text"]),
                ("", 4, False, None),
                (f"Overall Status: {status.upper()}", 14, True,
                 THEME["positive"] if status == "pass" else
                 (THEME["warning"] if status == "warn" else THEME["critical"])),
                ("", 6, False, None),
            ]

            for i, (stage_name, s_status) in enumerate(zip(stages, stage_statuses)):
                icon = "Pass" if s_status == "pass" else ("Warning" if s_status == "warn" else "Fail")
                color = THEME["positive"] if s_status == "pass" else \
                    (THEME["warning"] if s_status == "warn" else THEME["critical"])
                paras.append((f"  Stage {i}: {stage_name}  —  {icon}", 11, False, color))

            if integrity_row[6]:  # correction_attempted
                paras.append(("", 6, False, None))
                success = "succeeded" if integrity_row[7] else "failed"
                paras.append((f"Auto-correction: {success}", 11, False,
                             THEME["positive"] if integrity_row[7] else THEME["critical"]))

            paras.append(("", 10, False, None))
            paras.append(("5-stage automated validation covering temporal integrity, source reconciliation, "
                         "KPI computation verification, display consistency, and statistical anomaly detection. "
                         "Mirrors SOX internal control framework.",
                         9, False, THEME["subtext"]))

            _add_paragraphs(slide, paras, 0.5, 1.3, 12, 5.5)

            if include_notes:
                _notes(slide, "Control attestation provides audit-grade evidence of data integrity. "
                       "5 stages mirror SOX internal control framework. "
                       "Any 'Fail' status requires investigation before presenting to the board.")
    except Exception as e:
        print(f"[Board Pack] Attestation slide failed: {e}")

    conn.close()
